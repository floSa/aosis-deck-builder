"""Verify every shape stays within slide bounds [0, slide_w] × [0, slide_h].

Usage: python tests/check_bounds.py tests/out/roadmap_*.pptx
Exit code 0 if all in bounds; non-zero otherwise.
"""

import sys
from pathlib import Path
from pptx import Presentation


def emu_to_in(v):
    return v / 914400.0


def check(pptx_path, tol_emu=1000):
    prs = Presentation(pptx_path)
    sw, sh = prs.slide_width, prs.slide_height
    violations = []
    extremes = []  # (slide_idx, name, left_in, top_in, right_in, bottom_in)
    for si, slide in enumerate(prs.slides):
        for sh_ in slide.shapes:
            try:
                l, t, w, h = sh_.left, sh_.top, sh_.width, sh_.height
            except Exception:
                continue
            if l is None or w is None:
                continue
            right, bottom = l + w, t + h
            extremes.append((si + 1, sh_.name, l, t, right, bottom))
            if l < -tol_emu or t < -tol_emu or right > sw + tol_emu or bottom > sh + tol_emu:
                violations.append((si + 1, sh_.name, l, t, w, h, right, bottom))
    return prs, violations, extremes


def main(argv):
    rc = 0
    for path in argv[1:]:
        prs, violations, extremes = check(path)
        sw_in, sh_in = emu_to_in(prs.slide_width), emu_to_in(prs.slide_height)
        print(f"\n=== {Path(path).name} (slide {sw_in:.3f}\" x {sh_in:.3f}\") ===")
        if violations:
            print(f"  VIOLATIONS: {len(violations)}")
            for si, name, l, t, w, h, r, b in violations:
                print(f"    slide {si} shape={name!r} pos=({emu_to_in(l):.3f},{emu_to_in(t):.3f}) "
                      f"size=({emu_to_in(w):.3f}x{emu_to_in(h):.3f}) "
                      f"right={emu_to_in(r):.3f} bottom={emu_to_in(b):.3f}")
            rc = 1
        else:
            print(f"  OK — {len(extremes)} shapes all within bounds")
        # Print the 4 most extreme shapes (leftmost, topmost, rightmost, bottommost)
        if extremes:
            by_left   = min(extremes, key=lambda x: x[2])
            by_top    = min(extremes, key=lambda x: x[3])
            by_right  = max(extremes, key=lambda x: x[4])
            by_bottom = max(extremes, key=lambda x: x[5])
            print(f"  Most-left   : slide{by_left[0]}   {by_left[1]:25s} left  ={emu_to_in(by_left[2]):.3f}\"")
            print(f"  Most-top    : slide{by_top[0]}   {by_top[1]:25s} top   ={emu_to_in(by_top[3]):.3f}\"")
            print(f"  Most-right  : slide{by_right[0]}  {by_right[1]:25s} right ={emu_to_in(by_right[4]):.3f}\"")
            print(f"  Most-bottom : slide{by_bottom[0]} {by_bottom[1]:25s} bottom={emu_to_in(by_bottom[5]):.3f}\"")
    return rc


if __name__ == "__main__":
    sys.exit(main(sys.argv))
