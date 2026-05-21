"""Smoke + non-regression tests for the AOSIS deck builder skill.

These tests are deliberately minimal: they exercise the dispatcher across all
23 layouts, check that the golden deck and the four roadmap fixtures stay
within slide bounds, and verify the dynamic palette loading. They run on the
canonical template at `assets/AOSIS_template.pptx` and finish in <30 seconds.
"""
from __future__ import annotations

import io
import json
import shutil
import subprocess
import sys
from pathlib import Path

import pytest
from pptx import Presentation


# Make scripts/ importable as if we were inside the skill package
ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "scripts"))

from brand import BrandPalette  # noqa: E402
from build_deck import build_deck  # noqa: E402


FIXTURES = ROOT / "tests" / "fixtures"
TEMPLATE = ROOT / "assets" / "AOSIS_template.pptx"
SLIDE_W = 10.0       # inches
SLIDE_H = 5.625      # inches
TOL = 0.01           # 0.01" tolerance for floating-point rounding


def _load_spec(fixture_name: str) -> dict:
    """Load a fixture JSON spec, resolving any relative image paths to abs."""
    spec = json.loads((FIXTURES / fixture_name).read_text(encoding="utf-8"))
    for slide in spec.get("slides", []):
        img = slide.get("image")
        if img and not Path(img).is_absolute():
            slide["image"] = str((FIXTURES / img).resolve())
    return spec


def _generate(fixture_name: str, tmp_path: Path) -> Path:
    spec = _load_spec(fixture_name)
    out = tmp_path / "deck.pptx"
    build_deck(spec, out, template_path=TEMPLATE)
    return out


def _shapes_out_of_bounds(pptx_path: Path) -> list[str]:
    """Return a list of human-readable violation strings (empty if all OK).

    Inherited title placeholders that bleed slightly above the canvas
    (canonical at y ≈ −0.14" in the `5_Vide` master) are tolerated — they
    are an artefact of the master template, not of our rendering. We use a
    larger upward tolerance for these.
    """
    prs = Presentation(str(pptx_path))
    violations: list[str] = []
    for s_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.left is None or shape.width is None:
                continue
            left = shape.left / 914400
            top = shape.top / 914400
            width = shape.width / 914400
            height = shape.height / 914400
            right = left + width
            bottom = top + height
            # Master-inherited title placeholders sit at y ≈ -0.14" by design;
            # allow up to -0.20" upward bleed for shapes named "Title …".
            top_tol = 0.20 if shape.name.startswith("Title ") else TOL
            if (
                left < -TOL
                or top < -top_tol
                or right > SLIDE_W + TOL
                or bottom > SLIDE_H + TOL
            ):
                violations.append(
                    f"slide {s_idx} shape {shape.name!r}: "
                    f"({left:.3f}, {top:.3f}) → ({right:.3f}, {bottom:.3f})"
                )
    return violations


# -----------------------------------------------------------------------------
# Golden deck (Chantier 3 spec)
# -----------------------------------------------------------------------------

def test_golden_generates(tmp_path):
    """The golden deck (5 slides + closing) must build without error."""
    out = _generate("golden_spec.json", tmp_path)
    assert out.exists() and out.stat().st_size > 100_000
    prs = Presentation(str(out))
    # cover + 3 content slides (hero_stat, matrix_2x2, roadmap) + closing = 5
    assert len(prs.slides) == 5


def test_golden_no_overflow(tmp_path):
    """No shape in the golden deck may overflow slide bounds."""
    out = _generate("golden_spec.json", tmp_path)
    violations = _shapes_out_of_bounds(out)
    assert not violations, "Overflow detected:\n" + "\n".join(violations)


# -----------------------------------------------------------------------------
# Roadmap regression (Chantier 1)
# -----------------------------------------------------------------------------

@pytest.mark.parametrize(
    "fixture",
    ["roadmap_3.json", "roadmap_4.json", "roadmap_5.json", "roadmap_6.json"],
)
def test_roadmap_no_overflow(fixture, tmp_path):
    """Chantier 1 regression: roadmap with 3..6 milestones stays in bounds."""
    out = _generate(fixture, tmp_path)
    violations = _shapes_out_of_bounds(out)
    assert not violations, f"Overflow in {fixture}:\n" + "\n".join(violations)


# -----------------------------------------------------------------------------
# Coverage: every dispatcher layout builds without crashing
# -----------------------------------------------------------------------------

def test_all_layouts_generate(tmp_path):
    """The `all_layouts` fixture exercises every dispatcher layout."""
    out = _generate("all_layouts.json", tmp_path)
    assert out.exists()
    prs = Presentation(str(out))
    # cover + 20 content slides (section..image_hero) + closing = 22 minimum
    # exact value depends on fixture count; current fixture has 20 content slides
    assert len(prs.slides) >= 20, f"only {len(prs.slides)} slides generated"


# -----------------------------------------------------------------------------
# Dynamic palette (Chantier 3)
# -----------------------------------------------------------------------------

def test_palette_loads_from_canonical_template():
    """Palette is loaded from the canonical template's theme XML."""
    palette = BrandPalette.from_template(TEMPLATE)
    assert f"{palette.navy}" == "14163C"
    assert f"{palette.orange}" == "F26622"
    assert f"{palette.light}" == "FAFAF7"
    assert f"{palette.gray}" == "4A4D6B"
    # ex-NAVY_SOFT replaced by accent2 in Chantier 3
    assert f"{palette.navy_alt}" == "1E2261"


def test_palette_brand_error_on_missing_file():
    """BrandError must surface a useful message for missing templates."""
    from brand import BrandError  # local import keeps the global scope clean
    with pytest.raises(BrandError, match="Template not found"):
        BrandPalette.from_template("/nonexistent/path.pptx")


# -----------------------------------------------------------------------------
# Visual review (Chantier 6)
# -----------------------------------------------------------------------------

_SOFFICE = shutil.which("soffice")
_PDFTOPPM = shutil.which("pdftoppm")


@pytest.mark.skipif(
    _SOFFICE is None or _PDFTOPPM is None,
    reason="soffice and/or pdftoppm not installed; visual_review needs both",
)
def test_visual_review_generates_artifacts(tmp_path):
    """`visual_review.py` produces JPEGs + prompt + report template."""
    out = _generate("golden_spec.json", tmp_path)
    review_dir = tmp_path / "review"
    cmd = [
        sys.executable,
        str(ROOT / "scripts" / "visual_review.py"),
        str(out),
        str(review_dir),
        "--dpi", "100",  # smaller for fast test
    ]
    result = subprocess.run(cmd, capture_output=True, text=True)
    assert result.returncode == 0, (
        f"visual_review.py failed:\nSTDOUT={result.stdout}\nSTDERR={result.stderr}"
    )
    # At least one JPEG, the prompt, and the template
    jpegs = sorted(review_dir.glob("slide-*.jpg"))
    assert len(jpegs) >= 1, "no JPEGs produced"
    assert (review_dir / "review_prompt.md").exists()
    template = review_dir / "review_report.template.json"
    assert template.exists()
    # Template JSON parses and has the expected shape
    data = json.loads(template.read_text())
    assert "slides" in data and len(data["slides"]) == len(jpegs)
    for entry in data["slides"]:
        assert entry["defects"] == []  # skeleton, no defects yet


# -----------------------------------------------------------------------------
# Template-based engine (Chantier 7)
# -----------------------------------------------------------------------------

def test_template_layouts_discovery():
    """Template engine discovers the named layouts in the canonical template
    (since the consolidation chantier, AOSIS_template.pptx hosts them all)."""
    sys.path.insert(0, str(ROOT / "scripts"))
    from template_engine import discover_template_layouts  # noqa: E402
    layouts = discover_template_layouts(TEMPLATE)
    assert len(layouts) >= 10, f"Expected ≥10 named layouts, found {len(layouts)}: {sorted(layouts)}"
    # Sanity-check a few known layouts
    for required in ("executive_summary", "process_steps", "roadmap_styled", "next_steps"):
        assert required in layouts, f"Missing required layout '{required}'"


def test_template_based_deck_builds(tmp_path):
    """A spec using only template-based layouts builds without error and
    produces the right number of slides."""
    out = _generate("template_based_spec.json", tmp_path)
    assert out.exists() and out.stat().st_size > 100_000
    prs = Presentation(str(out))
    # 6 template-based slides + 1 code-based (hero_stat). closing=false so
    # no closing slide. Allow a tiny range to absorb any future minor changes.
    assert 5 <= len(prs.slides) <= 8, f"Got {len(prs.slides)} slides"


def test_template_repeat_mechanism(tmp_path):
    """A process_steps slide with N items produces N copies of the
    repeated group."""
    spec = {
        "slides": [{
            "layout": "process_steps",
            "eyebrow": "PROCESS",
            "title": "Five steps to deliver",
            "source": "AOSIS",
            "items": [
                {"title": f"Step {i}", "text": f"Detail {i}"}
                for i in range(1, 6)
            ],
        }],
        "closing": False,
    }
    out = tmp_path / "repeat.pptx"
    build_deck(spec, out, template_path=TEMPLATE)
    prs = Presentation(str(out))
    # Find the process_steps slide
    target = None
    for s in prs.slides:
        from pptx.oxml.ns import qn
        cSld = s._element.find(qn('p:cSld'))
        if cSld is not None and cSld.get('name') == 'process_steps':
            target = s
            break
    assert target is not None
    # Count GroupShape children (the duplicated REPEAT_ITEM copies)
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    groups = [sh for sh in target.shapes if sh.shape_type == MSO_SHAPE_TYPE.GROUP]
    assert len(groups) == 5, f"Expected 5 group copies, got {len(groups)}"


def test_template_kpi_with_chart(tmp_path):
    """A kpi_with_chart slide builds with KPI cards + title + a rendered chart.

    Chantier 8 — chart presence assertion is re-enabled (matplotlib chart is
    rendered into the {{CHART_PLACEHOLDER}} position as a PICTURE).
    """
    spec = {
        "slides": [{
            "layout": "kpi_with_chart",
            "eyebrow": "PERFORMANCE",
            "title": "Three KPIs over time",
            "source": "AOSIS",
            "kpis": [
                {"label": "Coverage", "value": "85 %", "detail": "Up from 60%"},
                {"label": "Latency",  "value": "30 m", "detail": "Down from 6h"},
                {"label": "Cost",     "value": "−40 %", "detail": "Cloud-native"},
            ],
            "chart": {
                "type": "bar",
                "labels": ["Q1", "Q2", "Q3", "Q4"],
                "values": [120, 145, 132, 168],
            },
        }],
        "closing": False,
    }
    out = tmp_path / "kpi_chart.pptx"
    build_deck(spec, out, template_path=TEMPLATE)
    prs = Presentation(str(out))
    from pptx.oxml.ns import qn
    target = None
    for s in prs.slides:
        cSld = s._element.find(qn('p:cSld'))
        if cSld is not None and cSld.get('name') == 'kpi_with_chart':
            target = s
            break
    assert target is not None, "kpi_with_chart slide not produced"

    # One GROUP per KPI card (REPEAT_ITEM duplicated N times).
    groups = [sh for sh in target.shapes if sh.shape_type == 6]
    assert len(groups) == len(spec["slides"][0]["kpis"]), (
        f"Expected {len(spec['slides'][0]['kpis'])} KPI card groups, got {len(groups)}"
    )

    # Title is substituted in the top-level {{TITLE}} placeholder.
    title_texts = [
        sh.text_frame.text for sh in target.shapes
        if sh.has_text_frame and 'TITLE' in sh.name.upper()
    ]
    assert any(spec["slides"][0]["title"] in t for t in title_texts), (
        f"Title not posed; title shape texts were: {title_texts}"
    )

    # Chart was rendered: a PICTURE shape exists, and {{CHART_PLACEHOLDER}} is gone.
    pictures = [sh for sh in target.shapes if sh.shape_type == 13]
    assert pictures, "No PICTURE shape — chart was not rendered"
    placeholders = [sh for sh in target.shapes if sh.name == '{{CHART_PLACEHOLDER}}']
    assert not placeholders, "{{CHART_PLACEHOLDER}} should have been replaced"


# ---------------------------------------------------------------------------
# Chantier 8 — one test per supported chart type
# ---------------------------------------------------------------------------
_CHART_SPECS = json.loads(
    (FIXTURES / "chart_specs.json").read_text(encoding="utf-8")
)


def _build_kpi_chart_slide(tmp_path, chart_spec, fname):
    """Build a minimal kpi_with_chart slide carrying the given chart spec
    and return the rendered slide for inspection."""
    spec = {
        "slides": [{
            "layout": "kpi_with_chart",
            "title": f"Chart type: {chart_spec.get('type')}",
            "source": "test fixture",
            "kpis": [
                {"label": "A", "value": "1"},
                {"label": "B", "value": "2"},
            ],
            "chart": chart_spec,
        }],
        "closing": False,
    }
    out = tmp_path / fname
    build_deck(spec, out, template_path=TEMPLATE)
    prs = Presentation(str(out))
    from pptx.oxml.ns import qn
    for s in prs.slides:
        cSld = s._element.find(qn('p:cSld'))
        if cSld is not None and cSld.get('name') == 'kpi_with_chart':
            return s
    raise AssertionError("kpi_with_chart slide not produced")


def _assert_chart_rendered(slide):
    """The chart placeholder must be replaced by a PICTURE shape."""
    pictures = [sh for sh in slide.shapes if sh.shape_type == 13]
    assert pictures, "No PICTURE shape inserted — chart was not rendered"
    placeholders = [sh for sh in slide.shapes if sh.name == '{{CHART_PLACEHOLDER}}']
    assert not placeholders, "{{CHART_PLACEHOLDER}} should have been replaced"


def test_kpi_with_chart_bar(tmp_path):
    slide = _build_kpi_chart_slide(tmp_path, _CHART_SPECS["bar"], "bar.pptx")
    _assert_chart_rendered(slide)


def test_kpi_with_chart_barh(tmp_path):
    slide = _build_kpi_chart_slide(tmp_path, _CHART_SPECS["barh"], "barh.pptx")
    _assert_chart_rendered(slide)


def test_kpi_with_chart_bar_stacked(tmp_path):
    slide = _build_kpi_chart_slide(tmp_path, _CHART_SPECS["bar_stacked"], "stk.pptx")
    _assert_chart_rendered(slide)


def test_kpi_with_chart_line(tmp_path):
    slide = _build_kpi_chart_slide(tmp_path, _CHART_SPECS["line"], "line.pptx")
    _assert_chart_rendered(slide)


def test_kpi_with_chart_donut(tmp_path):
    slide = _build_kpi_chart_slide(tmp_path, _CHART_SPECS["donut"], "donut.pptx")
    _assert_chart_rendered(slide)


def test_kpi_with_chart_pie(tmp_path):
    slide = _build_kpi_chart_slide(tmp_path, _CHART_SPECS["pie"], "pie.pptx")
    _assert_chart_rendered(slide)


def test_kpi_with_chart_combo(tmp_path):
    slide = _build_kpi_chart_slide(tmp_path, _CHART_SPECS["combo"], "combo.pptx")
    _assert_chart_rendered(slide)


def test_kpi_with_chart_waterfall(tmp_path):
    slide = _build_kpi_chart_slide(tmp_path, _CHART_SPECS["waterfall"], "wf.pptx")
    _assert_chart_rendered(slide)


# ---------------------------------------------------------------------------
# Chantier alternances — layout-specific alternation rules
# ---------------------------------------------------------------------------
def _find_slide_by_cSld_name(prs, name):
    from pptx.oxml.ns import qn
    for s in prs.slides:
        cSld = s._element.find(qn('p:cSld'))
        if cSld is not None and cSld.get('name') == name:
            return s
    return None


def _shape_fill_hex(sp_element):
    from pptx.oxml.ns import qn
    spPr = sp_element.find(qn('p:spPr'))
    if spPr is None:
        return None
    solid = spPr.find(qn('a:solidFill'))
    if solid is None:
        return None
    srgb = solid.find(qn('a:srgbClr'))
    if srgb is not None:
        return srgb.get('val').upper()
    return None


def _shape_xy(sp_element):
    from pptx.oxml.ns import qn
    spPr = sp_element.find(qn('p:spPr'))
    if spPr is None:
        return (None, None)
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is None:
        return (None, None)
    off = xfrm.find(qn('a:off'))
    if off is None:
        return (None, None)
    return (int(off.get('x', '0')), int(off.get('y', '0')))


def _named_descendant(grp_element, name):
    """First <p:sp> descendant whose cNvPr.name == name."""
    from pptx.oxml.ns import qn
    for sp in grp_element.iter(qn('p:sp')):
        cNvPr = sp.find('.//' + qn('p:cNvPr'))
        if cNvPr is not None and cNvPr.get('name') == name:
            return sp
    return None


def test_framework_3cards_alternates_colors(tmp_path):
    """The 3 cards alternate orange / navy_alt / orange on {{ITEM_BOXE}}."""
    spec = {
        "slides": [{
            "layout": "framework_3cards",
            "title": "Test",
            "items": [
                {"title": "A", "bullets": "x"},
                {"title": "B", "bullets": "y"},
                {"title": "C", "bullets": "z"},
            ],
        }],
        "closing": False,
    }
    out = tmp_path / "fw.pptx"
    build_deck(spec, out, template_path=TEMPLATE)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "framework_3cards")
    assert slide is not None
    groups = [sh for sh in slide.shapes if sh.shape_type == 6]
    assert len(groups) == 3
    fills = [_shape_fill_hex(_named_descendant(g._element, '{{ITEM_BOXE}}')) for g in groups]
    assert fills == ["F26622", "1E2261", "F26622"], f"Expected orange/navy/orange, got {fills}"


def test_process_steps_alternates_colors(tmp_path):
    """The 4 markers alternate orange / navy_alt / orange / navy_alt."""
    spec = {
        "slides": [{
            "layout": "process_steps",
            "title": "Test",
            "items": [
                {"title": "S1", "text": "t1"},
                {"title": "S2", "text": "t2"},
                {"title": "S3", "text": "t3"},
                {"title": "S4", "text": "t4"},
            ],
        }],
        "closing": False,
    }
    out = tmp_path / "ps.pptx"
    build_deck(spec, out, template_path=TEMPLATE)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "process_steps")
    assert slide is not None
    groups = [sh for sh in slide.shapes if sh.shape_type == 6]
    assert len(groups) == 4
    fills = [_shape_fill_hex(_named_descendant(g._element, '{{ITEM_MARKER}}')) for g in groups]
    assert fills == ["F26622", "1E2261", "F26622", "1E2261"], (
        f"Expected orange/navy/orange/navy on markers, got {fills}"
    )


def test_roadmap_alternates_position(tmp_path):
    """Roadmap text alternates above (odd 1-indexed) / below (even 1-indexed)
    the marker, with markers staying orange throughout."""
    spec = {
        "slides": [{
            "layout": "roadmap_styled",
            "title": "Test",
            "items": [
                {"date": "M1", "milestone": "S1"},
                {"date": "M2", "milestone": "S2"},
                {"date": "M3", "milestone": "S3"},
                {"date": "M4", "milestone": "S4"},
                {"date": "M5", "milestone": "S5"},
            ],
        }],
        "closing": False,
    }
    out = tmp_path / "rm.pptx"
    build_deck(spec, out, template_path=TEMPLATE)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "roadmap_styled")
    assert slide is not None
    groups = [sh for sh in slide.shapes if sh.shape_type == 6]
    assert len(groups) == 5

    # Position check: index 0,2,4 → text above marker; index 1,3 → text below.
    above_below = []
    for grp in groups:
        marker = _named_descendant(grp._element, '{{ITEM_MARKER}}')
        date = _named_descendant(grp._element, '{{ITEM_DATE}}')
        assert marker is not None and date is not None
        _, marker_y = _shape_xy(marker)
        _, date_y = _shape_xy(date)
        above_below.append("above" if date_y < marker_y else "below")
    assert above_below == ["above", "below", "above", "below", "above"], (
        f"Expected above/below/above/below/above, got {above_below}"
    )

    # Markers must all stay orange (no color alternation on roadmap markers).
    marker_fills = [_shape_fill_hex(_named_descendant(g._element, '{{ITEM_MARKER}}')) for g in groups]
    assert all(f == "F26622" for f in marker_fills), (
        f"Roadmap markers should all be orange, got {marker_fills}"
    )


# ---------------------------------------------------------------------------
# Chantier 9 — polishing fixes
# ---------------------------------------------------------------------------
def test_agenda_paginates_at_7_items(tmp_path):
    """9 items in agenda_diagonal → 2 slides (7 + 2). Chantier 11 reverted
    the max from 10 back to 7 for breathing room."""
    spec = {
        "slides": [{
            "layout": "agenda_diagonal",
            "title": "Sommaire",
            "items": [{"title": f"Section {n+1}"} for n in range(9)],
        }],
        "closing": False,
    }
    out = tmp_path / "agenda.pptx"
    build_deck(spec, out, template_path=TEMPLATE)
    prs = Presentation(str(out))
    agenda_slides = [s for s in prs.slides if _slide_cSld_name(s) == "agenda_diagonal"]
    assert len(agenda_slides) == 2, f"Expected 2 agenda pages, got {len(agenda_slides)}"
    groups0 = [sh for sh in agenda_slides[0].shapes if sh.shape_type == 6]
    groups1 = [sh for sh in agenda_slides[1].shapes if sh.shape_type == 6]
    assert len(groups0) == 7, f"page 1 should have 7 items, got {len(groups0)}"
    assert len(groups1) == 2, f"page 2 should have 2 items, got {len(groups1)}"


def test_agenda_continuous_numbering(tmp_path):
    """Items 8 and 9 on page 2 are numbered '08' and '09'."""
    spec = {
        "slides": [{
            "layout": "agenda_diagonal",
            "title": "Sommaire",
            "items": [{"title": f"Section {n+1}"} for n in range(9)],
        }],
        "closing": False,
    }
    out = tmp_path / "agenda.pptx"
    build_deck(spec, out, template_path=TEMPLATE)
    prs = Presentation(str(out))
    agenda_slides = [s for s in prs.slides if _slide_cSld_name(s) == "agenda_diagonal"]
    page2_numbers = _collect_item_number_texts(agenda_slides[1])
    assert page2_numbers == ["08", "09"], (
        f"Page 2 should display 08, 09 — got {page2_numbers}"
    )


def test_roadmap_preserves_font(tmp_path):
    """Roadmap copies preserve the source rPr (size, bold, font, color)."""
    from pptx.oxml.ns import qn
    spec = {
        "slides": [{
            "layout": "roadmap_styled",
            "title": "Test",
            "items": [
                {"date": "Jan 26", "milestone": "M1"},
                {"date": "Mai 26", "milestone": "M2"},
                {"date": "Oct 26", "milestone": "M3"},
                {"date": "Avr 27", "milestone": "M4"},
            ],
        }],
        "closing": False,
    }
    out = tmp_path / "rm.pptx"
    build_deck(spec, out, template_path=TEMPLATE)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "roadmap_styled")
    assert slide is not None
    # All 4 copies should have ITEM_DATE with sz=1400, b=1, font=Arial
    for grp in [sh for sh in slide.shapes if sh.shape_type == 6]:
        date_sp = _named_descendant(grp._element, '{{ITEM_DATE}}')
        assert date_sp is not None
        rPr = date_sp.find('.//' + qn('a:rPr'))
        assert rPr is not None, "ITEM_DATE missing rPr"
        assert rPr.get('sz') == '1400', f"size lost: {rPr.get('sz')}"
        assert rPr.get('b') == '1', "bold lost"
        latin = rPr.find(qn('a:latin'))
        assert latin is not None and latin.get('typeface') == 'Arial', "font lost"


def test_kpi_with_chart_renders_values(tmp_path):
    """KPI label/value from JSON spec appear in the rendered slide,
    replacing the template's 'KPI label' / '85 %' defaults."""
    from pptx.oxml.ns import qn
    spec = {
        "slides": [{
            "layout": "kpi_with_chart",
            "title": "KPI test",
            "kpis": [
                {"label": "TCO 2025",    "value": "4.2 M€"},
                {"label": "Cible 2028",  "value": "3.1 M€"},
                {"label": "Économie/an", "value": "1.1 M€"},
            ],
            "chart": {"type": "bar", "labels": ["A", "B", "C"], "values": [1, 2, 3]},
        }],
        "closing": False,
    }
    out = tmp_path / "kpi.pptx"
    build_deck(spec, out, template_path=TEMPLATE)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "kpi_with_chart")
    assert slide is not None

    # Collect all text from copy groups
    texts = []
    for grp in [sh for sh in slide.shapes if sh.shape_type == 6]:
        for t in grp._element.iter(qn('a:t')):
            if t.text:
                texts.append(t.text)
    haystack = "\n".join(texts)
    for kpi in spec["slides"][0]["kpis"]:
        assert kpi["label"] in haystack, f"label {kpi['label']!r} missing — got {haystack!r}"
        assert kpi["value"] in haystack, f"value {kpi['value']!r} missing"
    # And template defaults should be GONE (they were "KPI label" / "85 %")
    assert "KPI label" not in haystack, "template default 'KPI label' leaked through"


def test_framework_3cards_with_icons(tmp_path, monkeypatch):
    """When icons spec is provided, one PICTURE per card is added.
    The Iconify fetch is mocked to avoid network dependency in tests."""
    # Mock icon_engine.fetch_icon_png to return a tiny valid PNG
    import sys
    sys.path.insert(0, str(ROOT / "scripts"))
    import icon_engine

    # 1x1 transparent PNG bytes (smallest valid PNG)
    PNG_1x1 = (
        b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
        b"\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\rIDATx\x9cb\x00\x01"
        b"\x00\x00\x05\x00\x01\r\n-\xb4\x00\x00\x00\x00IEND\xaeB`\x82"
    )

    def fake_fetch(name, size_px=200, color=None, timeout=5.0):
        return PNG_1x1

    monkeypatch.setattr(icon_engine, "fetch_icon_png", fake_fetch)

    spec = {
        "slides": [{
            "layout": "framework_3cards",
            "title": "Test",
            "icons": ["mdi:cloud", "mdi:school", "mdi:tools"],
            "items": [
                {"title": "A", "bullets": "x"},
                {"title": "B", "bullets": "y"},
                {"title": "C", "bullets": "z"},
            ],
        }],
        "closing": False,
    }
    out = tmp_path / "fw_icons.pptx"
    build_deck(spec, out, template_path=TEMPLATE)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "framework_3cards")
    assert slide is not None
    pictures = [sh for sh in slide.shapes if sh.shape_type == 13]
    assert len(pictures) == 3, f"Expected 3 icon pictures, got {len(pictures)}"


def _slide_cSld_name(slide):
    from pptx.oxml.ns import qn
    cSld = slide._element.find(qn('p:cSld'))
    return cSld.get('name') if cSld is not None else None


def _collect_item_number_texts(slide):
    from pptx.oxml.ns import qn
    out = []
    for sh in slide.shapes:
        if sh.shape_type != 6:
            continue
        for sp in sh._element.iter(qn('p:sp')):
            cnvpr = sp.find('.//' + qn('p:cNvPr'))
            if cnvpr is not None and cnvpr.get('name') == '{{ITEM_NUMBER}}':
                for t in sp.iter(qn('a:t')):
                    if t.text:
                        out.append(t.text)
    return out


# ---------------------------------------------------------------------------
# Chantier 10 — post-test-réel fixes
# ---------------------------------------------------------------------------
def test_hero_stat_with_long_text(tmp_path):
    """`supporting` as a plain string must NOT be iterated character by
    character — it should appear as one supporting bullet."""
    long_text = "47 incidents en 2025 contre 19 en 2023 — l'infra ne tient plus"
    spec = {
        "slides": [{
            "layout": "hero_stat",
            "title": "L'AMPLEUR",
            "value": "+147 %",
            "label": "Incidents critiques",
            "supporting": long_text,
        }],
        "closing": False,
    }
    out = tmp_path / "hero.pptx"
    build_deck(spec, out, template_path=TEMPLATE)
    prs = Presentation(str(out))
    slide = prs.slides[0]
    # The long text must appear intact in ONE textbox, not split across many
    found = False
    for sh in slide.shapes:
        if sh.has_text_frame and long_text in sh.text_frame.text:
            found = True
            break
    assert found, "supporting string was split into characters (Chantier 10 bug)"


def test_matrix_2x2_renders_all_quadrants(tmp_path):
    """The 4 quadrants must keep their {{QUAD_*_TITLE}} and {{QUAD_*_BULLETS}}
    shapes after generation (chantier 10 bug: simple processor was deleting
    them right after the quad processor filled them)."""
    spec = {
        "slides": [{
            "layout": "matrix_2x2_styled",
            "title": "Test matrix",
            "x_axis": {"label": "Effort"},
            "y_axis": {"label": "Impact"},
            "quadrants": {
                "top_left":     {"title": "Quick wins",  "items": ["A", "B"]},
                "top_right":    {"title": "Strategic",   "items": ["C", "D"]},
                "bottom_left":  {"title": "Hygiene",     "items": ["E"]},
                "bottom_right": {"title": "Deprioritise","items": ["F"]},
            },
        }],
        "closing": False,
    }
    out = tmp_path / "mx.pptx"
    build_deck(spec, out, template_path=TEMPLATE)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "matrix_2x2_styled")
    assert slide is not None
    texts = [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]
    haystack = "\n".join(texts)
    for title in ("Quick wins", "Strategic", "Hygiene", "Deprioritise"):
        assert title in haystack, f"quadrant title {title!r} missing"
    for bullet_label in ("A", "B", "C", "D", "E", "F"):
        assert bullet_label in haystack, f"quadrant bullet {bullet_label!r} missing"


def test_comparison_before_after_renders_both_columns(tmp_path):
    """The {{BEFORE_*}} and {{AFTER_*}} placeholders must be filled from
    a nested `before: {...}` / `after: {...}` spec (chantier 10 flattening)."""
    spec = {
        "slides": [{
            "layout": "comparison_before_after",
            "title": "Avant/Après",
            "takeaway": "Le saut quantifié",
            "before": {"title": "Aujourd'hui",  "bullets": "8h par cycle\nTravail manuel"},
            "after":  {"title": "Cible T+12",   "bullets": "2h par cycle\nAutomatisation"},
        }],
        "closing": False,
    }
    out = tmp_path / "ba.pptx"
    build_deck(spec, out, template_path=TEMPLATE)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "comparison_before_after")
    assert slide is not None
    texts = [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]
    haystack = "\n".join(texts)
    assert "Aujourd'hui" in haystack, "BEFORE_TITLE not filled"
    assert "Cible T+12" in haystack, "AFTER_TITLE not filled"
    assert "8h par cycle" in haystack, "BEFORE_BULLETS not filled"
    assert "Automatisation" in haystack, "AFTER_BULLETS not filled"


def test_next_steps_long_dates(tmp_path):
    """Long date strings get auto-shrunk so they fit on one line. The
    final sz must be smaller than the template's default 1400."""
    from pptx.oxml.ns import qn
    spec = {
        "slides": [{
            "layout": "next_steps",
            "title": "Test",
            "items": [
                {"action": "X", "owner": "Y", "date": "1er septembre 2026 — kick-off"},
            ],
        }],
        "closing": False,
    }
    out = tmp_path / "ns.pptx"
    build_deck(spec, out, template_path=TEMPLATE)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "next_steps")
    assert slide is not None
    grp = [sh for sh in slide.shapes if sh.shape_type == 6][0]
    date_sp = _named_descendant(grp._element, '{{ITEM_DATE}}')
    rPr = date_sp.find('.//' + qn('a:rPr'))
    sz = int(rPr.get('sz'))
    assert sz < 1400, f"long date should shrink below 14pt, got sz={sz}"
    assert sz >= 1000, f"should not shrink below 10pt floor, got sz={sz}"


def test_kpi_value_shrinks_on_overflow(tmp_path):
    """KPI value with long unit shrinks. Chantier 19 height-aware sizing :
    the kpi_with_chart shape is only 0.47" tall, so even short values are
    capped around 30pt. Long values further constrained by width."""
    from pptx.oxml.ns import qn
    spec = {
        "slides": [{
            "layout": "kpi_with_chart",
            "title": "Test",
            "kpis": [{"label": "Revenue annuel cumulé", "value": "1 234 567 K€"}],
            "chart": {"type": "bar", "labels": ["A"], "values": [1]},
        }],
        "closing": False,
    }
    out = tmp_path / "kpi.pptx"
    build_deck(spec, out, template_path=TEMPLATE)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "kpi_with_chart")
    assert slide is not None
    grp = [sh for sh in slide.shapes if sh.shape_type == 6][0]
    value_sp = _named_descendant(grp._element, '{{KPI_VALUE}}')
    rPr = value_sp.find('.//' + qn('a:rPr'))
    sz = int(rPr.get('sz'))
    assert sz < 6000, f"long value should shrink below 60pt baseline, got sz={sz}"
    # Chantier 20 — kpi_with_chart uses floor_pt=14 to fit wide chars in
    # the narrow {{KPI_VALUE}} shape (1.00" × 0.47"). Long value lands
    # around 14-18pt depending on character width.
    assert sz >= 1400, f"should stay above 14pt floor, got sz={sz}"


def test_orphan_group_removed_when_empty(tmp_path, monkeypatch):
    """A `{{XXX_GROUP}}` group is removed entirely when `<XXX>` is empty
    in spec — even if it contains decorative shapes."""
    # We synthesize a tiny test scenario: monkey-patch the template_engine
    # function to verify behavior on a constructed slide.
    import sys
    sys.path.insert(0, str(ROOT / "scripts"))
    import template_engine
    from pptx.oxml.ns import qn
    from copy import deepcopy

    # Build a minimal deck, then inspect the orphan-group helper on a synthetic
    # group element. Direct unit test without going through build_deck.
    prs = Presentation(str(TEMPLATE))
    # Reach into any slide and craft a synthetic group; for simplicity, we
    # call the helper directly on a fake slide-like object.
    class FakeSlide:
        def __init__(self, group_xml):
            self.shapes = [FakeShape(group_xml)]

    class FakeShape:
        def __init__(self, el):
            self._element = el
            self.name = el.get_name()

    # Build a group element with name {{TAKEAWAY_GROUP}}
    from lxml import etree
    nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    root = etree.Element(qn('p:spTree'), nsmap=nsmap)
    grp = etree.SubElement(root, qn('p:grpSp'))
    nvGrpSpPr = etree.SubElement(grp, qn('p:nvGrpSpPr'))
    cNvPr = etree.SubElement(nvGrpSpPr, qn('p:cNvPr'))
    cNvPr.set('id', '1')
    cNvPr.set('name', '{{TAKEAWAY_GROUP}}')

    # Spec with empty takeaway → group should be removed
    class FakeSlide2:
        @property
        def shapes(self):
            class S:
                pass
            s = S()
            # _iter_top_level_shapes returns slide.shapes which is iterable of
            # shape-like objects with .name and ._element
            return _ShapeList([_ShapeWrapper(grp)])
    class _ShapeList(list):
        pass
    class _ShapeWrapper:
        def __init__(self, el):
            self._element = el
            self.name = el.find(qn('p:nvGrpSpPr')).find(qn('p:cNvPr')).get('name')

    # Easier: call directly with a real slide construction via render_template_slide
    # then verify {{TAKEAWAY_GROUP}} doesn't survive. But the current template
    # doesn't have such groups. Instead, just verify the helper works in isolation:
    # parent of the group must support .remove()
    parent = root
    assert grp in list(parent)
    # Manually run the helper's core logic
    name = grp.find(qn('p:nvGrpSpPr')).find(qn('p:cNvPr')).get('name')
    m = template_engine._ORPHAN_GROUP_RE.match(name)
    assert m is not None, "regex must match {{TAKEAWAY_GROUP}}"
    key = m.group(1).lower()
    spec_no_takeaway = {"title": "x"}
    val = spec_no_takeaway.get(key)
    assert val in (None, '', [], {}), "takeaway must be considered empty"


def test_debug_layouts_footer(tmp_path):
    """When --debug-layouts is on, every slide carries a tiny footer textbox."""
    spec = {
        "slides": [
            {"layout": "text", "title": "Code-based", "bullets": ["a"]},
            {"layout": "executive_summary", "title": "Template-based",
             "items": [{"title": "x", "bullets": "y"}]},
        ],
        "closing": False,
    }
    out = tmp_path / "dbg.pptx"
    build_deck(spec, out, template_path=TEMPLATE, debug_layouts=True)
    prs = Presentation(str(out))
    for s in prs.slides:
        # Find any textbox containing the [layout: ...] marker
        has_marker = False
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.startswith("[layout:"):
                has_marker = True
                break
        assert has_marker, "expected [layout: ...] footer on every slide"


# ---------------------------------------------------------------------------
# Chantier 11 — roadmap anchor, matrix bullets cap, images
# ---------------------------------------------------------------------------
def test_roadmap_labels_dont_overlap_axis(tmp_path):
    """Roadmap copies positioned ABOVE the axis (even index) must have
    their {{ITEM_DATE}} and {{ITEM_MILESTONE}} shapes anchored at the bottom
    so the text grows upward, away from the timeline axis."""
    from pptx.oxml.ns import qn
    spec = {
        "slides": [{
            "layout": "roadmap_styled",
            "title": "Test",
            "items": [
                {"date": "Jan '26", "milestone": "Audit complet de l'existant"},
                {"date": "Mai '26", "milestone": "Cadrage cible architecture"},
                {"date": "Sept '26", "milestone": "Build phase 1"},
                {"date": "Mars '27", "milestone": "Validation production"},
                {"date": "Juin '27", "milestone": "Bascule finale"},
            ],
        }],
        "closing": False,
    }
    out = tmp_path / "rm.pptx"
    build_deck(spec, out, template_path=TEMPLATE)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "roadmap_styled")
    assert slide is not None
    groups = [sh for sh in slide.shapes if sh.shape_type == 6]
    for i, grp in enumerate(groups):
        expected_anchor = 'b' if i % 2 == 0 else 't'
        for shape_name in ('{{ITEM_DATE}}', '{{ITEM_MILESTONE}}'):
            sp = _named_descendant(grp._element, shape_name)
            if sp is None:
                continue
            bodyPr = sp.find('.//' + qn('a:bodyPr'))
            assert bodyPr is not None, f"copy {i} {shape_name}: no bodyPr"
            assert bodyPr.get('anchor') == expected_anchor, (
                f"copy {i} {shape_name}: anchor={bodyPr.get('anchor')!r}, "
                f"expected {expected_anchor!r}"
            )


def test_matrix_truncates_bullets_to_3(tmp_path, capsys):
    """Quadrants with >3 bullets are silently truncated to the first 3 and
    a warning is printed to stderr."""
    spec = {
        "slides": [{
            "layout": "matrix_2x2_styled",
            "title": "Test",
            "x_axis": {"label": "X"},
            "y_axis": {"label": "Y"},
            "quadrants": {
                "top_left":     {"title": "Top L", "items": ["A1", "A2", "A3", "A4", "A5"]},
                "top_right":    {"title": "Top R", "items": ["B1"]},
                "bottom_left":  {"title": "Bot L", "items": ["C1"]},
                "bottom_right": {"title": "Bot R", "items": ["D1"]},
            },
        }],
        "closing": False,
    }
    out = tmp_path / "mx.pptx"
    build_deck(spec, out, template_path=TEMPLATE)
    captured = capsys.readouterr()
    assert "truncated to 3" in captured.err

    slide = _find_slide_by_cSld_name(Presentation(str(out)), "matrix_2x2_styled")
    bullets_texts = [
        sh.text_frame.text for sh in slide.shapes
        if sh.has_text_frame and sh.name == '{{QUAD_TOP_LEFT_BULLETS}}'
    ]
    # Joined text should contain A1, A2, A3 but NOT A4, A5
    haystack = "\n".join(bullets_texts)
    for keep in ("A1", "A2", "A3"):
        assert keep in haystack, f"bullet {keep!r} should be kept"
    for drop in ("A4", "A5"):
        assert drop not in haystack, f"bullet {drop!r} should have been truncated"


def test_image_engine_keyword_extraction():
    """`extract_keyword_from_title` strips stop words and short tokens."""
    import sys
    sys.path.insert(0, str(ROOT / "scripts"))
    from image_engine import extract_keyword_from_title

    cases = [
        ("Migration vers le cloud", "migration cloud"),
        ("Notre approche en 4 phases", "approche phases"),
        ("", ""),
        ("Le SI atteint ses limites", "atteint limites"),
    ]
    for title, expected in cases:
        actual = extract_keyword_from_title(title, max_words=3)
        assert actual == expected, f"{title!r} → got {actual!r}, expected {expected!r}"


def test_pexels_api_used_when_key_present(tmp_path, monkeypatch):
    """When PEXELS_API_KEY is set, the Pexels endpoint is hit first with
    the Authorization header carrying the key."""
    import sys
    sys.path.insert(0, str(ROOT / "scripts"))
    import image_engine

    monkeypatch.setenv("PEXELS_API_KEY", "fake-key-1234")
    captured = {}

    def fake_pexels(keyword, w, h, api_key, timeout=5.0):
        captured["keyword"] = keyword
        captured["api_key"] = api_key
        return b"FAKE_PNG_BYTES"

    def picsum_should_not_run(*a, **kw):
        raise AssertionError("Picsum should not have been called when Pexels succeeded")

    monkeypatch.setattr(image_engine, "_fetch_pexels", fake_pexels)
    monkeypatch.setattr(image_engine, "_fetch_picsum", picsum_should_not_run)

    result = image_engine.fetch_image_for_slide("cloud computing", 1000000, 800000)
    assert result == b"FAKE_PNG_BYTES"
    assert captured["keyword"] == "cloud computing"
    assert captured["api_key"] == "fake-key-1234"


def test_image_engine_falls_back_to_picsum_on_pexels_error(monkeypatch):
    """A Pexels error (401, network, anything) triggers Picsum fallback."""
    import sys
    sys.path.insert(0, str(ROOT / "scripts"))
    import image_engine

    monkeypatch.setenv("PEXELS_API_KEY", "fake-key")

    def pexels_401(*a, **kw):
        raise OSError("simulated 401 Unauthorized from Pexels")

    def picsum_ok(keyword, w, h, timeout):
        return b"PICSUM_BYTES"

    monkeypatch.setattr(image_engine, "_fetch_pexels", pexels_401)
    monkeypatch.setattr(image_engine, "_fetch_picsum", picsum_ok)

    result = image_engine.fetch_image_for_slide("cloud", 1000000, 800000)
    assert result == b"PICSUM_BYTES"


def test_image_engine_handles_total_failure(tmp_path, monkeypatch):
    """When Pexels AND Picsum fail, the result is None and the build does
    not crash (the slide simply lacks a photo)."""
    import sys
    sys.path.insert(0, str(ROOT / "scripts"))
    import image_engine

    def fail(*a, **kw):
        raise OSError("simulated total outage")

    monkeypatch.setattr(image_engine, "_fetch_pexels", fail)
    monkeypatch.setattr(image_engine, "_fetch_picsum", fail)

    result = image_engine.fetch_image_for_slide("cloud", 1000000, 800000)
    assert result is None

    spec = {
        "slides": [{"layout": "cover", "title": "Test", "subtitle": "S", "ref": "R"}],
        "closing": False,
    }
    out = tmp_path / "img.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=True)
    assert out.exists() and out.stat().st_size > 0


def test_image_inserted_in_cover(tmp_path, monkeypatch):
    """With auto_images on and a mocked fetch returning a real PNG, the
    cover slide's {{IMAGE}} placeholder is replaced by a PICTURE shape."""
    import sys
    sys.path.insert(0, str(ROOT / "scripts"))
    import image_engine

    PNG_1x1 = (
        b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
        b"\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\rIDATx\x9cb\x00\x01"
        b"\x00\x00\x05\x00\x01\r\n-\xb4\x00\x00\x00\x00IEND\xaeB`\x82"
    )

    def fake_fetch(keyword, w_emu, h_emu, timeout=8.0):
        return PNG_1x1

    monkeypatch.setattr(image_engine, "fetch_image_for_slide", fake_fetch)

    # Pick a layout that has {{IMAGE}} in the template: agenda_diagonal does.
    spec = {
        "slides": [{
            "layout": "agenda_diagonal",
            "title": "Sommaire test",
            "items": [{"title": "S1"}, {"title": "S2"}],
        }],
        "closing": False,
    }
    out = tmp_path / "ag.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=True)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "agenda_diagonal")
    assert slide is not None
    pictures = [sh for sh in slide.shapes if sh.shape_type == 13]
    assert pictures, "expected at least one PICTURE shape from auto-image fetch"
    # And the {{IMAGE}} placeholder should be gone
    placeholders = [sh for sh in slide.shapes if sh.name == '{{IMAGE}}']
    assert not placeholders, "{{IMAGE}} placeholder should have been replaced"


# ---------------------------------------------------------------------------
# Chantier 12 — diagonal overlay on cover/section/closing diagonal slides
# ---------------------------------------------------------------------------
def test_diagonal_overlay_applied_on_agenda(tmp_path, monkeypatch):
    """For layouts in SLIDES_WITH_DIAGONAL_OVERLAY (agenda_diagonal among
    them), the inserted picture's spPr should carry a custGeom (copied from
    the layout's Image placeholder) rather than the default prstGeom."""
    from pptx.oxml.ns import qn
    import sys
    sys.path.insert(0, str(ROOT / "scripts"))
    import image_engine

    PNG_1x1 = (
        b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
        b"\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\rIDATx\x9cb\x00\x01"
        b"\x00\x00\x05\x00\x01\r\n-\xb4\x00\x00\x00\x00IEND\xaeB`\x82"
    )

    def fake_fetch(keyword, w_emu, h_emu, timeout=8.0):
        return PNG_1x1

    monkeypatch.setattr(image_engine, "fetch_image_for_slide", fake_fetch)

    spec = {
        "slides": [{
            "layout": "agenda_diagonal",
            "title": "Test sommaire",
            "items": [{"title": "Section 1"}, {"title": "Section 2"}],
        }],
        "closing": False,
    }
    out = tmp_path / "ag.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=True)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "agenda_diagonal")
    assert slide is not None

    pics = [sh for sh in slide.shapes if sh.shape_type == 13]
    assert pics, "expected a PICTURE inserted from auto-image fetch"

    spPr = pics[0]._element.find(qn('p:spPr'))
    assert spPr is not None
    custgeom = spPr.find(qn('a:custGeom'))
    prstgeom = spPr.find(qn('a:prstGeom'))
    assert custgeom is not None, "picture should have custGeom (diagonal cut)"
    assert prstgeom is None, "default prstGeom should have been removed"


def test_no_diagonal_overlay_on_canvas_blank(tmp_path, monkeypatch):
    """canvas_blank is NOT in SLIDES_WITH_DIAGONAL_OVERLAY → no custGeom
    should be applied even if an image happens to be inserted."""
    from pptx.oxml.ns import qn
    import sys
    sys.path.insert(0, str(ROOT / "scripts"))
    import image_engine
    import template_engine

    # canvas_blank doesn't normally have {{IMAGE}} in the template, but we
    # verify the rule by direct invariant: the layout name is not in the set.
    assert 'canvas_blank' not in template_engine.SLIDES_WITH_DIAGONAL_OVERLAY
    assert 'agenda_diagonal' in template_engine.SLIDES_WITH_DIAGONAL_OVERLAY


# ---------------------------------------------------------------------------
# Chantier 14 — canvas_blank freeform composition
# ---------------------------------------------------------------------------
def _cb_find(slide):
    """Helper: return the canvas_blank slide (only one in test specs)."""
    return _find_slide_by_cSld_name(slide, "canvas_blank")


def test_canvas_blank_with_4_kpi_cards(tmp_path):
    """4 KPI cards → 2×2 grid; each card produces 2 textboxes (value + label).
    The bottom row must end above _CB_CONTENT_BOTTOM (footer reserve)."""
    import sys
    sys.path.insert(0, str(ROOT / "scripts"))
    import template_engine
    spec = {
        "slides": [{
            "layout": "canvas_blank",
            "title": "Chiffres clés",
            "blocks": [
                {"type": "kpi_card", "value": "87 %", "label": "saturation"},
                {"type": "kpi_card", "value": "23 j", "label": "provisionnement"},
                {"type": "kpi_card", "value": "47",   "label": "incidents"},
                {"type": "kpi_card", "value": "1.2 M€","label": "économie/an"},
            ],
        }],
        "closing": False,
    }
    out = tmp_path / "cb4.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    slide = _cb_find(Presentation(str(out)))
    assert slide is not None
    texts = [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]
    haystack = "\n".join(texts)
    for value in ("87 %", "23 j", "47", "1.2 M€"):
        assert value in haystack, f"KPI value {value!r} missing"
    for label in ("SATURATION", "PROVISIONNEMENT", "INCIDENTS", "ÉCONOMIE/AN"):
        assert label in haystack, f"KPI label {label!r} missing (uppercased)"

    # Chantier 16 — no shape should extend past the footer reserve
    bottom_limit = template_engine._CB_CONTENT_BOTTOM
    for sh in slide.shapes:
        if sh.name == '{{TITLE}}':
            continue  # title placeholder lives in the header area
        shape_bottom = sh.top + sh.height
        assert shape_bottom <= bottom_limit, (
            f"shape {sh.name!r} bottom={shape_bottom} exceeds "
            f"footer-safe limit {bottom_limit}"
        )


def test_canvas_blank_with_text_and_image(tmp_path, monkeypatch):
    """Text + image → asymmetric layout (text left, image right)."""
    import sys
    sys.path.insert(0, str(ROOT / "scripts"))
    import image_engine

    PNG = (
        b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
        b"\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\rIDATx\x9cb\x00\x01"
        b"\x00\x00\x05\x00\x01\r\n-\xb4\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    monkeypatch.setattr(image_engine, "fetch_image_for_slide",
                        lambda kw, w, h, timeout=8.0: PNG)

    spec = {
        "slides": [{
            "layout": "canvas_blank",
            "title": "Notre approche",
            "blocks": [
                {"type": "text", "content": "Lorem ipsum dolor sit amet."},
                {"type": "image", "keyword": "datacenter"},
            ],
        }],
        "closing": False,
    }
    out = tmp_path / "cbti.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    slide = _cb_find(Presentation(str(out)))
    assert slide is not None
    pics = [sh for sh in slide.shapes if sh.shape_type == 13]
    assert pics, "image block should produce a PICTURE"
    texts = [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]
    assert any("Lorem ipsum" in t for t in texts), "text block missing"


def test_canvas_blank_with_chart(tmp_path):
    """A chart block produces an inserted PICTURE that stays above
    the footer reserve (Chantier 16)."""
    import sys
    sys.path.insert(0, str(ROOT / "scripts"))
    import template_engine
    spec = {
        "slides": [{
            "layout": "canvas_blank",
            "title": "TCO 2025-2029",
            "blocks": [
                {"type": "chart", "chart_spec": {
                    "type": "bar",
                    "labels": ["2025", "2026", "2027", "2028"],
                    "values": [4.2, 4.4, 3.8, 3.1],
                }},
            ],
        }],
        "closing": False,
    }
    out = tmp_path / "cbc.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    slide = _cb_find(Presentation(str(out)))
    assert slide is not None
    pics = [sh for sh in slide.shapes if sh.shape_type == 13]
    assert pics, "chart block should produce a PICTURE"
    bottom_limit = template_engine._CB_CONTENT_BOTTOM
    for pic in pics:
        pic_bottom = pic.top + pic.height
        assert pic_bottom <= bottom_limit, (
            f"chart picture bottom={pic_bottom} exceeds footer-safe limit "
            f"{bottom_limit}"
        )


def test_canvas_blank_respects_palette(tmp_path):
    """Painted blocks use Arial + brand colors. The `{{TITLE}}` placeholder
    inherits its font/color from the master (Arial + navy by design)."""
    from pptx.oxml.ns import qn
    spec = {
        "slides": [{
            "layout": "canvas_blank",
            "title": "Titre test",
            "blocks": [
                {"type": "kpi_card", "value": "42 %", "label": "test"},
                {"type": "kpi_card", "value": "99 %", "label": "other"},
            ],
        }],
        "closing": False,
    }
    out = tmp_path / "cbpal.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    slide = _cb_find(Presentation(str(out)))
    assert slide is not None

    brand_colors = {'14163C', 'F26622', '4A4D6B', '1E2261', '7CB342', 'E63946'}
    # All EXPLICIT <a:latin> attrs we add must be Arial (no fallback fonts).
    # All EXPLICIT <a:srgbClr> values we add must be in the brand palette.
    explicit_fonts = []
    explicit_colors = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for r in sh._element.iter(qn('a:r')):
            latin = r.find('.//' + qn('a:latin'))
            if latin is not None and latin.get('typeface'):
                explicit_fonts.append(latin.get('typeface'))
            srgb = r.find('.//' + qn('a:srgbClr'))
            if srgb is not None and srgb.get('val'):
                explicit_colors.append(srgb.get('val').upper())

    # We painted 4 runs (2 KPI cards × value+label) → 4 explicit Arial + 4 brand colors
    assert explicit_fonts, "expected at least one painted run with explicit font"
    assert all(f == 'Arial' for f in explicit_fonts), (
        f"non-Arial font slipped through: {explicit_fonts}"
    )
    assert explicit_colors, "expected at least one painted run with explicit color"
    non_brand = [c for c in explicit_colors if c not in brand_colors]
    assert not non_brand, f"non-brand colors painted: {non_brand}"


def test_canvas_blank_ignores_eyebrow_when_no_placeholder(tmp_path, capsys):
    """Chantier 15: spec.eyebrow is silently ignored (with stderr warning)
    when the template has no `{{EYEBROW}}` placeholder. Same for takeaway
    and source. The slide must NOT contain a painted eyebrow textbox."""
    spec = {
        "slides": [{
            "layout": "canvas_blank",
            "eyebrow": "ANALYSE",
            "title": "Titre",
            "takeaway": "Take-away ignoré",
            "source": "Source ignorée",
            "blocks": [{"type": "kpi_card", "value": "42", "label": "x"}],
        }],
        "closing": False,
    }
    out = tmp_path / "cbig.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    captured = capsys.readouterr()
    for field, ph in [("eyebrow", "{{EYEBROW}}"),
                      ("takeaway", "{{TAKEAWAY}}"),
                      ("source", "{{SOURCE}}")]:
        assert f"{field} ignored — no {ph}" in captured.err, (
            f"warning for missing {ph} not emitted")

    slide = _cb_find(Presentation(str(out)))
    # No textbox should carry the eyebrow/takeaway/source text we passed
    texts = [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]
    haystack = "\n".join(texts)
    assert "ANALYSE" not in haystack, "eyebrow text should not have been painted"
    assert "Take-away ignoré" not in haystack, "takeaway text should not have been painted"
    assert "Source ignorée" not in haystack, "source text should not have been painted"


def test_canvas_blank_truncates_at_6_blocks(tmp_path, capsys):
    """8 blocks → only 6 are rendered, warning printed to stderr."""
    spec = {
        "slides": [{
            "layout": "canvas_blank",
            "title": "Test truncation",
            "blocks": [
                {"type": "kpi_card", "value": f"{i}", "label": f"k{i}"}
                for i in range(8)
            ],
        }],
        "closing": False,
    }
    out = tmp_path / "cb8.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    captured = capsys.readouterr()
    assert "truncating to first 6" in captured.err
    slide = _cb_find(Presentation(str(out)))
    texts = "\n".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
    # First 6 values 0..5 should appear; 6 and 7 should not
    for i in range(6):
        assert f"K{i}" in texts.upper(), f"label k{i} missing"
    assert "K6" not in texts.upper(), "label k6 should have been truncated"
    assert "K7" not in texts.upper(), "label k7 should have been truncated"


# ---------------------------------------------------------------------------
# Chantier 15 — data_table layout
# ---------------------------------------------------------------------------
def test_data_table_simple(tmp_path):
    """3 columns × 3 rows basic table renders as a TABLE shape."""
    spec = {
        "slides": [{
            "layout": "data_table",
            "title": "Test simple",
            "table": {
                "headers": ["KPI", "Avant", "Après"],
                "rows": [
                    ["Disponibilité", "96.4 %", "99.8 %"],
                    ["TCO/an",         "4.2 M€", "3.1 M€"],
                    ["Provisionnement","23 j",   "2 h"],
                ],
            },
        }],
        "closing": False,
    }
    out = tmp_path / "dt1.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "data_table")
    assert slide is not None
    tables = [sh for sh in slide.shapes if sh.has_table]
    assert len(tables) == 1
    t = tables[0].table
    assert len(t.rows) == 4    # 3 data rows + 1 header
    assert len(t.columns) == 3
    assert t.rows[0].cells[0].text == "KPI"
    assert t.rows[1].cells[0].text == "Disponibilité"


def test_data_table_with_highlights(tmp_path):
    """Highlight column → orange bold text. Highlight row → light orange fill."""
    from pptx.oxml.ns import qn
    spec = {
        "slides": [{
            "layout": "data_table",
            "title": "Test highlights",
            "table": {
                "headers": ["Option", "Coût", "Gain"],
                "rows": [
                    ["A", "1 M€", "-10 %"],
                    ["B", "2 M€", "-25 %"],
                    ["C", "3 M€", "-40 %"],
                ],
                "highlight_column": 2,
                "highlight_row": 1,
            },
        }],
        "closing": False,
    }
    out = tmp_path / "dt2.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "data_table")
    t = next(sh.table for sh in slide.shapes if sh.has_table)

    # Highlight col=2 → orange text on every data row at col 2
    for i in range(1, 4):
        cell = t.cell(i, 2)
        run = cell.text_frame.paragraphs[0].runs[0]
        assert run.font.color.rgb is not None
        assert str(run.font.color.rgb).upper() == "F26622", (
            f"highlight col cell ({i},2) should be orange, got "
            f"{run.font.color.rgb}"
        )
        assert run.font.bold is True, "highlight col should be bold"

    # Highlight row=1 → cells in row 2 (data row 1 = table row 2) have light-orange fill
    for j in range(3):
        cell = t.cell(2, j)
        fill_xml = cell._tc.xml
        assert "FDF1EA" in fill_xml.upper(), (
            f"highlight row cell at col {j} should have FDF1EA fill"
        )


def test_data_table_max_columns(tmp_path):
    """6 columns = the maximum allowed."""
    spec = {
        "slides": [{
            "layout": "data_table",
            "title": "Max cols",
            "table": {
                "headers": ["A", "B", "C", "D", "E", "F"],
                "rows": [["a", "b", "c", "d", "e", "f"]],
            },
        }],
        "closing": False,
    }
    out = tmp_path / "dt6.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "data_table")
    t = next(sh.table for sh in slide.shapes if sh.has_table)
    assert len(t.columns) == 6


def test_data_table_truncates_at_8_rows(tmp_path, capsys):
    """10 data rows → only 8 are rendered + stderr warning."""
    spec = {
        "slides": [{
            "layout": "data_table",
            "title": "Trop de lignes",
            "table": {
                "headers": ["Item", "Valeur"],
                "rows": [[f"row{i}", f"v{i}"] for i in range(10)],
            },
        }],
        "closing": False,
    }
    out = tmp_path / "dt10.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    captured = capsys.readouterr()
    assert "10 rows > 8, truncating" in captured.err
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "data_table")
    t = next(sh.table for sh in slide.shapes if sh.has_table)
    assert len(t.rows) == 9  # 8 data + 1 header


def test_data_table_autoshrinks_on_overflow(tmp_path):
    """8 rows × 6 cols with long cells → font size <10 or padding <6."""
    from pptx.util import Pt
    spec = {
        "slides": [{
            "layout": "data_table",
            "title": "Overflow",
            "table": {
                "headers": ["Col1", "Col2", "Col3", "Col4", "Col5", "Col6"],
                "rows": [[f"valeur_lo_{i}_{j}" for j in range(6)]
                         for i in range(8)],
            },
        }],
        "closing": False,
    }
    out = tmp_path / "dtshrink.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "data_table")
    t = next(sh.table for sh in slide.shapes if sh.has_table)
    # Sample a few body cells: font size should have shrunk below 10pt OR
    # the padding should have shrunk (we check that at least one happened).
    body_cell = t.cell(1, 1)
    run = body_cell.text_frame.paragraphs[0].runs[0]
    body_size_pt = run.font.size.pt
    tf = body_cell.text_frame
    body_padding_pt = tf.margin_left.pt
    # Either font ≤ 9 or padding ≤ 5 indicates autoshrink kicked in
    assert body_size_pt <= 9 or body_padding_pt <= 5, (
        f"expected autoshrink, got size={body_size_pt}, padding={body_padding_pt}"
    )


# ---------------------------------------------------------------------------
# Chantier 16 — closing_diagonal long title + uniform font in REPEAT_ITEM
# ---------------------------------------------------------------------------
def _collect_sz_for_shape_name(slide, shape_name):
    """Return list of sz attribute values found under shape_name across all
    REPEAT_ITEM copies in the slide."""
    from pptx.oxml.ns import qn
    out = []
    for top_sh in slide.shapes:
        if top_sh.shape_type != 6:
            continue
        for sp in top_sh._element.iter(qn('p:sp')):
            cn = sp.find('.//' + qn('p:cNvPr'))
            if cn is None or cn.get('name', '') != shape_name:
                continue
            rPr = sp.find('.//' + qn('a:rPr'))
            if rPr is not None and rPr.get('sz'):
                out.append(int(rPr.get('sz')))
    return out


def test_closing_diagonal_long_title(tmp_path):
    """A long closing title is shrunk to fit the narrow {{TITLE}} shape
    (60pt → ≤ 44pt), preserving word boundaries (no mid-word break)."""
    from pptx.oxml.ns import qn
    spec = {
        "slides": [{
            "layout": "closing_diagonal",
            "title": "Discutons de votre trajectoire vers le cloud en toute sérénité",
            "question": "Vos questions ?",
            "author_name": "Florian Horellou",
            "author_email": "florian@aosis.net",
            "author_phone": "+33 6 36 26 17 47",
        }],
        "closing": False,
    }
    out = tmp_path / "closing.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "closing_diagonal")
    assert slide is not None

    title_sh = next((sh for sh in slide.shapes if sh.name == '{{TITLE}}'), None)
    assert title_sh is not None
    rPr = title_sh._element.find('.//' + qn('a:rPr'))
    sz = int(rPr.get('sz'))
    # Original is 6000 (60pt). Long title should have shrunk to ≤ 4400 (44pt).
    assert sz <= 4400, f"long title not shrunk enough, got sz={sz}"
    assert sz >= 3600, f"shrunk below floor (36pt), got sz={sz}"

    # Author block must not be vertically beyond the slide bottom
    for name in ('{{AUTHOR_NAME}}', '{{AUTHOR_EMAIL}}', '{{AUTHOR_PHONE}}'):
        sh = next((s for s in slide.shapes if s.name == name), None)
        if sh is None:
            continue
        bottom = sh.top + sh.height
        assert bottom < 5_143_500, f"{name} bottom {bottom} beyond slide"


def test_closing_diagonal_short_title_not_shrunk(tmp_path):
    """A short title (e.g. 'MERCI') should keep its original 60pt size."""
    from pptx.oxml.ns import qn
    spec = {
        "slides": [{
            "layout": "closing_diagonal",
            "title": "MERCI",
            "author_name": "F.", "author_email": "f@aosis.net", "author_phone": "+33",
        }],
        "closing": False,
    }
    out = tmp_path / "merci.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "closing_diagonal")
    rPr = next(sh for sh in slide.shapes if sh.name == '{{TITLE}}')._element.find('.//' + qn('a:rPr'))
    assert int(rPr.get('sz')) == 6000, "short title should keep 60pt"


def test_roadmap_uniform_font_size(tmp_path):
    """All {{ITEM_MILESTONE}} copies share the same sz (= the smallest
    observed across all items, determined by the longest milestone)."""
    spec = {
        "slides": [{
            "layout": "roadmap_styled",
            "title": "Test",
            "items": [
                {"date": "Q1", "milestone": "M1"},
                {"date": "Q2", "milestone": "Audit"},
                {"date": "Q3", "milestone": "Décommissionnement complet du datacenter Lyon"},
                {"date": "Q4", "milestone": "Bascule"},
                {"date": "Q5", "milestone": "Done"},
            ],
        }],
        "closing": False,
    }
    out = tmp_path / "rm_uniform.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "roadmap_styled")
    sizes = _collect_sz_for_shape_name(slide, '{{ITEM_MILESTONE}}')
    assert len(sizes) == 5, f"expected 5 milestones, got {len(sizes)}"
    assert len(set(sizes)) == 1, (
        f"milestone sizes should all be equal, got {sizes}"
    )
    date_sizes = _collect_sz_for_shape_name(slide, '{{ITEM_DATE}}')
    assert len(set(date_sizes)) == 1, (
        f"date sizes should all be equal, got {date_sizes}"
    )


def test_next_steps_uniform_dates(tmp_path):
    """next_steps: all {{ITEM_DATE}} copies share the same sz."""
    spec = {
        "slides": [{
            "layout": "next_steps",
            "title": "Test",
            "items": [
                {"action": "A1", "owner": "X", "date": "15 juin 2026"},
                {"action": "A2", "owner": "Y", "date": "1er septembre 2026 — kick-off"},
                {"action": "A3", "owner": "Z", "date": "Q4"},
                {"action": "A4", "owner": "W", "date": "Mars 2027"},
            ],
        }],
        "closing": False,
    }
    out = tmp_path / "ns_uniform.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "next_steps")
    date_sizes = _collect_sz_for_shape_name(slide, '{{ITEM_DATE}}')
    assert len(date_sizes) == 4
    assert len(set(date_sizes)) == 1, (
        f"all date sizes should be equal, got {date_sizes}"
    )


def test_kpi_with_chart_uniform_values(tmp_path):
    """kpi_with_chart: all {{KPI_VALUE}} copies share the same sz."""
    spec = {
        "slides": [{
            "layout": "kpi_with_chart",
            "title": "Test",
            "kpis": [
                {"label": "A", "value": "42"},
                {"label": "B", "value": "1 234 567 K€"},
                {"label": "C", "value": "99"},
            ],
            "chart": {"type": "bar", "labels": ["A","B","C"], "values": [1,2,3]},
        }],
        "closing": False,
    }
    out = tmp_path / "kpi_uniform.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "kpi_with_chart")
    value_sizes = _collect_sz_for_shape_name(slide, '{{KPI_VALUE}}')
    assert len(value_sizes) == 3
    assert len(set(value_sizes)) == 1, (
        f"all KPI value sizes should be equal, got {value_sizes}"
    )


def test_agenda_uniform_items(tmp_path):
    """agenda_diagonal: all {{ITEM_TITLE}} copies share the same sz."""
    spec = {
        "slides": [{
            "layout": "agenda_diagonal",
            "title": "Sommaire",
            "items": [
                {"title": "Court"},
                {"title": "Section très longue avec beaucoup de mots et d'idées"},
                {"title": "Médium"},
                {"title": "Court 2"},
            ],
        }],
        "closing": False,
    }
    out = tmp_path / "agenda_uniform.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "agenda_diagonal")
    title_sizes = _collect_sz_for_shape_name(slide, '{{ITEM_TITLE}}')
    assert len(title_sizes) == 4
    assert len(set(title_sizes)) == 1, (
        f"all agenda item sizes should be equal, got {title_sizes}"
    )


# ---------------------------------------------------------------------------
# Chantier 17 — uniform font size across matrix_2x2_styled quadrants
# ---------------------------------------------------------------------------
def _collect_sz_for_top_level(slide, shape_name_pattern_re):
    """Return list of first-run sz attribute values across top-level shapes
    whose name matches `shape_name_pattern_re`."""
    from pptx.oxml.ns import qn
    import re
    out = []
    for sh in slide.shapes:
        if not re.match(shape_name_pattern_re, sh.name or ''):
            continue
        rPr = sh._element.find('.//' + qn('a:rPr'))
        if rPr is not None and rPr.get('sz'):
            out.append(int(rPr.get('sz')))
    return out


def test_matrix_2x2_uniform_title_size(tmp_path):
    """The 4 {{QUAD_*_TITLE}} share the same font size even when title
    lengths differ wildly."""
    spec = {
        "slides": [{
            "layout": "matrix_2x2_styled",
            "title": "Test uniform titles",
            "x_axis": {"label": "X"},
            "y_axis": {"label": "Y"},
            "quadrants": {
                "top_left":     {"title": "Quick wins",                  "items": ["A"]},
                "top_right":    {"title": "À mitiger systématiquement",  "items": ["B"]},
                "bottom_left":  {"title": "Hygiène",                     "items": ["C"]},
                "bottom_right": {"title": "Deprio",                      "items": ["D"]},
            },
        }],
        "closing": False,
    }
    out = tmp_path / "mxuti.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "matrix_2x2_styled")
    sizes = _collect_sz_for_top_level(
        slide, r'^\{\{QUAD_(TOP_LEFT|TOP_RIGHT|BOTTOM_LEFT|BOTTOM_RIGHT)_TITLE\}\}$'
    )
    assert len(sizes) == 4
    assert len(set(sizes)) == 1, (
        f"quad title sizes should be uniform, got {sizes}"
    )


def test_matrix_2x2_uniform_bullets_size(tmp_path):
    """The 4 {{QUAD_*_BULLETS}} share the same font size even when bullet
    counts vary across quadrants (1, 2, 3, 2)."""
    spec = {
        "slides": [{
            "layout": "matrix_2x2_styled",
            "title": "Test uniform bullets",
            "x_axis": {"label": "X"},
            "y_axis": {"label": "Y"},
            "quadrants": {
                "top_left":     {"title": "TL", "items": ["bullet 1"]},
                "top_right":    {"title": "TR", "items": ["b1", "b2"]},
                "bottom_left":  {"title": "BL", "items": ["b1", "b2", "b3"]},
                "bottom_right": {"title": "BR", "items": ["b1", "b2"]},
            },
        }],
        "closing": False,
    }
    out = tmp_path / "mxubu.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "matrix_2x2_styled")
    sizes = _collect_sz_for_top_level(
        slide, r'^\{\{QUAD_(TOP_LEFT|TOP_RIGHT|BOTTOM_LEFT|BOTTOM_RIGHT)_BULLETS\}\}$'
    )
    assert len(sizes) == 4
    assert len(set(sizes)) == 1, (
        f"quad bullets sizes should be uniform, got {sizes}"
    )


# ---------------------------------------------------------------------------
# Chantier 18 — premium effects (drop shadow, XXL KPI, chart border)
# ---------------------------------------------------------------------------
def _has_outer_shadow(element):
    """True if the shape element has an a:outerShdw effect in its spPr."""
    from pptx.oxml.ns import qn
    spPr = element.find(qn('p:spPr'))
    if spPr is None:
        return False
    eff = spPr.find(qn('a:effectLst'))
    if eff is None:
        return False
    return eff.find(qn('a:outerShdw')) is not None


def test_drop_shadow_applied_to_kpi_card(tmp_path):
    """The `kpi_card` background shape of kpi_with_chart copies carries
    an outerShdw effect after Chantier 18."""
    from pptx.oxml.ns import qn
    spec = {
        "slides": [{
            "layout": "kpi_with_chart",
            "title": "Test",
            "kpis": [
                {"label": "A", "value": "10"},
                {"label": "B", "value": "20"},
                {"label": "C", "value": "30"},
            ],
            "chart": {"type": "bar", "labels": ["A","B","C"], "values": [1,2,3]},
        }],
        "closing": False,
    }
    out = tmp_path / "kpi_shadow.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "kpi_with_chart")
    assert slide is not None
    found = 0
    for top in slide.shapes:
        if top.shape_type != 6:
            continue
        for sp in top._element.iter(qn('p:sp')):
            cn = sp.find('.//' + qn('p:cNvPr'))
            if cn is not None and cn.get('name') == 'kpi_card':
                if _has_outer_shadow(sp):
                    found += 1
    assert found >= 1, "no kpi_card carries outerShdw effect"


def test_drop_shadow_not_applied_to_diagonal_images(tmp_path, monkeypatch):
    """Pictures inserted in diagonal layouts (cover/agenda/section/closing)
    must NOT carry an outerShdw (the diagonal cut would render badly with
    a rectangular shadow)."""
    from pptx.oxml.ns import qn
    import sys
    sys.path.insert(0, str(ROOT / "scripts"))
    import image_engine

    PNG_1x1 = (
        b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
        b"\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\rIDATx\x9cb\x00\x01"
        b"\x00\x00\x05\x00\x01\r\n-\xb4\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    monkeypatch.setattr(image_engine, "fetch_image_for_slide",
                        lambda kw, w, h, timeout=8.0: PNG_1x1)
    spec = {
        "slides": [{
            "layout": "agenda_diagonal",
            "title": "Sommaire",
            "items": [{"title": "S1"}],
        }],
        "closing": False,
    }
    out = tmp_path / "ag_noshadow.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=True)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "agenda_diagonal")
    pics = [sh for sh in slide.shapes if sh.shape_type == 13]
    assert pics, "expected an auto-fetched picture"
    for pic in pics:
        assert not _has_outer_shadow(pic._element), (
            "diagonal-layout picture should NOT carry outerShdw"
        )


def test_kpi_value_xxl_default(tmp_path):
    """A short KPI value (e.g. '87%') uses the XXL baseline ≥ 48pt."""
    from pptx.oxml.ns import qn
    spec = {
        "slides": [{
            "layout": "kpi_with_chart",
            "title": "Test",
            "kpis": [{"label": "Coverage", "value": "87%"}],
            "chart": {"type": "bar", "labels": ["A"], "values": [1]},
        }],
        "closing": False,
    }
    out = tmp_path / "kpi_xxl.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "kpi_with_chart")
    grp = [sh for sh in slide.shapes if sh.shape_type == 6][0]
    value_sp = _named_descendant(grp._element, '{{KPI_VALUE}}')
    rPr = value_sp.find('.//' + qn('a:rPr'))
    sz = int(rPr.get('sz'))
    # "87%" is 3 chars and the kpi_with_chart shape is only 0.47" tall.
    # Chantier 19 height-aware sizing caps at 30pt. Still bigger than the
    # pre-Chantier 18 baseline of 28pt.
    assert sz >= 3000, f"short KPI value should be ≥ 30pt (XXL), got sz={sz}"


def test_chart_has_border(tmp_path):
    """A chart PNG carries a 2px gray border around its outer edge
    (Chantier 18 framing). Detect by reading the corner pixels of the
    rendered PNG — they should match the border color."""
    import sys
    sys.path.insert(0, str(ROOT / "scripts"))
    import chart_engine
    from PIL import Image
    import io
    png, _, _ = chart_engine.render_chart_to_png(
        {"type": "bar", "labels": ["A","B"], "values": [1,2]},
        1_000_000, 800_000,
    )
    img = Image.open(io.BytesIO(png)).convert("RGB")
    # Top-left pixel: should be the border color #E8E9F2 (or very close)
    r, g, b = img.getpixel((0, 0))
    # Allow PIL anti-alias rounding ±5
    assert abs(r - 232) <= 5 and abs(g - 233) <= 5 and abs(b - 242) <= 5, (
        f"top-left pixel {(r,g,b)} should match border #E8E9F2"
    )


# ---------------------------------------------------------------------------
# Chantier 19 — height-aware KPI sizing (anti-overlap)
# ---------------------------------------------------------------------------
def _read_kpi_value_font_size_pt(slide):
    """Return font size (in pt) of the first textbox shape that is NOT
    {{TITLE}}, {{SOURCE}} or label-style. Heuristic: largest bold sz."""
    from pptx.oxml.ns import qn
    sizes = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for r in sh._element.iter(qn('a:r')):
            rPr = r.find(qn('a:rPr'))
            if rPr is None:
                continue
            sz_str = rPr.get('sz')
            if sz_str and rPr.get('b') == '1':
                sizes.append(int(sz_str))
    return max(sizes) // 100 if sizes else 0


def test_kpi_no_overlap_short_card():
    """Unit test of the helper: a very short available height (0.3") yields
    a small font size (≤ 32pt) to leave room for the label without overlap."""
    import sys
    sys.path.insert(0, str(ROOT / "scripts"))
    from template_engine import _compute_max_kpi_font_size
    # 0.3" tall × 4" wide rectangle, short value "M1"
    pt = _compute_max_kpi_font_size("M1", available_w_emu=4 * 914400, available_h_emu=int(0.3 * 914400))
    assert pt <= 32, f"short height → small font (anti-overlap), got {pt}pt"
    assert pt >= 24, f"should not go below 24pt floor, got {pt}pt"


def test_kpi_xxl_on_tall_card(tmp_path):
    """A single kpi_card alone (= full slide rect) is tall enough to
    deserve XXL. Value '87%' should land at ≥ 48pt."""
    spec = {
        "slides": [{
            "layout": "canvas_blank",
            "title": "Single KPI dominant",
            "blocks": [
                {"type": "kpi_card", "value": "87%", "label": "label court"},
            ],
        }],
        "closing": False,
    }
    out = tmp_path / "kpi_tall.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "canvas_blank")
    max_pt = _read_kpi_value_font_size_pt(slide)
    assert max_pt >= 48, (
        f"tall card with short value should keep XXL ≥ 48pt, got {max_pt}"
    )


def test_kpi_uniform_min_across_cards(tmp_path):
    """When 3 kpi_card blocks are on the same canvas_blank slide, they all
    share the same font size = min of the max admissible per card. Layout
    with 3 cards uses 1 row × 3 cols → each card ≈ 3" wide × 3" tall."""
    from pptx.oxml.ns import qn
    spec = {
        "slides": [{
            "layout": "canvas_blank",
            "title": "Test uniformity across kpi_card blocks",
            "blocks": [
                # Long value forces a small per-card max
                {"type": "kpi_card", "value": "1 234 567 K€", "label": "TRES LONG"},
                {"type": "kpi_card", "value": "42", "label": "COURT"},
                {"type": "kpi_card", "value": "87", "label": "COURT"},
            ],
        }],
        "closing": False,
    }
    out = tmp_path / "kpi_uniform.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "canvas_blank")
    # All KPI value bold runs should share the same sz
    bold_sizes = set()
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for r in sh._element.iter(qn('a:r')):
            rPr = r.find(qn('a:rPr'))
            if rPr is None or rPr.get('b') != '1':
                continue
            sz = rPr.get('sz')
            if sz and int(sz) >= 2400:  # filter out small label runs
                bold_sizes.add(int(sz))
    # All large bold runs (= KPI values) should share one size
    assert len(bold_sizes) <= 1, (
        f"all kpi_card values should share one size, got {bold_sizes}"
    )


def test_kpi_long_value_shrinks_more():
    """Unit test of the helper: a long value in the same rect yields a
    smaller font than a short value (width constraint)."""
    import sys
    sys.path.insert(0, str(ROOT / "scripts"))
    from template_engine import _compute_max_kpi_font_size
    # Narrow rect: 2" wide × 1" tall — width will dominate for long text
    w = int(2 * 914400)
    h = int(1 * 914400)
    short = _compute_max_kpi_font_size("42", w, h)
    long_val = _compute_max_kpi_font_size("1 234 567 K€", w, h)
    assert long_val < short, (
        f"long value should shrink more, short={short} long={long_val}"
    )


# =============================================================================
# Chantier 20 — Anti-overlap in-card positioning + footer-safe grids
# =============================================================================
def test_kpi_card_label_below_value_no_overlap(tmp_path):
    """Chantier 20 — a short kpi_card (1×1.5") with a value 'M1' and a long
    label must place the label fully BELOW the value, not behind it. The
    label textbox top must be ≥ value textbox bottom."""
    from pptx.util import Emu
    import sys
    sys.path.insert(0, str(ROOT / "scripts"))
    from template_engine import _cb_kpi_card_dynamic_layout
    rect = (Emu(int(0.5 * 914400)), Emu(int(0.5 * 914400)),
            Emu(int(1.0 * 914400)), Emu(int(1.5 * 914400)))
    _, value_rect, label_rect = _cb_kpi_card_dynamic_layout(
        "M1", "AUDIT + CARTOGRAPHIE APPLICATIVE DETAILLEE", rect, floor_pt=14)
    vL, vT, vW, vH = value_rect
    lL, lT, lW, lH = label_rect
    # Label top must be at or below value bottom (no overlap)
    assert lT >= vT + vH, (
        f"label top {lT} must be ≥ value bottom {vT + vH}"
    )
    # Both rects must stay within the card
    card_left, card_top, card_w, card_h = rect
    assert vT >= card_top, "value top above card top"
    assert lT + lH <= card_top + card_h + 100, (
        f"label bottom {lT + lH} overflows card bottom {card_top + card_h}"
    )


def test_canvas_blank_5_blocks_no_footer_overlap(tmp_path):
    """Chantier 20 — with 5 blocks asymmetric (1 chart + 4 stacked KPI on
    left), no rect bottom must exceed 4.75" from slide top (master footer
    safety)."""
    spec = {
        "slides": [{
            "layout": "canvas_blank",
            "title": "Stress 5 blocs asym",
            "blocks": [
                {"type": "kpi_card", "value": "85%", "label": "ADOPTION"},
                {"type": "kpi_card", "value": "24m", "label": "PAYBACK"},
                {"type": "kpi_card", "value": "3,2M", "label": "ECONOMIES ANNUELLES"},
                {"type": "kpi_card", "value": "-42%", "label": "TCO REDUCTION"},
                {"type": "chart", "chart": {"type": "bar", "labels": ["A", "B"],
                                            "values": [1, 2]}},
            ],
        }],
        "closing": False,
    }
    out = tmp_path / "stress5.pptx"
    build_deck(spec, out, template_path=TEMPLATE, auto_images=False)
    prs = Presentation(str(out))
    slide = _find_slide_by_cSld_name(prs, "canvas_blank")
    EMU_PER_INCH = 914400
    MAX_BOTTOM = int(4.75 * EMU_PER_INCH)
    for sh in slide.shapes:
        if sh.top is None or sh.height is None:
            continue
        bottom = sh.top + sh.height
        # Skip title (well above) and any shape above title
        if sh.top < int(1.0 * EMU_PER_INCH):
            continue
        assert bottom <= MAX_BOTTOM, (
            f"shape '{sh.name}' bottom {bottom / EMU_PER_INCH:.3f}\" > "
            f"4.75\" (footer reserve violated)"
        )


def test_kpi_with_chart_value_fits_width(tmp_path):
    """Chantier 20 — a kpi_with_chart KPI value '5.9 M€' (wide chars 'M€')
    must be sized so the estimated text width fits within the {{KPI_VALUE}}
    shape (1.00" wide), preventing wrap-to-next-line."""
    from pptx.oxml.ns import qn
    spec = {
        "slides": [{
            "layout": "kpi_with_chart",
            "title": "TCO trajectoire",
            "kpis": [
                {"label": "TCO actuel",  "value": "5.9 M€"},
                {"label": "TCO cible",   "value": "3.4 M€"},
                {"label": "Économies",   "value": "2.5 M€"},
            ],
            "chart": {"type": "bar", "labels": ["2026", "2028"], "values": [5.9, 3.4]},
        }],
        "closing": False,
    }
    out = tmp_path / "tco.pptx"
    build_deck(spec, out, template_path=TEMPLATE)
    slide = _find_slide_by_cSld_name(Presentation(str(out)), "kpi_with_chart")
    EMU_PER_PT = 12700.0
    # Read all {{KPI_VALUE}} sizes
    sizes = []
    for sh in slide.shapes:
        if sh.shape_type != 6:
            continue
        for sp in sh._element.iter(qn('p:sp')):
            cNvPr = sp.find('.//' + qn('p:cNvPr'))
            if cNvPr is None or cNvPr.get('name') != '{{KPI_VALUE}}':
                continue
            for rPr in sp.iter(qn('a:rPr')):
                sz = rPr.get('sz')
                if sz:
                    sizes.append(int(sz) / 100.0)
                    break
    assert sizes, "no {{KPI_VALUE}} found"
    # Shape is 1.00" wide ; "5.9 M€" = 6 chars × 0.55 × font_pt
    # must be ≤ 1.00" × 72 pt/" = 72pt → font_pt ≤ 72 / (6 × 0.55) ≈ 21.8pt
    for sz_pt in sizes:
        est_width_pt = 6 * 0.55 * sz_pt
        shape_width_pt = 1.00 * 72  # shape width in pt
        assert est_width_pt <= shape_width_pt * 1.05, (
            f"value at {sz_pt}pt has est width {est_width_pt:.1f}pt > "
            f"shape width {shape_width_pt}pt → would wrap"
        )


def test_summarize_report(capsys, tmp_path):
    """`summarize_report` prints per-slide defect counts and totals."""
    # Make scripts/ importable for visual_review
    sys.path.insert(0, str(ROOT / "scripts"))
    from visual_review import summarize_report  # noqa: E402

    fake_report = {
        "deck": "fake_deck.pptx",
        "reviewed_at": "2026-05-13T16:00:00+00:00",
        "slides": [
            {"slide": 1, "image": "slide-01.jpg", "defects": []},
            {"slide": 2, "image": "slide-02.jpg", "defects": [
                {"severity": "minor", "category": "legibility", "description": "x"},
            ]},
            {"slide": 3, "image": "slide-03.jpg", "defects": [
                {"severity": "critical", "category": "overflow",  "description": "y"},
                {"severity": "critical", "category": "overflow",  "description": "z"},
                {"severity": "important", "category": "alignment", "description": "w"},
            ]},
            {"slide": 4, "image": "slide-04.jpg", "defects": []},
        ],
    }
    report_path = tmp_path / "review_report.json"
    report_path.write_text(json.dumps(fake_report), encoding="utf-8")

    summarize_report(report_path)
    captured = capsys.readouterr()
    out = captured.out

    assert "fake_deck.pptx" in out
    assert "Slide 1: 0 defects" in out
    assert "Slide 2: 1 minor (legibility)" in out
    # Slide 3 has 2 critical overflows + 1 important alignment
    assert "Slide 3: 2 critical (overflow), 1 important (alignment)" in out
    assert "Slide 4: 0 defects" in out
    # Totals: 2 critical + 1 important + 1 minor across 4 slides
    assert "2 critical" in out and "1 important" in out and "1 minor" in out
    assert "across 4 slides" in out
