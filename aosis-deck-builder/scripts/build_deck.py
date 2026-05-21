"""
build_deck.py — Generate an AOSIS-branded .pptx from a JSON spec.

Always starts from the official AOSIS template (assets/AOSIS_template.pptx)
so the charte graphique (navy/orange palette, logo, masters, footers) is
inherited by construction. On top of the four template layouts, this script
composes rich visual slide types (stat grids, timelines, cards, comparisons,
auto-generated charts, quotes, etc.) using python-pptx shapes coloured with
the AOSIS palette.

Usage:
    python build_deck.py <spec.json> <output.pptx>

See SKILL.md for the full JSON spec schema and layout catalogue.
"""

import argparse
import json
import os
import sys
import tempfile
from pathlib import Path


def _load_dotenv() -> None:
    """Load environment variables from a `.env` file at the project root.

    Walks up from this script's directory looking for `.env`. If found, its
    `KEY=VALUE` lines are imported into os.environ — without overwriting
    variables already set externally (so explicit `KEY=... python ...` always
    wins). No third-party dependency; minimal parser handling `export `
    prefix, quoted values, and `#` comments.
    """
    here = Path(__file__).resolve()
    for parent in [here.parent, *here.parents]:
        candidate = parent / ".env"
        if candidate.exists():
            try:
                for raw in candidate.read_text(encoding="utf-8").splitlines():
                    line = raw.strip()
                    if not line or line.startswith("#"):
                        continue
                    if line.startswith("export "):
                        line = line[len("export "):].lstrip()
                    if "=" not in line:
                        continue
                    key, _, value = line.partition("=")
                    key = key.strip()
                    value = value.strip()
                    # Strip matching single or double quotes
                    if len(value) >= 2 and value[0] == value[-1] and value[0] in ("'", '"'):
                        value = value[1:-1]
                    if key and key not in os.environ:
                        os.environ[key] = value
            except OSError:
                pass
            return  # Stop at the first .env found


_load_dotenv()

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

from brand import BrandPalette  # noqa: E402  (sibling module in scripts/)
from template_engine import (  # noqa: E402
    discover_template_layouts,
    render_template_slide,
    strip_sample_slides,
)

# Slide geometry (10" × 5.62", 16:9)
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.62)
CONTENT_TOP = Inches(1.1)
CONTENT_BOTTOM = Inches(4.95)
CONTENT_LEFT = Inches(0.4)
CONTENT_RIGHT = Inches(9.6)
CONTENT_W = Inches(9.2)
CONTENT_H = Inches(3.85)

SCRIPT_DIR = Path(__file__).resolve().parent
TEMPLATE_PATH = SCRIPT_DIR.parent / "assets" / "AOSIS_template.pptx"

# Consolidation chantier — AOSIS_template.pptx now hosts BOTH the brand
# (masters, layouts, theme) AND the named template slides (16 layouts as of
# 2026-05-18). The build base is a clone of AOSIS_template.pptx with the
# sample slides stripped so masters/layouts/themes are shared by both
# template-based (cross-deck deepcopy) and code-based renders.
# `exhibits.pptx` is kept in the repo as an archive of the prior state.
EXHIBITS_PATH = TEMPLATE_PATH

# Chantier 3 — palette read dynamically from the active template's theme XML.
BRAND = BrandPalette.from_template(TEMPLATE_PATH)

# Chantier 8 — inject the active palette into chart_engine so matplotlib
# charts in `kpi_with_chart` track the brand. Optional import keeps the
# skill usable without matplotlib (charts then degrade gracefully).
try:
    import chart_engine  # noqa: E402
    chart_engine.set_brand(BRAND)
except ImportError:
    pass

# Inject the active palette into template_engine so layout alternation
# rules (framework_3cards, process_steps, roadmap_styled) resolve their
# semantic color names (orange, navy_alt, …) against the active theme.
import template_engine as _template_engine  # noqa: E402
if hasattr(_template_engine, "set_brand"):
    _template_engine.set_brand(BRAND)

# Code-based layouts use these LAYOUT_MAP indices; resolution is by name in
# AOSIS_template.pptx (4 layouts across 2 masters).
LAYOUT_MAP = {
    "cover":   (0, 0),   # 'Cover'           (master 0)
    "section": (0, 1),   # 'Closing'         (master 0)
    "closing": (0, 1),   # 'Closing'         (master 0)
    "content": (1, 0),   # 'Contenu + texte' (master 1)
    "text":    (1, 1),   # 'Texte'           (master 1)
}

# Discover template-based layouts dynamically at module load. Refreshed when
# build_deck() reloads the palette if a custom --template is used.
try:
    TEMPLATE_BASED_LAYOUTS = set(discover_template_layouts(EXHIBITS_PATH).keys())
except Exception:
    # If exhibits.pptx is missing or malformed, fall back to empty set — the
    # skill remains usable with code-based layouts.
    TEMPLATE_BASED_LAYOUTS = set()


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------
def _resolve_layout(prs, kind):
    if kind not in LAYOUT_MAP:
        raise ValueError(f"Unknown base layout '{kind}'. Valid: {sorted(LAYOUT_MAP)}")
    m, l = LAYOUT_MAP[kind]
    return prs.slide_masters[m].slide_layouts[l]


def _placeholder_by_idx(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _set_text(placeholder, text):
    if placeholder is None or text is None:
        return
    placeholder.text_frame.text = str(text)


def _remove_body_placeholder(slide):
    body = _placeholder_by_idx(slide, 10)
    if body is not None:
        body._element.getparent().remove(body._element)


def _blank_canvas(prs, title):
    """Title + AOSIS footer kept; body placeholder removed to free the content area."""
    slide = prs.slides.add_slide(_resolve_layout(prs, "text"))
    _set_text(_placeholder_by_idx(slide, 0), title)
    _remove_body_placeholder(slide)
    return slide


def _add_text(slide, left, top, width, height, text, *,
              size=14, bold=False, color=None, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP, font="Arial"):
    # `color=None` resolves to BRAND.navy at call time so a rebound BRAND
    # (via build_deck --template) propagates correctly.
    if color is None:
        color = BRAND.navy
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    lines = str(text).split("\n") if text is not None else [""]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return tb


def _add_rect(slide, left, top, width, height, *, fill=None, line=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        if line_width is not None:
            shape.line.width = line_width
    shape.text_frame.margin_left = Emu(0)
    shape.text_frame.margin_right = Emu(0)
    shape.text_frame.margin_top = Emu(0)
    shape.text_frame.margin_bottom = Emu(0)
    return shape


def _add_rounded_rect(slide, left, top, width, height, *,
                      fill=None, line=None, line_width=None, corner=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    shape.adjustments[0] = corner
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        if line_width is not None:
            shape.line.width = line_width
    shape.text_frame.margin_left = Emu(0)
    shape.text_frame.margin_right = Emu(0)
    shape.text_frame.margin_top = Emu(0)
    shape.text_frame.margin_bottom = Emu(0)
    return shape


def _add_circle_number(slide, left, top, size, n, *, fill=None, font_size=13):
    # Late-bound default (see _add_text) so BRAND rebinding propagates.
    if fill is None:
        fill = BRAND.orange
    ell = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    ell.shadow.inherit = False
    ell.fill.solid()
    ell.fill.fore_color.rgb = fill
    ell.line.fill.background()
    tf = ell.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = str(n)
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = BRAND.white
    return ell


# ---------------------------------------------------------------------------
# Base layouts (placeholder-based)
# ---------------------------------------------------------------------------
def add_cover(prs, title, ref=None):
    slide = prs.slides.add_slide(_resolve_layout(prs, "cover"))
    _set_text(_placeholder_by_idx(slide, 0), title)
    _set_text(_placeholder_by_idx(slide, 10), ref or "")
    return slide


def add_section(prs, title, ref=None):
    return add_cover(prs, title, ref)


def add_closing(prs):
    return prs.slides.add_slide(_resolve_layout(prs, "closing"))


def add_text_slide(prs, title, bullets=None):
    slide = prs.slides.add_slide(_resolve_layout(prs, "text"))
    _set_text(_placeholder_by_idx(slide, 0), title)
    body = _placeholder_by_idx(slide, 10)
    if body is None or not bullets:
        return slide
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        text, level = (item.get("text", ""), int(item.get("level", 0))) if isinstance(item, dict) else (str(item), 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = max(0, min(4, level))
        for run in p.runs:
            run.font.size = Pt(16 if level == 0 else 14)
    return slide


def add_content_slide(prs, title, bullets=None, image=None):
    slide = prs.slides.add_slide(_resolve_layout(prs, "content"))
    _set_text(_placeholder_by_idx(slide, 0), title)
    if image and bullets:
        ph = _placeholder_by_idx(slide, 10)
        if ph is not None:
            ph._element.getparent().remove(ph._element)
        tb = slide.shapes.add_textbox(Inches(0.4), Inches(1.1), Inches(4.3), Inches(3.9))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(bullets):
            text, level = (item.get("text", ""), int(item.get("level", 0))) if isinstance(item, dict) else (str(item), 0)
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = ("• " if level == 0 else "    – ") + text
            for run in p.runs:
                run.font.size = Pt(15)
                run.font.color.rgb = BRAND.navy
        slide.shapes.add_picture(image, Inches(5.0), Inches(1.1),
                                 width=Inches(4.6), height=Inches(3.9))
    elif image:
        ph = _placeholder_by_idx(slide, 10)
        if ph is not None:
            try:
                ph.insert_picture(image)
            except (AttributeError, ValueError):
                left, top, width, height = ph.left, ph.top, ph.width, ph.height
                ph._element.getparent().remove(ph._element)
                slide.shapes.add_picture(image, left, top, width, height)
    else:
        body = _placeholder_by_idx(slide, 10)
        if body and bullets:
            tf = body.text_frame
            tf.clear()
            tf.word_wrap = True
            for i, item in enumerate(bullets):
                text, level = (item.get("text", ""), int(item.get("level", 0))) if isinstance(item, dict) else (str(item), 0)
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = text
                p.level = max(0, min(4, level))
                for run in p.runs:
                    run.font.size = Pt(16 if level == 0 else 14)
    return slide


# ---------------------------------------------------------------------------
# Rich visual layouts
# ---------------------------------------------------------------------------
def add_stat_grid(prs, title, stats, footnote=None):
    """Up to 4 big-number callouts in a row.
    stats = [{"value": "1.2 M€", "label": "Budget total", "accent": "orange"|"navy"}]"""
    slide = _blank_canvas(prs, title)
    n = max(1, min(len(stats), 4))
    stats = stats[:n]

    gap = Inches(0.25)
    card_w = Emu((CONTENT_W - gap * (n - 1)) / n)
    card_h = Inches(2.8)
    top = Inches(1.4)

    for i, st in enumerate(stats):
        accent = (st.get("accent") or ("orange" if i % 2 == 0 else "navy")).lower()
        bg = BRAND.orange if accent == "orange" else BRAND.navy
        left = Emu(CONTENT_LEFT + (card_w + gap) * i)
        _add_rounded_rect(slide, left, top, card_w, card_h, fill=bg, corner=0.06)

        # Adaptive value font size: scale down for longer strings so they
        # never wrap inside the card. Calibrated against an Arial Bold metric
        # where ~13 chars at 44pt fits a 2"-wide card.
        value = str(st.get("value", ""))
        if len(value) <= 4:
            value_size = 48
        elif len(value) <= 6:
            value_size = 40
        elif len(value) <= 9:
            value_size = 32
        else:
            value_size = 26

        _add_text(slide, left, Emu(top + Inches(0.5)), card_w, Inches(1.2),
                  value, size=value_size, bold=True, color=BRAND.white,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text(slide, Emu(left + Inches(0.2)), Emu(top + Inches(1.75)),
                  Emu(card_w - Inches(0.4)), Inches(0.9), st.get("label", ""),
                  size=13, color=BRAND.white, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    if footnote:
        _add_text(slide, CONTENT_LEFT, Inches(4.45), CONTENT_W, Inches(0.35),
                  footnote, size=10, color=BRAND.gray)
    return slide


def add_cards(prs, title, cards, columns=None):
    """Grid of cards. cards = [{"title": "...", "body": "...", "badge": "01"}]"""
    slide = _blank_canvas(prs, title)
    n = len(cards)
    if columns is None:
        columns = 2 if n <= 4 else 3
    columns = max(1, min(columns, 4))
    rows = (n + columns - 1) // columns

    gap = Inches(0.2)
    card_w = Emu((CONTENT_W - gap * (columns - 1)) / columns)
    avail_h = CONTENT_H
    card_h = Emu((avail_h - gap * (rows - 1)) / rows) if rows > 0 else avail_h

    for idx, card in enumerate(cards):
        r, c = divmod(idx, columns)
        left = Emu(CONTENT_LEFT + (card_w + gap) * c)
        top = Emu(CONTENT_TOP + (card_h + gap) * r)
        _add_rounded_rect(slide, left, top, card_w, card_h, fill=BRAND.gray_light, corner=0.04)
        _add_rect(slide, left, top, Inches(0.08), card_h, fill=BRAND.orange)

        pad_x = Inches(0.25)
        title_left = Emu(left + pad_x)
        title_width = Emu(card_w - pad_x * 2)

        badge = card.get("badge")
        if badge:
            b_size = Inches(0.42)
            b_left = Emu(left + card_w - b_size - Inches(0.2))
            b_top = Emu(top + Inches(0.18))
            _add_circle_number(slide, b_left, b_top, b_size, badge, font_size=11)
            title_width = Emu(title_width - b_size - Inches(0.1))

        _add_text(slide, title_left, Emu(top + Inches(0.18)),
                  title_width, Inches(0.5),
                  card.get("title", ""), size=16, bold=True, color=BRAND.navy)
        body_top = Emu(top + Inches(0.75))
        body_h = Emu(card_h - Inches(0.9))
        _add_text(slide, title_left, body_top, title_width, body_h,
                  card.get("body", ""), size=12, color=BRAND.gray)
    return slide


def add_comparison(prs, title, left, right):
    """Two columns. left/right = {"title": "...", "subtitle": "...", "items": [...]}"""
    slide = _blank_canvas(prs, title)
    col_w = Inches(4.45)
    col_h = Inches(3.7)
    top = Inches(1.15)
    left_x = Inches(0.4)
    right_x = Inches(5.15)

    def _draw_col(x, data, accent_color):
        _add_rounded_rect(slide, x, top, col_w, col_h, fill=BRAND.light,
                          line=BRAND.gray_light, corner=0.04)
        _add_rect(slide, x, top, col_w, Inches(0.6), fill=accent_color)
        _add_text(slide, Emu(x + Inches(0.3)), Emu(top + Inches(0.08)),
                  Emu(col_w - Inches(0.6)), Inches(0.45),
                  data.get("title", ""), size=18, bold=True, color=BRAND.white,
                  anchor=MSO_ANCHOR.MIDDLE)
        items_top = Inches(0.85)
        subtitle = data.get("subtitle")
        if subtitle:
            _add_text(slide, Emu(x + Inches(0.3)), Emu(top + Inches(0.7)),
                      Emu(col_w - Inches(0.6)), Inches(0.35),
                      subtitle, size=11, color=BRAND.gray)
            items_top = Inches(1.1)
        y = Emu(top + items_top)
        for item in data.get("items", []):
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                         Emu(x + Inches(0.3)), Emu(y + Inches(0.09)),
                                         Inches(0.12), Inches(0.12))
            dot.shadow.inherit = False
            dot.fill.solid()
            dot.fill.fore_color.rgb = accent_color
            dot.line.fill.background()
            _add_text(slide, Emu(x + Inches(0.55)), y,
                      Emu(col_w - Inches(0.85)), Inches(0.55),
                      str(item), size=13, color=BRAND.navy)
            y = Emu(y + Inches(0.5))

    _draw_col(left_x, left, BRAND.orange)
    _draw_col(right_x, right, BRAND.navy)
    return slide


def add_timeline(prs, title, phases):
    """Horizontal phase boxes. phases = [{"name": "...", "duration": "...", "detail": "..."}]"""
    slide = _blank_canvas(prs, title)
    n = max(1, min(len(phases), 6))
    phases = phases[:n]

    gap = Inches(0.12)
    box_w = Emu((CONTENT_W - gap * (n - 1)) / n)
    box_h = Inches(1.6)
    top = Inches(1.55)

    for i, ph in enumerate(phases):
        left = Emu(CONTENT_LEFT + (box_w + gap) * i)
        bg = BRAND.navy if i % 2 == 0 else BRAND.navy_alt
        _add_rounded_rect(slide, left, top, box_w, box_h, fill=bg, corner=0.06)

        b_size = Inches(0.45)
        _add_circle_number(slide, Emu(left + Inches(0.2)), Emu(top + Inches(0.2)),
                           b_size, i + 1, fill=BRAND.orange, font_size=13)

        _add_text(slide, Emu(left + Inches(0.2)), Emu(top + Inches(0.75)),
                  Emu(box_w - Inches(0.4)), Inches(0.45),
                  ph.get("name", ""), size=14, bold=True, color=BRAND.white)
        _add_text(slide, Emu(left + Inches(0.2)), Emu(top + Inches(1.18)),
                  Emu(box_w - Inches(0.4)), Inches(0.35),
                  ph.get("duration", ""), size=11, color=BRAND.orange)

        detail = ph.get("detail")
        if detail:
            _add_text(slide, left, Emu(top + box_h + Inches(0.2)),
                      box_w, Inches(1.6), detail, size=10, color=BRAND.gray)
    return slide


def add_process(prs, title, steps):
    """Vertical numbered process. steps = [{"title": "...", "detail": "..."}]"""
    slide = _blank_canvas(prs, title)
    n = max(1, min(len(steps), 5))
    steps = steps[:n]

    row_h = Emu(CONTENT_H / n)
    top0 = CONTENT_TOP
    c_size = Inches(0.6)
    c_left = Inches(0.7)

    for i, st in enumerate(steps):
        y = Emu(top0 + row_h * i)
        c_top = Emu(y + Inches(0.1))
        _add_circle_number(slide, c_left, c_top, c_size, i + 1, font_size=18)
        if i < n - 1:
            line_x = Emu(c_left + c_size / 2)
            line_top = Emu(c_top + c_size + Inches(0.05))
            line_bottom = Emu(top0 + row_h * (i + 1) + Inches(0.1))
            ln = slide.shapes.add_connector(1, line_x, line_top, line_x, line_bottom)
            ln.line.color.rgb = BRAND.orange
            ln.line.width = Pt(1.5)

        text_left = Emu(c_left + c_size + Inches(0.35))
        text_w = Emu(CONTENT_RIGHT - text_left)
        _add_text(slide, text_left, Emu(y + Inches(0.05)), text_w, Inches(0.5),
                  st.get("title", ""), size=18, bold=True, color=BRAND.navy)
        detail = st.get("detail")
        if detail:
            _add_text(slide, text_left, Emu(y + Inches(0.55)), text_w,
                      Emu(row_h - Inches(0.55)), detail, size=12, color=BRAND.gray)
    return slide


def add_quote(prs, text, author=None):
    """Editorial pull-quote slide: huge typography, orange side bar."""
    slide = prs.slides.add_slide(_resolve_layout(prs, "text"))
    _set_text(_placeholder_by_idx(slide, 0), "")
    _remove_body_placeholder(slide)
    # Thick orange vertical bar on the left (visual anchor)
    _add_rect(slide, Inches(0.8), Inches(1.1), Inches(0.12), Inches(3.4), fill=BRAND.orange)
    # Adaptive font size depending on quote length
    n = len(str(text))
    if n <= 80:
        qsize = 36
    elif n <= 140:
        qsize = 30
    elif n <= 220:
        qsize = 24
    else:
        qsize = 20
    _add_text(slide, Inches(1.2), Inches(1.1), Inches(7.8), Inches(3.4),
              text, size=qsize, bold=True, color=BRAND.navy, anchor=MSO_ANCHOR.MIDDLE)
    if author:
        _add_text(slide, Inches(1.2), Inches(4.55), Inches(7.8), Inches(0.4),
                  author.upper(), size=11, color=BRAND.orange, bold=True)
    return slide


def add_image_hero(prs, image, title=None, subtitle=None):
    """Full-bleed image with optional title/subtitle band at the bottom."""
    slide = prs.slides.add_slide(_resolve_layout(prs, "text"))
    _set_text(_placeholder_by_idx(slide, 0), "")
    _remove_body_placeholder(slide)
    slide.shapes.add_picture(image, 0, 0, SLIDE_W, SLIDE_H)
    if title or subtitle:
        band_h = Inches(1.4)
        band_top = Emu(SLIDE_H - band_h)
        _add_rect(slide, 0, band_top, SLIDE_W, band_h, fill=BRAND.navy)
        if title:
            _add_text(slide, Inches(0.6), Emu(band_top + Inches(0.2)),
                      Inches(8.8), Inches(0.7), title, size=24, bold=True, color=BRAND.white)
        if subtitle:
            _add_text(slide, Inches(0.6), Emu(band_top + Inches(0.85)),
                      Inches(8.8), Inches(0.45), subtitle, size=14, color=BRAND.orange)
    return slide


# ---------------------------------------------------------------------------
# Inspirational layouts (high-impact, editorial-style compositions)
# ---------------------------------------------------------------------------
def add_hero_stat(prs, title, value, label, context=None, supporting=None):
    """One massive number on the left, optional supporting list on the right.
    Magazine-cover style: heavy left-anchored composition with whitespace.

    title: small eyebrow tag at the top (can be a category like "IMPACT")
    value: the giant number (string — e.g. "75%", "1,2 M€", "12 mois")
    label: short bold caption under the value
    context: optional small gray text under the label
    supporting: optional list of bullet strings shown on the right
    """
    slide = _blank_canvas(prs, title)
    # Adaptive font size for the value based on length
    n = len(str(value))
    if n <= 3:
        v_size = 150
    elif n <= 5:
        v_size = 120
    elif n <= 7:
        v_size = 95
    else:
        v_size = 75

    # Big number
    big_w = Inches(5.4) if supporting else Inches(9.0)
    _add_text(slide, Inches(0.4), Inches(1.1), big_w, Inches(2.2),
              value, size=v_size, bold=True, color=BRAND.navy,
              anchor=MSO_ANCHOR.TOP)
    # Thin orange accent under the number
    _add_rect(slide, Inches(0.4), Inches(3.4), Inches(0.7), Inches(0.06), fill=BRAND.orange)
    # Label — adaptive font size based on length so an action-style
    # sentence fits without overflowing into the context line.
    label_len = len(str(label))
    if label_len <= 40:
        l_size = 18
    elif label_len <= 80:
        l_size = 16
    else:
        l_size = 14
    _add_text(slide, Inches(0.4), Inches(3.6), big_w, Inches(1.0),
              label, size=l_size, bold=True, color=BRAND.navy)
    if context:
        _add_text(slide, Inches(0.4), Inches(4.65), big_w, Inches(0.35),
                  context, size=11, color=BRAND.gray)
    # Supporting list on the right. Accept a string (becomes 1-item list)
    # or a list; iterating a bare string would split it character by character.
    if supporting:
        if isinstance(supporting, str):
            supporting = [supporting]
        y = Inches(1.4)
        for item in supporting[:5]:
            _add_rect(slide, Inches(6.1), Emu(y + Inches(0.12)),
                      Inches(0.18), Inches(0.04), fill=BRAND.orange)
            _add_text(slide, Inches(6.4), y, Inches(3.2), Inches(0.6),
                      str(item), size=13, color=BRAND.navy)
            y = Emu(y + Inches(0.65))
    return slide


def add_big_idea(prs, idea, title=None, supports=None, attribution=None):
    """A bold thesis statement layout. Strong typography on the left,
    optional supporting bullets on the right. Use this to anchor the
    deck's central argument."""
    slide = _blank_canvas(prs, title or "")
    # Orange accent line above the idea
    _add_rect(slide, Inches(0.4), Inches(1.2), Inches(0.6), Inches(0.06), fill=BRAND.orange)

    # Adaptive size for the idea based on length
    n = len(str(idea))
    if n <= 80:
        i_size = 40
    elif n <= 160:
        i_size = 32
    elif n <= 260:
        i_size = 26
    else:
        i_size = 22

    idea_width = Inches(5.4) if supports else Inches(9.0)
    _add_text(slide, Inches(0.4), Inches(1.5), idea_width, Inches(3.0),
              idea, size=i_size, bold=True, color=BRAND.navy, anchor=MSO_ANCHOR.TOP)
    if attribution:
        _add_text(slide, Inches(0.4), Inches(4.55), idea_width, Inches(0.35),
                  attribution.upper(), size=11, color=BRAND.orange, bold=True)

    if supports:
        # A subtle separator
        _add_rect(slide, Inches(6.0), Inches(1.3), Inches(0.02), Inches(3.4),
                  fill=BRAND.gray_light)
        y = Inches(1.4)
        for item in supports[:5]:
            if isinstance(item, dict):
                head = item.get("title", "")
                detail = item.get("detail", "")
            else:
                head, detail = str(item), None
            # Small orange tick
            _add_rect(slide, Inches(6.25), Emu(y + Inches(0.13)),
                      Inches(0.15), Inches(0.04), fill=BRAND.orange)
            _add_text(slide, Inches(6.5), y, Inches(3.1), Inches(0.4),
                      head, size=14, bold=True, color=BRAND.navy)
            if detail:
                _add_text(slide, Inches(6.5), Emu(y + Inches(0.35)),
                          Inches(3.1), Inches(0.5), detail, size=11, color=BRAND.gray)
                y = Emu(y + Inches(0.85))
            else:
                y = Emu(y + Inches(0.55))
    return slide


def add_matrix_2x2(prs, title, x_axis, y_axis, quadrants):
    """Strategic 2×2 matrix (BCG-style). Two axes, four quadrants.
    x_axis / y_axis = {"label": "Impact", "low": "Faible", "high": "Élevé"}
    quadrants = {
        "top_left":     {"title": "...", "items": [...]},
        "top_right":    {"title": "...", "items": [...]},
        "bottom_left":  {"title": "...", "items": [...]},
        "bottom_right": {"title": "...", "items": [...]},
    }
    Convention: top-right is typically the "star" quadrant.
    """
    slide = _blank_canvas(prs, title)

    grid_left = Inches(1.8)
    grid_top = Inches(1.15)
    grid_w = Inches(7.4)
    grid_h = Inches(3.5)
    quad_w = Emu(grid_w / 2)
    quad_h = Emu(grid_h / 2)

    # Y axis label (vertical-ish — rendered horizontally on the left, sorry)
    _add_text(slide, Inches(0.2), Emu(grid_top + grid_h / 2 - Inches(0.2)),
              Inches(1.5), Inches(0.4), y_axis.get("label", ""), size=12,
              bold=True, color=BRAND.navy, align=PP_ALIGN.CENTER)
    # Y axis ticks
    _add_text(slide, Inches(0.2), grid_top, Inches(1.5), Inches(0.3),
              y_axis.get("high", "Élevé"), size=10, color=BRAND.gray,
              align=PP_ALIGN.RIGHT)
    _add_text(slide, Inches(0.2), Emu(grid_top + grid_h - Inches(0.3)),
              Inches(1.5), Inches(0.3), y_axis.get("low", "Faible"),
              size=10, color=BRAND.gray, align=PP_ALIGN.RIGHT)
    # X axis label
    _add_text(slide, grid_left, Emu(grid_top + grid_h + Inches(0.5)),
              grid_w, Inches(0.35), x_axis.get("label", ""), size=12,
              bold=True, color=BRAND.navy, align=PP_ALIGN.CENTER)
    # X axis ticks
    _add_text(slide, grid_left, Emu(grid_top + grid_h + Inches(0.15)),
              Inches(2), Inches(0.3), x_axis.get("low", "Faible"),
              size=10, color=BRAND.gray, align=PP_ALIGN.LEFT)
    _add_text(slide, Emu(grid_left + grid_w - Inches(2)),
              Emu(grid_top + grid_h + Inches(0.15)), Inches(2),
              Inches(0.3), x_axis.get("high", "Élevé"), size=10,
              color=BRAND.gray, align=PP_ALIGN.RIGHT)

    # Quadrant background colours: top-right (star) is orange, others are
    # neutral with subtle differentiation
    quad_defs = [
        ("top_left",     grid_left,                  grid_top,                  BRAND.light,      BRAND.navy),
        ("top_right",    Emu(grid_left + quad_w),    grid_top,                  BRAND.orange,     BRAND.white),
        ("bottom_left",  grid_left,                  Emu(grid_top + quad_h),    BRAND.gray_light, BRAND.navy),
        ("bottom_right", Emu(grid_left + quad_w),    Emu(grid_top + quad_h),    BRAND.light,      BRAND.navy),
    ]
    for key, left, top, bg, text_color in quad_defs:
        q = quadrants.get(key, {})
        _add_rect(slide, left, top, quad_w, quad_h, fill=bg,
                  line=BRAND.gray_light, line_width=Pt(0.5))
        _add_text(slide, Emu(left + Inches(0.2)), Emu(top + Inches(0.12)),
                  Emu(quad_w - Inches(0.4)), Inches(0.4),
                  q.get("title", ""), size=13, bold=True, color=text_color)
        y = Emu(top + Inches(0.55))
        for item in q.get("items", [])[:4]:
            _add_text(slide, Emu(left + Inches(0.25)), y,
                      Emu(quad_w - Inches(0.4)), Inches(0.3),
                      "• " + str(item), size=10, color=text_color)
            y = Emu(y + Inches(0.32))
    return slide


def add_funnel(prs, title, stages):
    """N horizontal bars decreasing in width — a conversion funnel.
    stages = [{"name": "Leads", "value": "10000", "detail": "...optional..."}]
    """
    slide = _blank_canvas(prs, title)
    n = max(1, min(len(stages), 6))
    stages = stages[:n]

    total_h = Inches(3.7)
    bar_h = Emu(total_h / n - Inches(0.08))
    top0 = Inches(1.2)
    max_w = Inches(6.2)
    min_w = Inches(2.0)

    for i, st in enumerate(stages):
        ratio = 1 - (i / max(1, n - 1)) * 0.7 if n > 1 else 1.0
        w = Emu(min_w + (max_w - min_w) * ratio)
        x_center = Inches(3.5)  # funnel is centered on left half
        left = Emu(x_center - w / 2)
        y = Emu(top0 + (bar_h + Inches(0.08)) * i)
        # Color gradient: bottom = orange (the conversion!), top = navy
        if i == n - 1:
            color = BRAND.orange
        else:
            # Lerp navy → navy_soft
            color = BRAND.navy if i % 2 == 0 else BRAND.navy_alt
        _add_rounded_rect(slide, left, y, w, bar_h, fill=color, corner=0.15)
        _add_text(slide, left, y, w, bar_h,
                  st.get("name", ""), size=14, bold=True, color=BRAND.white,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Value to the right of the funnel
        value = st.get("value", "")
        if value:
            _add_text(slide, Inches(7.2), y, Inches(2.4), bar_h,
                      str(value), size=20, bold=True,
                      color=(BRAND.orange if i == n - 1 else BRAND.navy),
                      align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        detail = st.get("detail")
        if detail:
            _add_text(slide, Inches(7.2), Emu(y + Inches(0.35)),
                      Inches(2.4), bar_h, detail, size=10, color=BRAND.gray,
                      align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    return slide


# Chantier 1 (fix overflow) — approche 2 : amplitude horizontale réduite.
# Contrainte limitante = textbox `detail` de largeur 2.3" centrée sur le marker
# (label_w 1.9 + Inches(0.4) de débord lat.). Pour qu'elle reste dans [0, 10"],
# il faut marker_x ∈ [1.15", 8.85"]. On prend [1.2", 8.8"] (0.05" de cushion).
# La symétrie centre/centre des labels est préservée — pas de clamp asymétrique.
def add_roadmap(prs, title, milestones):
    """Horizontal timeline with diamond milestone markers.
    milestones = [{"date": "Jun '26", "name": "Audit", "detail": "..."}]
    Use when the story is about *when* things happen (vs. timeline which
    focuses on *what* happens in each phase).
    """
    slide = _blank_canvas(prs, title)
    n = max(2, min(len(milestones), 6))
    milestones = milestones[:n]

    line_y = Inches(2.9)
    line_left = Inches(1.2)
    line_right = Inches(8.8)
    span = line_right - line_left

    # Main horizontal line
    line = slide.shapes.add_connector(1, line_left, line_y, line_right, line_y)
    line.line.color.rgb = BRAND.navy
    line.line.width = Pt(2)

    for i, m in enumerate(milestones):
        x = Emu(line_left + (span * i) // (n - 1))
        # Diamond milestone marker
        diam_size = Inches(0.34)
        diam = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                      Emu(x - diam_size / 2),
                                      Emu(line_y - diam_size / 2),
                                      diam_size, diam_size)
        diam.shadow.inherit = False
        diam.fill.solid()
        diam.fill.fore_color.rgb = BRAND.orange if i % 2 == 0 else BRAND.navy
        diam.line.color.rgb = BRAND.white
        diam.line.width = Pt(2)

        # Alternate label position: even above, odd below — creates a
        # cleaner rhythm and avoids label collisions
        above = (i % 2 == 0)
        label_w = Inches(1.9)
        label_left = Emu(x - label_w / 2)

        if above:
            # Date above (orange small), name below date
            _add_text(slide, label_left, Inches(1.05), label_w, Inches(0.3),
                      m.get("date", "").upper(), size=11, bold=True,
                      color=BRAND.orange, align=PP_ALIGN.CENTER)
            _add_text(slide, label_left, Inches(1.4), label_w, Inches(0.4),
                      m.get("name", ""), size=14, bold=True,
                      color=BRAND.navy, align=PP_ALIGN.CENTER)
            detail = m.get("detail")
            if detail:
                _add_text(slide, Emu(label_left - Inches(0.2)),
                          Inches(1.85), Emu(label_w + Inches(0.4)),
                          Inches(0.8), detail, size=10, color=BRAND.gray,
                          align=PP_ALIGN.CENTER)
        else:
            _add_text(slide, label_left, Inches(3.35), label_w, Inches(0.3),
                      m.get("date", "").upper(), size=11, bold=True,
                      color=BRAND.orange, align=PP_ALIGN.CENTER)
            _add_text(slide, label_left, Inches(3.7), label_w, Inches(0.4),
                      m.get("name", ""), size=14, bold=True,
                      color=BRAND.navy, align=PP_ALIGN.CENTER)
            detail = m.get("detail")
            if detail:
                _add_text(slide, Emu(label_left - Inches(0.2)),
                          Inches(4.15), Emu(label_w + Inches(0.4)),
                          Inches(0.8), detail, size=10, color=BRAND.gray,
                          align=PP_ALIGN.CENTER)
    return slide


def _generate_abstract_background(output_path, mode="dots"):
    """Generate a subtle decorative background image using matplotlib.
    Used by hero/cinematic slides that ask for `decoration="abstract"`."""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import numpy as np
    except ImportError as e:
        raise ImportError(
            "matplotlib is required for the abstract-background hero decoration. "
            "Install with: pip install matplotlib"
        ) from e

    fig, ax = plt.subplots(figsize=(10, 5.62), dpi=120)
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 5.62)
    ax.set_aspect("equal")
    ax.axis("off")
    fig.patch.set_facecolor(f"#{BRAND.light}")

    rng = np.random.default_rng(seed=42)

    if mode == "dots":
        # Scattered navy dots fading from dense (right) to sparse (left)
        n = 500
        x = rng.uniform(0, 10, n)
        sizes = rng.uniform(2, 30, n)
        # density bias toward right
        weights = x / 10
        keep = rng.uniform(0, 1, n) < (weights ** 1.5)
        x = x[keep]
        y = rng.uniform(0, 5.62, len(x))
        sizes = sizes[keep]
        colors = [f"#{BRAND.navy}" if rng.random() > 0.15 else f"#{BRAND.orange}" for _ in x]
        ax.scatter(x, y, s=sizes, c=colors, alpha=0.18)
    elif mode == "lines":
        for i in range(30):
            xs = rng.uniform(0, 10, 2)
            ys = rng.uniform(0, 5.62, 2)
            ax.plot(xs, ys, color=f"#{BRAND.navy}", alpha=0.06, linewidth=0.8)
    else:  # gradient
        # Vertical gradient white → very light navy
        grad = np.linspace(0, 1, 100).reshape(-1, 1)
        ax.imshow(grad, extent=[0, 10, 0, 5.62], cmap="bone_r",
                  aspect="auto", alpha=0.2)

    plt.subplots_adjust(left=0, right=1, top=1, bottom=0)
    plt.savefig(output_path, dpi=120, facecolor=f"#{BRAND.light}",
                bbox_inches="tight", pad_inches=0)
    plt.close(fig)


def add_swot(prs, title, strengths, weaknesses, opportunities, threats):
    """SWOT analysis 2×2.
    Each parameter is a dict {"title": "...", "items": [...]}.
    Convention: S/O are positive (filled), W/T are neutral (outlined).
    Layout vs matrix_2x2: SWOT has fixed quadrant semantics (no axes), so
    it gets its own visual language with large S/W/O/T letters.
    """
    slide = _blank_canvas(prs, title)

    grid_left = Inches(0.8)
    grid_top = Inches(1.15)
    grid_w = Inches(8.4)
    grid_h = Inches(3.7)
    cell_w = Emu(grid_w / 2)
    cell_h = Emu(grid_h / 2)

    cells = [
        # (data,            letter, left,                      top,                       bg,         fg,    accent_for_letter)
        (strengths,         "S",    grid_left,                 grid_top,                  BRAND.navy,       BRAND.white, BRAND.orange),
        (weaknesses,        "W",    Emu(grid_left + cell_w),   grid_top,                  BRAND.light,      BRAND.navy,  BRAND.navy_alt),
        (opportunities,     "O",    grid_left,                 Emu(grid_top + cell_h),    BRAND.orange,     BRAND.white, BRAND.navy),
        (threats,           "T",    Emu(grid_left + cell_w),   Emu(grid_top + cell_h),    BRAND.gray_light, BRAND.navy,  BRAND.orange),
    ]

    for data, letter, left, top, bg, fg, accent in cells:
        _add_rect(slide, left, top, cell_w, cell_h, fill=bg,
                  line=BRAND.gray_light, line_width=Pt(0.5))
        # Big letter
        _add_text(slide, Emu(left + Inches(0.15)), Emu(top + Inches(0.08)),
                  Inches(0.7), Inches(0.7), letter,
                  size=44, bold=True, color=accent,
                  anchor=MSO_ANCHOR.TOP)
        # Title
        _add_text(slide, Emu(left + Inches(0.9)), Emu(top + Inches(0.2)),
                  Emu(cell_w - Inches(1.0)), Inches(0.4),
                  (data or {}).get("title", ""), size=14, bold=True, color=fg)
        # Items — cap at 3, tighter spacing, smaller font so the cell never overflows
        y = Emu(top + Inches(0.78))
        for item in (data or {}).get("items", [])[:3]:
            _add_text(slide, Emu(left + Inches(0.25)), y,
                      Emu(cell_w - Inches(0.4)), Inches(0.32),
                      "• " + str(item), size=11, color=fg)
            y = Emu(y + Inches(0.34))
    return slide


def add_pyramid(prs, title, levels, inverted=False):
    """Hierarchical pyramid. levels listed bottom-to-top by default
    (base first, apex last) — set inverted=True to flip.

    Use for Maslow-style frameworks, value chains, or any "this rests on
    that" structure. Detail text for each level is rendered to the right.
    """
    slide = _blank_canvas(prs, title)
    n = max(1, min(len(levels), 5))
    levels = levels[:n]

    available_h = Inches(3.55)
    gap = Inches(0.08)
    band_h = Emu((available_h - gap * (n - 1)) / n)
    top0 = Inches(1.2)
    max_w = Inches(6.4)
    min_w = Inches(2.0)
    x_center = Inches(3.6)  # pyramid is centered on left ~70% of slide

    for render_pos in range(n):
        # Visual position 0 = top of slide (apex of pyramid)
        # Visual position n-1 = bottom of slide (base of pyramid)
        # When inverted=False (default), levels[0] is the base → at the bottom
        # When inverted=True, levels[0] is at the top
        if inverted:
            lev = levels[render_pos]
            ratio_from_apex = render_pos / max(1, n - 1)
        else:
            lev = levels[n - 1 - render_pos]
            ratio_from_apex = render_pos / max(1, n - 1)
        w = Emu(min_w + (max_w - min_w) * ratio_from_apex)
        left = Emu(x_center - w / 2)
        y = Emu(top0 + (band_h + gap) * render_pos)
        color = BRAND.navy if render_pos % 2 == 0 else BRAND.navy_alt
        # Highlight apex/base in orange to draw the eye
        if (not inverted and render_pos == 0) or (inverted and render_pos == n - 1):
            color = BRAND.orange
        _add_rect(slide, left, y, w, band_h, fill=color)
        _add_text(slide, left, y, w, band_h,
                  lev.get("name", ""), size=13, bold=True, color=BRAND.white,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Detail to the right of the pyramid
        detail = lev.get("detail")
        if detail:
            _add_text(slide, Inches(7.4), y, Inches(2.4), band_h,
                      detail, size=11, color=BRAND.gray,
                      align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    return slide


def add_org_chart(prs, title, leader, reports):
    """Two-level organisation chart (leader + direct reports).
    Each report can optionally list `members` (a short list of names) shown
    as a stacked group under that report's box.

    leader  = {"name": "...", "role": "..."}
    reports = [{"name": "...", "role": "...", "members": ["...", "..."]}, ...]

    Up to 5 reports — beyond that the row gets too cramped on 16:9.
    """
    slide = _blank_canvas(prs, title)

    # Leader box centered at the top
    leader_w = Inches(3.0)
    leader_h = Inches(0.95)
    leader_left = Emu(SLIDE_W / 2 - leader_w / 2)
    leader_top = Inches(1.15)
    _add_rounded_rect(slide, leader_left, leader_top, leader_w, leader_h,
                      fill=BRAND.navy, corner=0.05)
    _add_text(slide, leader_left, Emu(leader_top + Inches(0.15)),
              leader_w, Inches(0.35),
              (leader or {}).get("name", ""), size=13, bold=True, color=BRAND.white,
              align=PP_ALIGN.CENTER)
    _add_text(slide, leader_left, Emu(leader_top + Inches(0.5)),
              leader_w, Inches(0.35),
              (leader or {}).get("role", ""), size=10, color=BRAND.orange,
              align=PP_ALIGN.CENTER)

    n = min(len(reports), 5) if reports else 0
    if n == 0:
        return slide
    reports = reports[:n]

    report_w = Inches(1.7)
    report_h = Inches(0.85)
    report_top = Inches(2.85)
    gap = Inches(0.2)
    total_w = report_w * n + gap * (n - 1)
    reports_start_left = Emu(SLIDE_W / 2 - total_w / 2)

    # Connector line geometry
    leader_bottom_y = Emu(leader_top + leader_h)
    connector_y = Emu(leader_bottom_y + Inches(0.4))
    slide_center_x = Emu(SLIDE_W / 2)

    # Vertical line from leader bottom down to the horizontal connector
    v_root = slide.shapes.add_connector(1, slide_center_x, leader_bottom_y,
                                        slide_center_x, connector_y)
    v_root.line.color.rgb = BRAND.gray
    v_root.line.width = Pt(1)

    if n > 1:
        first_center_x = Emu(reports_start_left + report_w / 2)
        last_center_x = Emu(reports_start_left + (report_w + gap) * (n - 1) + report_w / 2)
        h_bus = slide.shapes.add_connector(1, first_center_x, connector_y,
                                           last_center_x, connector_y)
        h_bus.line.color.rgb = BRAND.gray
        h_bus.line.width = Pt(1)

    for i, r in enumerate(reports):
        left = Emu(reports_start_left + (report_w + gap) * i)
        center_x = Emu(left + report_w / 2)
        # Drop from horizontal bus to report box
        drop = slide.shapes.add_connector(1, center_x, connector_y,
                                          center_x, report_top)
        drop.line.color.rgb = BRAND.gray
        drop.line.width = Pt(1)
        # Report box (white-on-orange-border for visual hierarchy)
        _add_rounded_rect(slide, left, report_top, report_w, report_h,
                          fill=BRAND.light, line=BRAND.orange, line_width=Pt(1.5), corner=0.05)
        _add_text(slide, left, Emu(report_top + Inches(0.12)),
                  report_w, Inches(0.32), r.get("name", ""),
                  size=11, bold=True, color=BRAND.navy, align=PP_ALIGN.CENTER)
        _add_text(slide, left, Emu(report_top + Inches(0.44)),
                  report_w, Inches(0.32), r.get("role", ""),
                  size=9, color=BRAND.gray, align=PP_ALIGN.CENTER)
        # Optional members stack underneath
        members = r.get("members") or []
        if members:
            members = members[:3]
            mem_top = Inches(4.0)
            mem_h = Inches(0.32)
            box_bottom_y = Emu(report_top + report_h)
            link = slide.shapes.add_connector(1, center_x, box_bottom_y,
                                              center_x, mem_top)
            link.line.color.rgb = BRAND.gray
            link.line.width = Pt(0.5)
            y = mem_top
            for m in members:
                name = m if isinstance(m, str) else m.get("name", "")
                _add_rect(slide, Emu(left + Inches(0.1)), y,
                          Emu(report_w - Inches(0.2)), mem_h,
                          fill=BRAND.gray_light, line=BRAND.gray_light)
                _add_text(slide, Emu(left + Inches(0.1)), y,
                          Emu(report_w - Inches(0.2)), mem_h,
                          name, size=9, color=BRAND.navy,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                y = Emu(y + mem_h + Inches(0.05))
    return slide


def add_agenda(prs, title, items):
    """Numbered table-of-contents / agenda layout.
    items = [{"title": "...", "detail": "..."}]. Up to 6 items.

    Use as the second slide of any deck longer than 8 slides so the audience
    knows the path; also great as a section-recap slide.
    """
    slide = _blank_canvas(prs, title)
    n = max(1, min(len(items), 6))
    items = items[:n]

    available_h = Inches(3.85)
    gap = Inches(0.05)
    row_h = Emu((available_h - gap * (n - 1)) / n)
    top0 = Inches(1.15)

    for i, item in enumerate(items):
        y = Emu(top0 + (row_h + gap) * i)
        # Big orange number
        _add_text(slide, Inches(0.5), y, Inches(1.4), row_h,
                  f"{i + 1:02d}", size=46, bold=True, color=BRAND.orange,
                  anchor=MSO_ANCHOR.MIDDLE)
        # Thin vertical separator
        sep_x = Inches(2.0)
        sep = slide.shapes.add_connector(1, sep_x, Emu(y + Inches(0.15)),
                                         sep_x, Emu(y + row_h - Inches(0.15)))
        sep.line.color.rgb = BRAND.gray_light
        sep.line.width = Pt(1)
        # Item title (left) and detail (right side of the same row)
        detail = item.get("detail")
        text_left = Inches(2.25)
        if detail:
            # Reserve ~1.6" on the right for the detail
            title_w = Inches(5.7)
            _add_text(slide, text_left, y, title_w, row_h,
                      item.get("title", ""), size=17, bold=True, color=BRAND.navy,
                      anchor=MSO_ANCHOR.MIDDLE)
            _add_text(slide, Inches(8.0), y, Inches(1.5), row_h,
                      detail, size=12, color=BRAND.orange, bold=True,
                      align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        else:
            _add_text(slide, text_left, y, Inches(7.2), row_h,
                      item.get("title", ""), size=18, bold=True, color=BRAND.navy,
                      anchor=MSO_ANCHOR.MIDDLE)
    return slide


def add_dashboard(prs, title, stats=None, chart=None, chart_title=None):
    """Mixed executive dashboard: compact stats row at the top + chart below.

    stats = up to 4 entries, each {"value": "...", "label": "..."}.
    chart = same shape as for the `chart` layout (type/labels/data/series).
    chart_title is an optional small label above the chart.

    The styling is intentionally restrained — no big coloured boxes — so the
    slide reads like a Financial Times / Bloomberg dashboard rather than
    decoration.
    """
    slide = _blank_canvas(prs, title)

    stats = (stats or [])[:4]
    n_stats = max(1, len(stats)) if stats else 0

    if n_stats:
        gap = Inches(0.2)
        stat_w = Emu((CONTENT_W - gap * (n_stats - 1)) / n_stats)
        stat_top = Inches(1.15)
        for i, st in enumerate(stats):
            left = Emu(CONTENT_LEFT + (stat_w + gap) * i)
            # Value
            _add_text(slide, Emu(left + Inches(0.1)), stat_top,
                      Emu(stat_w - Inches(0.2)), Inches(0.6),
                      st.get("value", ""), size=30, bold=True, color=BRAND.navy)
            # Thin orange underline
            _add_rect(slide, Emu(left + Inches(0.1)), Emu(stat_top + Inches(0.65)),
                      Inches(0.4), Inches(0.04), fill=BRAND.orange)
            # Label
            _add_text(slide, Emu(left + Inches(0.1)), Emu(stat_top + Inches(0.78)),
                      Emu(stat_w - Inches(0.2)), Inches(0.4),
                      st.get("label", ""), size=11, color=BRAND.gray)

    # Chart fills the remaining space
    chart_top = Inches(2.55) if n_stats else Inches(1.15)
    chart_h = Emu(Inches(5.0) - chart_top)

    if chart:
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        tmp.close()
        try:
            _render_chart_png(chart, tmp.name)
            if chart_title:
                _add_text(slide, CONTENT_LEFT, chart_top, CONTENT_W, Inches(0.3),
                          chart_title.upper(), size=10, bold=True, color=BRAND.gray)
                chart_top = Emu(chart_top + Inches(0.3))
                chart_h = Emu(chart_h - Inches(0.3))
            slide.shapes.add_picture(tmp.name, Inches(0.8), chart_top,
                                     width=Inches(8.4), height=chart_h)
        finally:
            try:
                os.unlink(tmp.name)
            except OSError:
                pass
    return slide


# ---------------------------------------------------------------------------
# Auto-generated charts (matplotlib → PNG → image)
# ---------------------------------------------------------------------------
def _render_chart_png(chart_spec, output_path):
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        from matplotlib import rcParams
    except ImportError as e:
        raise ImportError(
            "matplotlib is required for `chart` and `dashboard` layouts. "
            "Install with: pip install matplotlib"
        ) from e

    rcParams["font.family"] = "DejaVu Sans"
    rcParams["axes.spines.top"] = False
    rcParams["axes.spines.right"] = False
    rcParams["axes.edgecolor"] = f"#{BRAND.gray}"
    rcParams["axes.labelcolor"] = f"#{BRAND.navy}"
    rcParams["xtick.color"] = f"#{BRAND.gray}"
    rcParams["ytick.color"] = f"#{BRAND.gray}"
    rcParams["axes.titlecolor"] = f"#{BRAND.navy}"
    rcParams["axes.titleweight"] = "bold"

    palette_main = f"#{BRAND.navy}"
    palette_accent = f"#{BRAND.orange}"
    palette_series = [f"#{BRAND.navy}", f"#{BRAND.orange}", f"#{BRAND.gray}", f"#{BRAND.accent5}", f"#{BRAND.accent4}"]

    kind = chart_spec.get("type", "bar")
    labels = chart_spec.get("labels", [])
    series = chart_spec.get("series")
    data = chart_spec.get("data")

    fig, ax = plt.subplots(figsize=(7.5, 4.2), dpi=140)

    if kind in ("bar", "column"):
        if series:
            import numpy as np
            x = np.arange(len(labels))
            w = 0.8 / max(1, len(series))
            for i, s in enumerate(series):
                ax.bar(x + i * w - 0.4 + w / 2, s["values"], w,
                       label=s.get("name", ""),
                       color=palette_series[i % len(palette_series)])
            ax.set_xticks(x)
            ax.set_xticklabels(labels)
            ax.legend(frameon=False, loc="best")
        else:
            ax.bar(labels, data, color=palette_main)
            # Optional explicit highlight (index of the bar to colour in
            # accent). If "highlight" is "max" or "min", auto-pick. Default:
            # no highlight — all bars in the main palette.
            hl = chart_spec.get("highlight")
            if hl == "max" and data:
                ax.patches[data.index(max(data))].set_facecolor(palette_accent)
            elif hl == "min" and data:
                ax.patches[data.index(min(data))].set_facecolor(palette_accent)
            elif isinstance(hl, int) and 0 <= hl < len(data):
                ax.patches[hl].set_facecolor(palette_accent)
    elif kind == "barh":
        if data:
            ax.barh(labels, data, color=palette_main)
            hl = chart_spec.get("highlight")
            if hl == "max":
                ax.patches[data.index(max(data))].set_facecolor(palette_accent)
            elif hl == "min":
                ax.patches[data.index(min(data))].set_facecolor(palette_accent)
            elif isinstance(hl, int) and 0 <= hl < len(data):
                ax.patches[hl].set_facecolor(palette_accent)
    elif kind == "line":
        if series:
            for i, s in enumerate(series):
                ax.plot(labels, s["values"], marker="o", linewidth=2.2,
                        label=s.get("name", ""),
                        color=palette_series[i % len(palette_series)])
            ax.legend(frameon=False, loc="best")
        else:
            ax.plot(labels, data, marker="o", linewidth=2.4, color=palette_main)
    elif kind == "pie":
        wedges, texts, autotexts = ax.pie(
            data, labels=labels,
            colors=palette_series[: len(data)],
            autopct="%1.0f%%", startangle=90,
            wedgeprops={"linewidth": 1.5, "edgecolor": "white"},
        )
        for t in autotexts:
            t.set_color("white")
            t.set_fontweight("bold")
    else:
        ax.text(0.5, 0.5, f"chart type '{kind}' not supported",
                ha="center", va="center", transform=ax.transAxes)

    if chart_spec.get("ylabel"):
        ax.set_ylabel(chart_spec["ylabel"])
    if chart_spec.get("xlabel"):
        ax.set_xlabel(chart_spec["xlabel"])

    plt.tight_layout()
    plt.savefig(output_path, dpi=140, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def add_chart_slide(prs, title, chart_spec, commentary=None):
    """Generate a chart from data. If commentary present → side-by-side bullets+chart."""
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    tmp.close()
    try:
        _render_chart_png(chart_spec, tmp.name)
        if commentary:
            return add_content_slide(prs, title, bullets=commentary, image=tmp.name)
        slide = _blank_canvas(prs, title)
        slide.shapes.add_picture(tmp.name, Inches(0.8), Inches(1.05),
                                 width=Inches(8.4), height=Inches(3.9))
        return slide
    finally:
        try:
            os.unlink(tmp.name)
        except OSError:
            pass


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------
DISPATCH = {
    "cover":      lambda prs, s: add_cover(prs, s.get("title", ""), s.get("ref")),
    "section":    lambda prs, s: add_section(prs, s.get("title", ""), s.get("ref")),
    "closing":    lambda prs, s: add_closing(prs),
    "text":       lambda prs, s: add_text_slide(prs, s.get("title", ""), s.get("bullets")),
    "content":    lambda prs, s: add_content_slide(prs, s.get("title", ""),
                                                   s.get("bullets"), s.get("image")),
    "stat_grid":  lambda prs, s: add_stat_grid(prs, s.get("title", ""),
                                               s.get("stats", []), s.get("footnote")),
    "cards":      lambda prs, s: add_cards(prs, s.get("title", ""),
                                           s.get("cards", []), s.get("columns")),
    "comparison": lambda prs, s: add_comparison(prs, s.get("title", ""),
                                                s.get("left", {}), s.get("right", {})),
    "timeline":   lambda prs, s: add_timeline(prs, s.get("title", ""), s.get("phases", [])),
    "process":    lambda prs, s: add_process(prs, s.get("title", ""), s.get("steps", [])),
    "quote":      lambda prs, s: add_quote(prs, s.get("text", ""), s.get("author")),
    "image_hero": lambda prs, s: add_image_hero(prs, s.get("image"),
                                                s.get("title"), s.get("subtitle")),
    "chart":      lambda prs, s: add_chart_slide(prs, s.get("title", ""),
                                                 s.get("chart", {}),
                                                 s.get("commentary")),
    # Inspirational layouts
    "hero_stat":  lambda prs, s: add_hero_stat(prs, s.get("title", ""),
                                               s.get("value", ""), s.get("label", ""),
                                               s.get("context"), s.get("supporting")),
    "big_idea":   lambda prs, s: add_big_idea(prs, s.get("idea", ""),
                                              s.get("title"), s.get("supports"),
                                              s.get("attribution")),
    "matrix_2x2": lambda prs, s: add_matrix_2x2(prs, s.get("title", ""),
                                                s.get("x_axis", {}), s.get("y_axis", {}),
                                                s.get("quadrants", {})),
    "funnel":     lambda prs, s: add_funnel(prs, s.get("title", ""), s.get("stages", [])),
    "roadmap":    lambda prs, s: add_roadmap(prs, s.get("title", ""),
                                             s.get("milestones", [])),
    "swot":       lambda prs, s: add_swot(prs, s.get("title", ""),
                                          s.get("strengths", {}),
                                          s.get("weaknesses", {}),
                                          s.get("opportunities", {}),
                                          s.get("threats", {})),
    "pyramid":    lambda prs, s: add_pyramid(prs, s.get("title", ""),
                                             s.get("levels", []),
                                             s.get("inverted", False)),
    "org_chart":  lambda prs, s: add_org_chart(prs, s.get("title", ""),
                                               s.get("leader", {}),
                                               s.get("reports", [])),
    "agenda":     lambda prs, s: add_agenda(prs, s.get("title", ""),
                                            s.get("items", [])),
    "dashboard":  lambda prs, s: add_dashboard(prs, s.get("title", ""),
                                               s.get("stats"),
                                               s.get("chart"),
                                               s.get("chart_title")),
}


def _add_layout_debug_footer(slide, layout_name: str, is_code_based: bool) -> None:
    """Stamp a tiny grey footer on a slide showing its layout name. Used
    only when build_deck is invoked with --debug-layouts."""
    tag = f"[layout: {layout_name}{' (code-based)' if is_code_based else ''}]"
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(5.35),
                                   Inches(3.5), Inches(0.20))
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = tag
    run.font.size = Pt(6)
    run.font.italic = True
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0xA0, 0xA0, 0xA0)


def build_deck(spec, output_path, template_path=None, debug_layouts=False,
               auto_images=True):
    """Build a deck from a JSON spec.

    Routing:
      - If `spec.slides[i].layout` is in TEMPLATE_BASED_LAYOUTS (discovered
        dynamically from AOSIS_template.pptx), the slide is rendered by
        `template_engine.render_template_slide` (cross-deck deepcopy + fill).
      - Otherwise, it is rendered by DISPATCH[layout] (code-based path).

    The build base is always AOSIS_template.pptx (with sample slides
    stripped) so that templated slides keep their master-level decorations
    (filigrane, diagonal backgrounds, …) and code-based slides land on the
    same brand surface.
    """
    template_path = Path(template_path or TEMPLATE_PATH)
    if not template_path.exists():
        raise FileNotFoundError(f"AOSIS template not found at {template_path}")
    # Rebind the module-level palette so all add_* helpers (which look up
    # BRAND.* at call time via module globals) match the template actually
    # in use — useful when --template overrides the default.
    global BRAND
    BRAND = BrandPalette.from_template(template_path)
    # Rebind chart palette so charts track --template overrides.
    try:
        import chart_engine
        chart_engine.set_brand(BRAND)
    except ImportError:
        pass
    # Rebind template_engine palette likewise (alternation rules).
    if hasattr(_template_engine, "set_brand"):
        _template_engine.set_brand(BRAND)

    # Build base = stripped clone of AOSIS_template.pptx. This guarantees
    # that templated slides inherit the same masters/layouts as the source.
    prs = Presentation(str(template_path))
    strip_sample_slides(prs)

    def _stamp_new_slides(prev_n: int, layout_name: str, is_code: bool):
        """Apply the debug footer to every slide added since prev_n."""
        if not debug_layouts:
            return
        for k in range(prev_n, len(prs.slides)):
            _add_layout_debug_footer(prs.slides[k], layout_name, is_code)

    cover = spec.get("cover") or {}
    if cover:
        prev = len(prs.slides)
        add_cover(prs, cover.get("title", ""), cover.get("ref"))
        _stamp_new_slides(prev, "cover", True)

    for i, s in enumerate(spec.get("slides", [])):
        layout = s.get("layout", "text")
        prev = len(prs.slides)
        # Inject the auto_images flag into each slide spec so render_template_slide
        # can decide whether to call image_engine for `{{IMAGE}}` placeholders.
        s = dict(s)
        s.setdefault("_auto_images", auto_images)
        if layout in TEMPLATE_BASED_LAYOUTS:
            # Paginate layouts that have a max-per-page limit (agenda).
            paginate_limit = _template_engine.PAGINATED_LAYOUTS.get(layout)
            if paginate_limit and len(s.get("items", [])) > paginate_limit:
                for page_spec in _paginate_slide(s, paginate_limit):
                    render_template_slide(prs, template_path, layout, page_spec)
            else:
                render_template_slide(prs, template_path, layout, s)
            _stamp_new_slides(prev, layout, False)
            continue
        if layout not in DISPATCH:
            raise ValueError(
                f"slides[{i}].layout='{layout}' is invalid. "
                f"Valid code-based: {sorted(DISPATCH)}. "
                f"Valid template-based: {sorted(TEMPLATE_BASED_LAYOUTS)}."
            )
        DISPATCH[layout](prs, s)
        _stamp_new_slides(prev, layout, True)

    if spec.get("closing", True):
        prev = len(prs.slides)
        add_closing(prs)
        _stamp_new_slides(prev, "closing", True)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def _paginate_slide(slide_spec, max_per_page):
    """Yield per-page spec dicts for a layout that supports pagination.
    Each page carries an `_number_offset` for continued numbering."""
    items = slide_spec.get("items", [])
    base_title = slide_spec.get("title", "")
    for start in range(0, len(items), max_per_page):
        page = dict(slide_spec)
        page["items"] = items[start:start + max_per_page]
        page["_number_offset"] = start
        if start > 0 and base_title and "(suite)" not in base_title:
            page["title"] = f"{base_title} (suite)"
        yield page


def main(argv=None):
    p = argparse.ArgumentParser(description="Build an AOSIS-branded .pptx from a JSON spec.")
    p.add_argument("spec", nargs="?")
    p.add_argument("output", nargs="?")
    p.add_argument("--template", default=None)
    p.add_argument("--debug-layouts", action="store_true",
                   help="Stamp a tiny grey footer on each slide showing its layout name.")
    p.add_argument("--no-images", action="store_true",
                   help="Skip auto-fetching stock photos for {{IMAGE}} placeholders.")
    p.add_argument("--no-cache-images", action="store_true",
                   help="Bypass the local Pexels image cache and fetch fresh.")
    p.add_argument("--clear-image-cache", action="store_true",
                   help="Wipe the local image cache and exit (no deck built).")
    args = p.parse_args(argv)

    # Chantier 22 — utility action: clear cache and exit
    if args.clear_image_cache:
        sys.path.insert(0, str(Path(__file__).parent))
        from image_engine import clear_image_cache, CACHE_DIR
        n = clear_image_cache()
        print(f"OK — cleared {n} files from {CACHE_DIR}")
        return 0

    # Chantier 22 — apply --no-cache-images before fetch
    if args.no_cache_images:
        sys.path.insert(0, str(Path(__file__).parent))
        from image_engine import set_cache_enabled
        set_cache_enabled(False)

    if not args.spec or not args.output:
        p.error("spec and output are required unless using --clear-image-cache")

    with open(args.spec, encoding="utf-8") as f:
        spec = json.load(f)
    out = build_deck(spec, args.output, template_path=args.template,
                     debug_layouts=args.debug_layouts,
                     auto_images=not args.no_images)
    print(f"OK — wrote {out} ({out.stat().st_size:,} bytes)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
