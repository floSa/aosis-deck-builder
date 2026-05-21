"""template_engine.py — Render slides from named templates in the AOSIS template.

The engine discovers available layouts dynamically by scanning slide names
(`cSld@name`) in the template pptx (canonically `assets/AOSIS_template.pptx`),
then renders a requested layout by:

  1. Loading the template pptx (read-only).
  2. Finding the slide whose name matches the requested layout.
  3. Adding a new slide to the target deck on the equivalent slide layout
     (looked up by name — assumes the target deck shares the same masters).
  4. Deep-copying the source slide's shape XML into the new slide.
  5. Filling placeholders marked with `{{KEY}}` from the spec dict.
  6. Duplicating `{{REPEAT_ITEM}}` groups based on `spec.items`.

Conventions
-----------
- Shape names follow `{{UPPER_SNAKE_CASE}}`.
- Simple placeholder: `{{KEY}}` → replaced by `spec[key.lower()]`.
- Group repeat: `{{REPEAT_ITEM}}` (GROUP shape containing `{{ITEM_*}}` children)
  → duplicated once per element in `spec.items` (or another list key,
  inferred by layout name).
- Image placeholder: `{{IMAGE}}` → replaced by image from `spec.image`.
- Indexed placeholder: `{{TAG_<N>_<KEY>}}` → filled from
  `spec.<inferred_list>[N-1][key.lower()]`. Items beyond the spec length
  cause the shape to be removed.
- Quadrant placeholder: `{{QUAD_<position>_<KEY>}}` → filled from
  `spec.quadrants[position][key]` (matrix_2x2_styled convention).

Architectural note: the target deck MUST derive from the template pptx (i.e.,
be a stripped clone) so masters/layouts/themes are shared. Otherwise
deep-copied shapes may render incorrectly (missing master backgrounds, broken
layout references). This is enforced by `build_deck.py` which clones the
template and strips its sample slides at build start.
"""
from __future__ import annotations

import re
from copy import deepcopy
from pathlib import Path
from typing import Optional

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches


# -----------------------------------------------------------------------------
# Patterns & conventions
# -----------------------------------------------------------------------------
# Brand palette injection (optional) — used by alternation rules to resolve
# semantic color names (e.g. "orange" → BRAND.orange). Build_deck.py calls
# `set_brand(BRAND)` at load time and on --template overrides.
_BRAND = None


def set_brand(palette) -> None:
    """Inject the active brand palette so alternation rules can resolve
    semantic color names. Safe no-op if never called."""
    global _BRAND
    _BRAND = palette


_COLOR_FALLBACK = {
    "orange":   "F26622",
    "navy_alt": "1E2261",
    "navy":     "14163C",
    "accent4":  "F9B233",
    "accent5":  "7CB342",
    "accent6":  "E63946",
}


def _resolve_color(name: str) -> str:
    """Resolve a palette attribute name → hex string (no '#').
    Falls back to canonical AOSIS hexes when no brand was injected."""
    if _BRAND is not None and hasattr(_BRAND, name):
        v = getattr(_BRAND, name)
        if v is not None:
            return str(v).upper().lstrip('#')
    return _COLOR_FALLBACK.get(name, "F26622")


# -----------------------------------------------------------------------------
# Alternation rules — applied to each REPEAT_ITEM copy by index (0-based).
# -----------------------------------------------------------------------------
ALTERNATION_RULES = {
    "framework_3cards": {
        "type": "fill_color",
        # Card background shape inside each REPEAT_ITEM copy
        "shape_names": ["{{ITEM_BOXE}}"],
        # Cycle of palette attribute names; resolved live via _resolve_color
        "colors": ["orange", "navy_alt"],
    },
    "process_steps": {
        "type": "fill_color",
        # Numbered marker disk inside each copy
        "shape_names": ["{{ITEM_MARKER}}"],
        "colors": ["orange", "navy_alt"],
    },
    "roadmap_styled": {
        "type": "vertical_flip",
        # Text shapes that flip above/below the marker on odd indices
        "text_shape_names": ["{{ITEM_DATE}}", "{{ITEM_MILESTONE}}"],
        # Anchor shape used to compute the flip baseline
        "anchor_shape_name": "{{ITEM_MARKER}}",
        # EMU gap between marker edge and text block when flipped below
        "margin_emu": 150000,
    },
}


PLACEHOLDER_RE = re.compile(r'^\{\{([A-Z][A-Z0-9_]*)\}\}$')
INDEXED_PLACEHOLDER_RE = re.compile(r'^\{\{([A-Z]+(?:_[A-Z]+)*)_(\d+)_([A-Z][A-Z0-9_]*)\}\}$')
# Item-level sub-placeholders inside a {{REPEAT_ITEM}} group. Accept the
# canonical "ITEM" prefix plus layout-specific prefixes (e.g. "KPI" for the
# `kpi_with_chart` template which uses `{{KPI_LABEL}}` / `{{KPI_VALUE}}`).
ITEM_PLACEHOLDER_RE = re.compile(r'^\{\{(?:ITEM|KPI)_([A-Z][A-Z0-9_]*)\}\}$')

# Per-layout overrides for how `{{REPEAT_ITEM}}` copies are laid out.
# Layouts not listed fall back to 'horizontal' distribution (centred row).
DISTRIBUTION_BY_LAYOUT = {
    'next_steps': 'vertical',
    'kpi_with_chart': 'vertical_left',  # KPI cards on left, chart fills right
    'roadmap_styled': 'horizontal_alternating',
    'agenda_diagonal': 'single_column',  # Chantier 9 — was 'grid_2cols'
}

# Layouts that support multi-page pagination (auto-split when items > max).
PAGINATED_LAYOUTS = {
    'agenda_diagonal': 7,  # Chantier 11 — back to 7 for breathing room
}

# Tag → spec list key (used by indexed placeholders `{{TAG_<N>_<KEY>}}`).
TAG_TO_LIST_KEY = {
    'KPI': 'kpis',
    'STEP': 'items',
    'COL': 'items',
    'MILESTONE': 'items',
    'ACTION': 'items',
    'CARD': 'items',
    'PHASE': 'items',
    'MONTH': 'months',
}


class TemplateError(Exception):
    """Raised when a template layout cannot be found or rendered."""


# -----------------------------------------------------------------------------
# Build-base preparation
# -----------------------------------------------------------------------------
def strip_sample_slides(prs) -> int:
    """Remove every slide from `prs` while keeping its masters, layouts and
    themes. Returns the number of slides removed.

    Uses python-pptx's `drop_rel` to ensure the corresponding slide parts
    are unlinked AND removed from the package — otherwise the resulting
    .pptx contains duplicate slide-XML entries that PowerPoint rejects
    (HRESULT 0x80CB4404).

    This is the canonical way to turn `exhibits.pptx` (which ships with
    sample slides demonstrating each template) into the empty build base
    used by `build_deck()`.
    """
    sldIdLst = prs.slides._sldIdLst
    sld_id_elements = list(sldIdLst)
    for sld_id_el in sld_id_elements:
        rId = sld_id_el.get(qn('r:id'))
        sldIdLst.remove(sld_id_el)
        prs.part.drop_rel(rId)
    return len(sld_id_elements)


# -----------------------------------------------------------------------------
# Discovery
# -----------------------------------------------------------------------------
def discover_template_layouts(exhibits_path: "str | Path") -> dict[str, int]:
    """Return mapping `{cSld.name: slide_index}` for every named slide in
    exhibits.pptx. Slides with no `cSld@name` are skipped.

    Raises TemplateError if the file is missing or unreadable.
    """
    path = Path(exhibits_path)
    if not path.exists():
        raise TemplateError(f"Exhibits file not found: {path}")
    try:
        prs = Presentation(str(path))
    except Exception as e:
        raise TemplateError(f"Cannot open exhibits file {path}: {e}") from e
    out: dict[str, int] = {}
    for i, slide in enumerate(prs.slides):
        cSld = slide._element.find(qn('p:cSld'))
        if cSld is not None:
            name = cSld.get('name')
            if name:
                out[name] = i
    return out


# -----------------------------------------------------------------------------
# Public entry point
# -----------------------------------------------------------------------------
def render_template_slide(
    prs,
    exhibits_path: "str | Path",
    layout_name: str,
    spec: dict,
):
    """Render a templated slide into `prs`, sourced from exhibits.pptx.

    The target deck MUST derive from exhibits.pptx (clone + strip samples)
    so the masters/layouts are shared.
    """
    src_prs = Presentation(str(exhibits_path))
    src_slide = _find_slide_by_name(src_prs, layout_name)
    if src_slide is None:
        available = sorted(discover_template_layouts(exhibits_path).keys())
        raise TemplateError(
            f"Template layout '{layout_name}' not found.\n"
            f"Available named layouts: {available}"
        )

    # Reproduce the source slide's layout in the target by name
    src_layout_name = src_slide.slide_layout.name
    target_layout = _find_layout_by_name(prs, src_layout_name)
    if target_layout is None:
        # Fallback: use the first slide layout of the first master
        target_layout = prs.slide_masters[0].slide_layouts[0]

    # Chantier 10 — spec preprocessing: flatten nested dicts like
    # `before: {title, bullets}` into `before_title`, `before_bullets` keys
    # so they match flat `{{BEFORE_TITLE}}` placeholders in the template.
    spec = _flatten_nested_groups(spec)

    new_slide = prs.slides.add_slide(target_layout)

    # Remove placeholders auto-injected by the layout (they'll be replaced
    # by the deep-copied shapes from the source slide).
    for ph in list(new_slide.placeholders):
        ph._element.getparent().remove(ph._element)

    # Deep-copy each source shape into the new slide's shape tree
    for sh in src_slide.shapes:
        new_el = deepcopy(sh._element)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    # Tag the new slide with the layout name (for traceability)
    new_cSld = new_slide._element.find(qn('p:cSld'))
    if new_cSld is not None:
        new_cSld.set('name', layout_name)

    # Now process placeholders. Order matters:
    #   - REPEAT_ITEM groups first (they may contain {{ITEM_*}} sub-placeholders)
    #   - then quadrant placeholders ({{QUAD_TOP_LEFT_TITLE}} etc.)
    #   - then indexed placeholders ({{COL_1_TITLE}} etc.)
    #   - then chart placeholder (must run before simple placeholders so the
    #     `{{CHART_PLACEHOLDER}}` shape is replaced with an image rather than
    #     deleted as "unknown {{KEY}}")
    #   - then simple placeholders ({{TITLE}} etc.)
    #   - finally image placeholders
    _process_repeat_items(new_slide, layout_name, spec)
    _process_quad_placeholders(new_slide, spec)
    _process_indexed_placeholders(new_slide, spec)
    _process_chart_placeholder(new_slide, spec)
    # Chantier 10 — remove `{{XXX_GROUP}}` whose inner `<XXX>` is empty
    # in spec, BEFORE simple processor handles individual placeholders.
    _process_orphan_groups(new_slide, spec)
    _process_simple_placeholders(new_slide, spec)
    _process_image_placeholders(new_slide, spec)

    # Chantier 14 — freeform composition on canvas_blank when `blocks` provided
    if layout_name == 'canvas_blank' and spec.get('blocks'):
        _render_canvas_blank_freeform(new_slide, spec)

    # Chantier 15 — data_table dynamic rendering from spec.table
    if layout_name == 'data_table' and spec.get('table'):
        _process_data_table(new_slide, spec)

    # Chantier 16 — closing_diagonal: aggressive title shrink for long titles
    if layout_name == 'closing_diagonal':
        _shrink_closing_title(new_slide)

    return new_slide


# -----------------------------------------------------------------------------
# Source-slide lookups
# -----------------------------------------------------------------------------
def _flatten_nested_groups(spec: dict) -> dict:
    """Pre-process spec: for ergonomic keys like `before: {title, bullets}`,
    flatten them into `before_title`, `before_bullets` so the simple
    placeholder processor can match `{{BEFORE_TITLE}}` etc.

    Only flattens 1 level deep, only for nested dicts whose values are
    scalars or lists. The original keys are preserved alongside the flat
    aliases (so existing flat-key callers still work).
    """
    out = dict(spec)
    for outer_key, outer_val in list(spec.items()):
        if not isinstance(outer_val, dict):
            continue
        # Skip dicts that the engine handles via dedicated processors
        if outer_key in ('quadrants', 'chart', 'cover'):
            continue
        for inner_key, inner_val in outer_val.items():
            flat_key = f"{outer_key}_{inner_key}"
            # Don't overwrite if user already provided the flat key
            if flat_key in out:
                continue
            if isinstance(inner_val, (str, int, float, list)):
                out[flat_key] = inner_val
    return out


def _find_slide_by_name(prs, name: str):
    for s in prs.slides:
        cSld = s._element.find(qn('p:cSld'))
        if cSld is not None and cSld.get('name') == name:
            return s
    return None


def _find_layout_by_name(prs, name: str):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    return None


# -----------------------------------------------------------------------------
# Simple placeholder processing — {{KEY}} → spec[key.lower()]
# -----------------------------------------------------------------------------
def _process_simple_placeholders(slide, spec: dict) -> None:
    """Replace text in every {{KEY}} shape. Shapes whose key is missing
    in spec are removed for a clean output.

    Shapes handled by specialised processors (quad, indexed, image, chart)
    are skipped here so we don't second-guess them.
    """
    to_remove = []
    for sh in list(_iter_top_level_shapes(slide)):
        m = PLACEHOLDER_RE.match(sh.name)
        if not m:
            continue
        key = m.group(1).lower()
        if key == 'image':
            continue  # handled by _process_image_placeholders
        if key == 'chart_placeholder':
            continue  # handled by _process_chart_placeholder
        # Specialised placeholders already handled above — don't touch them
        # (chantier 10 fix: simple processor was deleting QUAD shapes after
        # the quad processor filled them).
        if key.startswith('quad_'):
            continue
        # Indexed placeholders like COL_1_TITLE are handled by the indexed
        # processor; their key has a digit segment.
        if re.search(r'_\d+_', key):
            continue
        if key in spec:
            value = spec[key]
            if value is None:
                to_remove.append(sh)
            elif isinstance(value, (str, int, float)):
                _replace_text_keep_style(sh, str(value))
            elif isinstance(value, list):
                # Lists become bulleted text
                bullets = '\n'.join(f'• {v}' for v in value)
                _replace_text_keep_style(sh, bullets)
        else:
            to_remove.append(sh)
    for sh in to_remove:
        sh._element.getparent().remove(sh._element)


# -----------------------------------------------------------------------------
# Indexed placeholder processing — {{TAG_<N>_<KEY>}} → spec.<list>[N-1][key]
# -----------------------------------------------------------------------------
def _process_indexed_placeholders(slide, spec: dict) -> None:
    to_remove = []
    for sh in list(_iter_top_level_shapes(slide)):
        m = INDEXED_PLACEHOLDER_RE.match(sh.name)
        if not m:
            continue
        tag = m.group(1)
        idx = int(m.group(2)) - 1  # 1-based to 0-based
        key = m.group(3).lower()

        # Special-case: QUAD_<position>_<KEY> reads from spec.quadrants[position]
        if tag == 'QUAD':
            quad_position = m.group(2).lower()  # but m.group(2) is the digit…
            # Adjust: QUAD doesn't use indexed digits — it uses position names.
            # The regex above only matches when the middle is digits, so QUAD
            # cases like {{QUAD_TOP_LEFT_TITLE}} won't match this regex.
            continue

        list_key = TAG_TO_LIST_KEY.get(tag, 'items')
        items = spec.get(list_key, [])

        if 0 <= idx < len(items):
            item = items[idx]
            if isinstance(item, dict) and key in item:
                _replace_text_keep_style(sh, str(item[key]))
            elif isinstance(item, str) and key == 'text':
                _replace_text_keep_style(sh, item)
            else:
                to_remove.append(sh)
        else:
            to_remove.append(sh)
    for sh in to_remove:
        sh._element.getparent().remove(sh._element)


# Separate regex for quad placeholders: {{QUAD_<position>_<key>}} where
# position is e.g. TOP_LEFT, BOTTOM_RIGHT, etc.
QUAD_PLACEHOLDER_RE = re.compile(
    r'^\{\{QUAD_(TOP_LEFT|TOP_RIGHT|BOTTOM_LEFT|BOTTOM_RIGHT)_([A-Z][A-Z0-9_]*)\}\}$'
)


_ORPHAN_GROUP_RE = re.compile(r'^\{\{([A-Z][A-Z0-9_]*)_GROUP\}\}$')


def _process_orphan_groups(slide, spec: dict) -> None:
    """Remove top-level groups named `{{<KEY>_GROUP}}` whose `<KEY>` value
    is absent or empty in `spec`. Use this to bundle decorative shapes
    (e.g. an orange bar) with their associated text placeholder so they
    disappear together when the text is empty.

    No-op for templates that don't use the convention.
    """
    to_remove = []
    for sh in list(_iter_top_level_shapes(slide)):
        m = _ORPHAN_GROUP_RE.match(sh.name or '')
        if not m:
            continue
        key = m.group(1).lower()
        val = spec.get(key)
        if val in (None, '', [], {}):
            to_remove.append(sh)
    for sh in to_remove:
        sh._element.getparent().remove(sh._element)


def _process_quad_placeholders(slide, spec: dict) -> None:
    """Fill {{QUAD_<position>_<KEY>}} from spec.quadrants[position][key]."""
    quadrants = spec.get('quadrants', {})
    to_remove = []
    for sh in list(_iter_top_level_shapes(slide)):
        m = QUAD_PLACEHOLDER_RE.match(sh.name)
        if not m:
            continue
        position = m.group(1).lower()
        key = m.group(2).lower()
        quad = quadrants.get(position, {})
        # Accept `items` as alias for `bullets` (and vice versa).
        value = None
        if isinstance(quad, dict):
            if key in quad:
                value = quad[key]
            elif key == 'bullets' and 'items' in quad:
                value = quad['items']
            elif key == 'items' and 'bullets' in quad:
                value = quad['bullets']
        if value is not None:
            # Chantier 11 — max 3 bullets per quadrant
            if isinstance(value, list):
                if key in ('bullets', 'items') and len(value) > 3:
                    import sys as _sys
                    print(
                        f"matrix_2x2_styled: quad {position!r} has "
                        f"{len(value)} bullets, truncated to 3",
                        file=_sys.stderr,
                    )
                    value = value[:3]
                value = '\n'.join(f'• {v}' for v in value)
            _replace_text_keep_style(sh, str(value))
            # Chantier 11 — bottom-left anchor + auto-shrink for bullets
            if key in ('bullets', 'items'):
                _set_text_anchor(sh._element, 'b')
                _maybe_shrink_to_fit(sh._element, str(value), min_sz=1000)
        else:
            to_remove.append(sh)
    for sh in to_remove:
        sh._element.getparent().remove(sh._element)

    # Chantier 17 — harmonize font sizes across the 4 quadrants so all titles
    # share one size and all bullets share another, regardless of per-quadrant
    # content length.
    _apply_uniform_font_size_to_quads(slide)


def _apply_uniform_font_size_to_quads(slide) -> None:
    """For `matrix_2x2_styled`, normalize the font size across the 4
    `{{QUAD_<position>_TITLE}}` shapes (= take the smallest observed size
    and apply it to all four). Same for the 4 `{{QUAD_<position>_BULLETS}}`
    shapes. Per-quadrant auto-shrink (Chantier 11) may have produced
    different sizes; this pass restores visual uniformity.
    """
    title_pattern_re = re.compile(
        r'^\{\{QUAD_(TOP_LEFT|TOP_RIGHT|BOTTOM_LEFT|BOTTOM_RIGHT)_TITLE\}\}$'
    )
    bullets_pattern_re = re.compile(
        r'^\{\{QUAD_(TOP_LEFT|TOP_RIGHT|BOTTOM_LEFT|BOTTOM_RIGHT)_BULLETS\}\}$'
    )

    title_shapes = []
    bullets_shapes = []
    for sh in _iter_top_level_shapes(slide):
        nm = sh.name or ''
        if title_pattern_re.match(nm):
            title_shapes.append(sh)
        elif bullets_pattern_re.match(nm):
            bullets_shapes.append(sh)

    for shape_group in (title_shapes, bullets_shapes):
        if len(shape_group) < 2:
            continue
        # Collect all rPr.sz across the group
        rPr_per_shape = []
        min_sz = None
        for sh in shape_group:
            rPrs = list(sh._element.iter(qn('a:rPr')))
            if not rPrs:
                continue
            rPr_per_shape.append(rPrs)
            for rPr in rPrs:
                sz_str = rPr.get('sz')
                if sz_str:
                    sz = int(sz_str)
                    if min_sz is None or sz < min_sz:
                        min_sz = sz
        if min_sz is None:
            continue
        # Apply min sz to all rPrs in the group
        target = str(min_sz)
        for rPrs in rPr_per_shape:
            for rPr in rPrs:
                rPr.set('sz', target)


# -----------------------------------------------------------------------------
# REPEAT_ITEM processing — duplicate group N times based on spec.items
# -----------------------------------------------------------------------------
def _process_repeat_items(slide, layout_name: str, spec: dict) -> None:
    """Find {{REPEAT_ITEM}} groups, duplicate them per spec.items, fill
    each copy's {{ITEM_*}} sub-placeholders, and remove the template."""
    repeats = [sh for sh in _iter_top_level_shapes(slide)
               if sh.name == '{{REPEAT_ITEM}}']
    if not repeats:
        return

    # Most slides have ONE {{REPEAT_ITEM}}. comparison_2cols/before_after
    # could have two (left/right) but the current exhibits.pptx has one
    # per slide (the user duplicates manually).
    template_shape = repeats[0]

    # Identify the source spec key:
    spec_key = _infer_repeat_spec_key(layout_name, spec)
    items = spec.get(spec_key, [])

    parent = template_shape._element.getparent()

    if not items:
        # Nothing to render — remove the template
        parent.remove(template_shape._element)
        return

    # Read template geometry (read at runtime — never hardcoded)
    base_left = template_shape.left
    base_top = template_shape.top
    base_w = template_shape.width
    base_h = template_shape.height

    # Slide dimensions (also read at runtime)
    slide_w = slide.part.package.presentation_part.presentation.slide_width
    slide_h = slide.part.package.presentation_part.presentation.slide_height

    distribution = DISTRIBUTION_BY_LAYOUT.get(layout_name, 'horizontal')

    positions = _compute_positions(
        distribution=distribution,
        n=len(items),
        base_left=base_left,
        base_top=base_top,
        base_w=base_w,
        base_h=base_h,
        slide_w=slide_w,
        slide_h=slide_h,
    )

    # Capture the original element index so duplicates appear after it
    template_el = template_shape._element
    template_idx = list(parent).index(template_el)

    # Chantier 9 — `_number_offset` lets the caller continue auto-numbering
    # across multiple pages (used by agenda_diagonal pagination).
    number_offset = int(spec.get('_number_offset', 0))
    icons = spec.get('icons', [])

    for i, item in enumerate(items):
        new_el = deepcopy(template_el)
        # Rename the copy so subsequent placeholder passes don't re-match it
        # as the unknown `{{REPEAT_ITEM}}` simple placeholder.
        _rename_group(new_el, f'repeat_item_copy_{i+1}')
        left, top, w, h = positions[i]
        dx = left - base_left
        dy = top - base_top
        _shift_group(new_el, dx=dx, dy=dy)
        _fill_item_placeholders(new_el, item, i + number_offset, distribution)
        _apply_alternation(new_el, layout_name, i)
        # Chantier 11 — per-layout postprocess for the deepcopy (anchor, shrink…)
        _apply_layout_postprocess(new_el, layout_name, i)
        parent.insert(template_idx + 1 + i, new_el)

        # Chantier 9 — optional icon injection (framework_3cards et al.)
        if i < len(icons) and icons[i]:
            _inject_icon_on_copy(slide, new_el, icons[i], dx, dy)

    # Remove the template
    parent.remove(template_el)

    # Chantier 18 — drop shadow on kpi_card backgrounds.
    # Chantier 19 — height-AND-width-aware KPI value sizing : compute the
    # max font size that fits each KPI's own shape, take the MIN across the
    # 3 KPIs (uniformity), apply that one size to all. Prevents value/label
    # overlap when the kpi_card height is small.
    if layout_name == 'kpi_with_chart':
        kpi_value_shapes = []
        for top_sh in _iter_top_level_shapes(slide):
            if top_sh.shape_type != 6:
                continue
            for sp in top_sh._element.iter(qn('p:sp')):
                cNvPr = sp.find('.//' + qn('p:cNvPr'))
                sp_name = cNvPr.get('name', '') if cNvPr is not None else ''
                if sp_name == 'kpi_card':
                    _apply_drop_shadow(sp)
                elif sp_name == '{{KPI_VALUE}}':
                    kpi_value_shapes.append(sp)

        if kpi_value_shapes:
            per_kpi_max = []
            for sp in kpi_value_shapes:
                w_emu, h_emu = _read_sp_dimensions(sp)
                if w_emu and h_emu:
                    text = ''.join(t.text or '' for t in sp.iter(qn('a:t')))
                    # Chantier 20 — floor_pt=14: the kpi_with_chart shape is
                    # narrow (1.00" × 0.47"), values with wide chars like "M€"
                    # would wrap if the floor stayed at 24pt. Allow shrinking
                    # below to keep single-line rendering.
                    per_kpi_max.append(
                        _compute_max_kpi_font_size(text, w_emu, h_emu,
                                                   floor_pt=14)
                    )
            if per_kpi_max:
                unified_sz = str(min(per_kpi_max) * 100)  # OOXML sz = pt × 100
                for sp in kpi_value_shapes:
                    for rPr in sp.iter(qn('a:rPr')):
                        rPr.set('sz', unified_sz)

    # Chantier 16 — uniform font size across all repeated copies for layouts
    # where per-item auto-shrink would produce visually inconsistent sizing.
    uniform_names = _UNIFORM_REPEAT_SHAPES.get(layout_name)
    if uniform_names:
        _apply_uniform_font_size_to_repeats(slide, uniform_names)


# -----------------------------------------------------------------------------
# Alternation application — colors and vertical position
# -----------------------------------------------------------------------------
def _apply_layout_postprocess(grp_element, layout_name: str, index: int) -> None:
    """Layout-specific touch-ups applied to each REPEAT_ITEM copy AFTER its
    text substitution and alternation. Currently handles the roadmap_styled
    case: anchor text labels by the bottom when positioned above the
    timeline axis (odd 1-indexed → above), by the top when below.
    """
    if layout_name == 'roadmap_styled':
        # 0-indexed: copies 0, 2, 4 sit above the axis (text grows upward)
        # → anchor by bottom so the text stays glued to its bottom edge.
        # copies 1, 3 sit below → anchor by top (the default).
        is_above = (index % 2 == 0)
        anchor = 'b' if is_above else 't'
        for shape_name in ('{{ITEM_DATE}}', '{{ITEM_MILESTONE}}'):
            sp = _find_first_named_sp(grp_element, shape_name)
            if sp is None:
                continue
            _set_text_anchor(sp, anchor)
            # Auto-shrink: roadmap milestones can have long names; reduce font
            # by 2pt steps until estimated width fits, floor 10pt.
            txt = _read_sp_text(sp)
            if txt:
                _maybe_shrink_to_fit(sp, txt, min_sz=1000)


def _read_sp_text(sp_element) -> str:
    """Concatenate the text content of a `<p:sp>` (used to drive auto-shrink)."""
    txBody = sp_element.find(qn('p:txBody'))
    if txBody is None:
        return ''
    parts = []
    for t in txBody.iter(qn('a:t')):
        if t.text:
            parts.append(t.text)
    return ' '.join(parts)


def _apply_alternation(grp_element, layout_name: str, index: int) -> None:
    """Apply the layout's alternation rule (if any) to a freshly duplicated
    REPEAT_ITEM copy. `index` is 0-based."""
    rule = ALTERNATION_RULES.get(layout_name)
    if rule is None:
        return
    kind = rule.get("type")
    if kind == "fill_color":
        colors = rule["colors"]
        hex_color = _resolve_color(colors[index % len(colors)])
        for shape_name in rule["shape_names"]:
            for sp in _find_named_sp(grp_element, shape_name):
                _set_solid_srgb_fill(sp, hex_color)
    elif kind == "vertical_flip":
        # Even index (0, 2, 4…) keeps the template's default position.
        # Odd index (1, 3, 5…) flips text shapes below the anchor.
        if index % 2 == 0:
            return
        anchor = _find_first_named_sp(grp_element, rule["anchor_shape_name"])
        if anchor is None:
            return
        anchor_y, anchor_h = _read_sp_y_h(anchor)
        if anchor_y is None:
            return
        anchor_bottom = anchor_y + anchor_h
        margin = int(rule.get("margin_emu", 150000))
        # Find current topmost text y to compute the uniform delta
        text_sps = []
        for shape_name in rule["text_shape_names"]:
            text_sps.extend(_find_named_sp(grp_element, shape_name))
        if not text_sps:
            return
        ys = []
        for sp in text_sps:
            ty, _ = _read_sp_y_h(sp)
            if ty is not None:
                ys.append(ty)
        if not ys:
            return
        top_text_y = min(ys)
        new_top_text_y = anchor_bottom + margin
        delta_y = new_top_text_y - top_text_y
        for sp in text_sps:
            _shift_sp_y(sp, delta_y)


def _find_named_sp(grp_element, target_name: str):
    """Yield every <p:sp> descendant whose cNvPr@name == target_name."""
    for sp in grp_element.iter(qn('p:sp')):
        cNvPr = sp.find('.//' + qn('p:cNvPr'))
        if cNvPr is not None and cNvPr.get('name', '') == target_name:
            yield sp


def _find_first_named_sp(grp_element, target_name: str):
    for sp in _find_named_sp(grp_element, target_name):
        return sp
    return None


def _read_sp_y_h(sp_element):
    """Return (y, height) in EMU for a <p:sp>, or (None, None)."""
    spPr = sp_element.find(qn('p:spPr'))
    if spPr is None:
        return (None, None)
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is None:
        return (None, None)
    off = xfrm.find(qn('a:off'))
    ext = xfrm.find(qn('a:ext'))
    if off is None or ext is None:
        return (None, None)
    return (int(off.get('y', '0')), int(ext.get('cy', '0')))


def _shift_sp_y(sp_element, dy: int) -> None:
    spPr = sp_element.find(qn('p:spPr'))
    if spPr is None:
        return
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is None:
        return
    off = xfrm.find(qn('a:off'))
    if off is None:
        return
    off.set('y', str(int(off.get('y', '0')) + dy))


def _set_solid_srgb_fill(sp_element, hex_color: str) -> None:
    """Replace (or create) the solidFill of <p:sp> with an srgbClr of the
    given hex value (no leading '#'). Idempotent."""
    spPr = sp_element.find(qn('p:spPr'))
    if spPr is None:
        return
    solid = spPr.find(qn('a:solidFill'))
    if solid is not None:
        # Drop any existing color child and inject a fresh srgbClr
        for c in list(solid):
            solid.remove(c)
        srgb = etree.SubElement(solid, qn('a:srgbClr'))
        srgb.set('val', hex_color.upper())
        return
    # No solidFill present — insert one after the geometry child (or end)
    solid = etree.SubElement(spPr, qn('a:solidFill'))
    srgb = etree.SubElement(solid, qn('a:srgbClr'))
    srgb.set('val', hex_color.upper())


def _inject_icon_on_copy(slide, grp_element, icon_name: str, dx: int, dy: int) -> None:
    """Fetch an Iconify icon and overlay it on the {{ITEM_ICON}} shape of a
    deep-copied REPEAT_ITEM group. Silent no-op on any failure (network,
    missing cairosvg, unknown icon).

    Position is computed at the slide level (not inside the group) so we can
    use python-pptx's high-level add_picture API. The math assumes the
    template's group offset == chOffset (no scaling) — true for all current
    AOSIS template groups.
    """
    try:
        import icon_engine
    except ImportError:
        return

    icon_sp = None
    for sp in grp_element.iter(qn('p:sp')):
        cNvPr = sp.find('.//' + qn('p:cNvPr'))
        if cNvPr is not None and cNvPr.get('name', '') == '{{ITEM_ICON}}':
            icon_sp = sp
            break
    if icon_sp is None:
        return

    spPr = icon_sp.find(qn('p:spPr'))
    xfrm = spPr.find(qn('a:xfrm')) if spPr is not None else None
    if xfrm is None:
        return
    off = xfrm.find(qn('a:off'))
    ext = xfrm.find(qn('a:ext'))
    if off is None or ext is None:
        return
    abs_left = int(off.get('x', '0')) + dx
    abs_top = int(off.get('y', '0')) + dy
    w = int(ext.get('cx', '0'))
    h = int(ext.get('cy', '0'))
    if w == 0 or h == 0:
        return

    # Inset the picture inside the circle so the icon doesn't touch the edge.
    pad = int(min(w, h) * 0.18)
    pic_left = abs_left + pad
    pic_top = abs_top + pad
    pic_w = w - 2 * pad
    pic_h = h - 2 * pad

    # Default icon color = navy (dark) so it reads on the white circle bg.
    icon_color = '#' + _resolve_color('navy')

    try:
        png_bytes = icon_engine.fetch_icon_png(
            icon_name, size_px=256, color=icon_color, timeout=5.0,
        )
    except Exception:
        # Network down, cairosvg missing, icon not found — fail silently
        return

    import io
    slide.shapes.add_picture(io.BytesIO(png_bytes), pic_left, pic_top,
                             width=pic_w, height=pic_h)


def _rename_group(grp_element, new_name: str) -> None:
    """Set the group's `<p:cNvPr@name>` attribute."""
    nvGrpSpPr = grp_element.find(qn('p:nvGrpSpPr'))
    if nvGrpSpPr is None:
        return
    cNvPr = nvGrpSpPr.find(qn('p:cNvPr'))
    if cNvPr is not None:
        cNvPr.set('name', new_name)


def _infer_repeat_spec_key(layout_name: str, spec: dict) -> str:
    """Heuristically pick the spec list key that feeds REPEAT_ITEM."""
    # Layouts with non-'items' conventions
    if layout_name == 'kpi_with_chart':
        return 'kpis' if 'kpis' in spec else 'items'
    if layout_name == 'gantt_phases':
        return 'phases' if 'phases' in spec else 'items'
    return 'items'


def _compute_positions(
    distribution: str,
    n: int,
    base_left: int,
    base_top: int,
    base_w: int,
    base_h: int,
    slide_w: int,
    slide_h: int,
) -> list[tuple[int, int, int, int]]:
    """Return a list of (left, top, w, h) tuples for N copies. All values in
    EMU. Computation is purely geometric — no palette knowledge."""
    margin = Inches(0.40)

    if distribution == 'horizontal':
        # N items centred across the slide horizontally, all on base_top
        available = slide_w - 2 * margin
        gap = Inches(0.20)
        if n == 1:
            item_w = base_w
            left0 = (slide_w - item_w) // 2
            return [(left0, base_top, item_w, base_h)]
        item_w = (available - gap * (n - 1)) // n
        return [
            (margin + (item_w + gap) * i, base_top, item_w, base_h)
            for i in range(n)
        ]

    if distribution == 'vertical':
        # N items stacked vertically below base_top, full width preserved
        available = slide_h - base_top - Inches(0.50)  # leave room for source
        gap = Inches(0.12)
        if n == 1:
            return [(base_left, base_top, base_w, base_h)]
        item_h = (available - gap * (n - 1)) // n
        return [
            (base_left, base_top + (item_h + gap) * i, base_w, item_h)
            for i in range(n)
        ]

    if distribution == 'vertical_left':
        # N items stacked vertically in the left half (kpi_with_chart pattern)
        # Keep the template's left/width; only stack vertically
        available = slide_h - base_top - Inches(0.50)
        gap = Inches(0.15)
        if n == 1:
            return [(base_left, base_top, base_w, base_h)]
        item_h = (available - gap * (n - 1)) // n
        return [
            (base_left, base_top + (item_h + gap) * i, base_w, item_h)
            for i in range(n)
        ]

    if distribution == 'horizontal_alternating':
        # Items along a horizontal axis, alternating above (even i) / below (odd i)
        # Spread from the template's left to (slide_w - margin)
        # The vertical offset is encoded in the template by where the user
        # placed the marker — we just shift X.
        left_x = margin
        right_x = slide_w - margin - base_w
        if n == 1:
            return [((left_x + right_x) // 2, base_top, base_w, base_h)]
        span = right_x - left_x
        return [
            (left_x + (span * i) // (n - 1), base_top, base_w, base_h)
            for i in range(n)
        ]

    if distribution == 'single_column':
        # N items stacked vertically in 1 column, full width preserved.
        # Caller paginates when N > the layout's max-per-page (typically 7).
        if n == 1:
            return [(base_left, base_top, base_w, base_h)]
        # Chantier 11 — spacious distribution: aim for gap = 0.5 × item_h
        # (so vertical pitch = 1.5 × item_h). Compress only if needed.
        available_h = slide_h - base_top - Inches(0.50)
        target_gap = base_h // 2
        total_with_target = n * base_h + (n - 1) * target_gap
        if total_with_target <= available_h:
            gap = target_gap
            item_h = base_h
        else:
            # Reduce gap before compressing items
            gap = max(Inches(0.05), (available_h - n * base_h) // max(1, n - 1))
            per_item_max = (available_h - gap * (n - 1)) // n
            item_h = min(base_h, per_item_max)
        return [
            (base_left, base_top + (item_h + gap) * i, base_w, item_h)
            for i in range(n)
        ]

    if distribution == 'grid_2cols':
        # Two columns, items 1..ceil(n/2) on the left, rest on the right
        col_w = base_w
        gap_x = Inches(0.50)
        rows = (n + 1) // 2
        row_h = base_h + Inches(0.20)
        positions = []
        for i in range(n):
            col = i % 2
            row = i // 2
            x = base_left + col * (col_w + gap_x)
            y = base_top + row * row_h
            positions.append((x, y, base_w, base_h))
        return positions

    # Default fallback
    return [(base_left + Inches(0.5) * i, base_top, base_w, base_h)
            for i in range(n)]


# -----------------------------------------------------------------------------
# XML-level positioning of groups
# -----------------------------------------------------------------------------
def _shift_group(grp_element, dx: int, dy: int) -> None:
    """Shift a `<p:grpSp>` by (dx, dy) EMU.

    Only the group's `<a:off>` is updated — NOT `<a:chOff>`. The OOXML
    group transform formula renders children as
        screen = grp.off + (child.off − grp.chOff) × (grp.ext / grp.chExt)
    so shifting both `off` and `chOff` by the same amount cancels out
    (children stay in place). Only `off` should move.
    """
    grpSpPr = grp_element.find(qn('p:grpSpPr'))
    if grpSpPr is None:
        # Standalone shape — update its own spPr/xfrm
        _shift_shape(grp_element, dx, dy)
        return
    xfrm = grpSpPr.find(qn('a:xfrm'))
    if xfrm is None:
        return
    off = xfrm.find(qn('a:off'))
    if off is not None:
        cur_x = int(off.get('x', '0'))
        cur_y = int(off.get('y', '0'))
        off.set('x', str(cur_x + dx))
        off.set('y', str(cur_y + dy))


def _shift_shape(sp_element, dx: int, dy: int) -> None:
    """Shift a standalone shape (`<p:sp>`) by (dx, dy)."""
    spPr = sp_element.find(qn('p:spPr'))
    if spPr is None:
        return
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is None:
        return
    off = xfrm.find(qn('a:off'))
    if off is not None:
        cur_x = int(off.get('x', '0'))
        cur_y = int(off.get('y', '0'))
        off.set('x', str(cur_x + dx))
        off.set('y', str(cur_y + dy))


# -----------------------------------------------------------------------------
# Fill {{ITEM_*}} sub-placeholders inside a duplicated REPEAT_ITEM group
# -----------------------------------------------------------------------------
def _fill_item_placeholders(grp_element, item, index: int, distribution: str) -> None:
    """Walk the deepcopied group XML, find descendants whose `<p:cNvPr@name>`
    matches `{{ITEM_<KEY>}}`, replace text from `item[key.lower()]`.

    For decoration shapes (no text content, e.g., {{ITEM_MARKER}},
    {{ITEM_ICON}}, {{ITEM_BOXE}}), the shape is kept as-is regardless of
    whether spec provides a value — the template's fill/outline IS the
    intended content.

    For text shapes whose key is missing from the spec, the placeholder is
    left untouched (template default text is shown).
    """
    for sp in grp_element.iter(qn('p:sp')):
        nvSpPr = sp.find(qn('p:nvSpPr'))
        if nvSpPr is None:
            continue
        cNvPr = nvSpPr.find(qn('p:cNvPr'))
        if cNvPr is None:
            continue
        name = cNvPr.get('name', '')
        m = ITEM_PLACEHOLDER_RE.match(name)
        if not m:
            continue
        key = m.group(1).lower()

        # Decoration shapes (marker / icon / boxe / bg / bar) have no text
        # content to replace — leave them alone.
        if key in ('marker', 'icon', 'boxe', 'box', 'bg', 'bar', 'background'):
            continue

        value = _resolve_item_value(item, key, index, distribution)
        if value is None:
            # Leave the template's default text in place; do not delete.
            continue
        _set_sp_text(sp, str(value))
        # Chantier 10 — auto-shrink for keys prone to overflow (KPI values,
        # dates, action labels). Reads shape width at runtime, reduces sz
        # in 2pt steps until the heuristic estimate fits.
        if key in SHRINKABLE_ITEM_KEYS:
            _maybe_shrink_to_fit(sp, str(value))


# Keys whose value frequently overflows the template shape width (KPI values,
# dates, etc.). These get auto-shrink applied after substitution.
SHRINKABLE_ITEM_KEYS = {
    # Keys with short single-line content that the moteur should auto-shrink
    # when the spec value is too long for the template shape width.
    # `milestone`, `title`, `text`, `bullets` are intentionally NOT in this
    # list because their template shapes have `<a:spAutoFit/>` and are
    # expected to wrap onto multiple lines.
    'value', 'date', 'action', 'owner',
}


def _set_text_anchor(sp_element, anchor: str) -> None:
    """Set the vertical text anchor on a `<p:sp>` by patching its `<a:bodyPr>`.

    Args:
        anchor: 't' (top), 'ctr' (center), 'b' (bottom).
    """
    txBody = sp_element.find(qn('p:txBody'))
    if txBody is None:
        return
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is None:
        return
    bodyPr.set('anchor', anchor)


def _maybe_shrink_to_fit(sp_element, text: str, min_sz: int = 1000) -> None:
    """If the heuristic estimates `text` overflows the shape's width at the
    current font size, reduce the `sz` attribute on every run in steps of
    200 (= 2pt) until it fits, with a floor at `min_sz` (default 10pt).

    The heuristic is approximate: it estimates char width as
    ``sz_pt × 0.55`` (Arial bold) or ``sz_pt × 0.50`` (regular). Acceptable
    for short single-line content (numbers, dates, short labels) — not for
    multi-line bullets.
    """
    spPr = sp_element.find(qn('p:spPr'))
    if spPr is None:
        return
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is None:
        return
    ext = xfrm.find(qn('a:ext'))
    if ext is None:
        return
    shape_w_emu = int(ext.get('cx', '0'))
    if shape_w_emu == 0:
        return

    # Collect all rPr elements (multi-run text)
    rPr_list = sp_element.findall('.//' + qn('a:rPr'))
    if not rPr_list:
        return
    first_rPr = rPr_list[0]
    sz_str = first_rPr.get('sz')
    if not sz_str:
        return
    sz = int(sz_str)
    is_bold = first_rPr.get('b') == '1'
    factor = 0.55 if is_bold else 0.50

    def width_for(sz_val: int) -> int:
        # sz is in hundredths of a point; 1pt = 12700 EMU
        char_emu = (sz_val / 100.0) * factor * 12700
        return int(len(text) * char_emu)

    target = int(shape_w_emu * 0.95)  # 5% safety margin
    while sz > min_sz and width_for(sz) > target:
        sz -= 200

    # Apply new sz to all runs
    for rPr in rPr_list:
        rPr.set('sz', str(sz))


def _resolve_item_value(item, key: str, index: int, distribution: str):
    """Look up `key` in `item`. Special keys: 'number' auto-fills index+1
    if not present; 'marker' is intentionally ignored (styling)."""
    if not isinstance(item, dict):
        # If item is a plain string, the single text slot is filled with it
        if key in ('text', 'title', 'label', 'name'):
            return item
        return None
    if key in item:
        return item[key]
    # Auto-fill for {{ITEM_NUMBER}}
    if key == 'number':
        return f'{index + 1:02d}'
    return None


def _set_sp_text(sp_element, text: str) -> None:
    """Replace the text content of a `<p:sp>`, preserving the first run's
    formatting (font, size, color, bold, italic) of the original text."""
    txBody = sp_element.find(qn('p:txBody'))
    if txBody is None:
        return
    # Capture the first run's `<a:rPr>` as a style template
    first_p = txBody.find(qn('a:p'))
    if first_p is None:
        return
    first_r = first_p.find(qn('a:r'))
    rPr_template = None
    if first_r is not None:
        rPr = first_r.find(qn('a:rPr'))
        if rPr is not None:
            rPr_template = deepcopy(rPr)
    # Remove all <a:p> elements
    for p in list(txBody.findall(qn('a:p'))):
        txBody.remove(p)
    # Split text by newlines into paragraphs (preserve bullets if multi-line)
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = etree.SubElement(txBody, qn('a:p'))
        # Preserve original pPr if present in the first paragraph of source
        # (kept for first paragraph; subsequent paragraphs use default)
        if i == 0 and first_p is not None:
            src_pPr = first_p.find(qn('a:pPr'))
            if src_pPr is not None:
                p.insert(0, deepcopy(src_pPr))
        r = etree.SubElement(p, qn('a:r'))
        if rPr_template is not None:
            r.append(deepcopy(rPr_template))
        t = etree.SubElement(r, qn('a:t'))
        t.text = line


# -----------------------------------------------------------------------------
# Replace top-level shape text (preserve formatting)
# -----------------------------------------------------------------------------
def _replace_text_keep_style(shape, text: str) -> None:
    """Replace the text of a top-level shape, preserving formatting."""
    if not shape.has_text_frame:
        return
    _set_sp_text(shape._element, text)


# -----------------------------------------------------------------------------
# Chart placeholder — generate matplotlib PNG and insert at placeholder pos
# -----------------------------------------------------------------------------
def _process_chart_placeholder(slide, spec: dict) -> None:
    """If both the `{{CHART_PLACEHOLDER}}` shape and `spec.chart` exist,
    render a matplotlib chart and insert it at the placeholder's position.

    - Placeholder present + chart spec present → render and replace.
    - Placeholder present + no chart spec → leave for simple-placeholder
      processor (it'll be deleted).
    - Chart spec present + no placeholder → no-op (anomaly, silently ignored).
    """
    chart_spec = spec.get('chart')
    if not chart_spec:
        return

    placeholder = None
    for sh in list(_iter_top_level_shapes(slide)):
        if sh.name == '{{CHART_PLACEHOLDER}}':
            placeholder = sh
            break
    if placeholder is None:
        return

    # Lazy import to keep matplotlib optional for users who never render charts
    try:
        import chart_engine
    except ImportError:
        return
    try:
        png_bytes, render_w, render_h = chart_engine.render_chart_to_png(
            chart_spec, placeholder.width, placeholder.height
        )
    except ImportError:
        # matplotlib missing — leave placeholder for normal cleanup
        return

    # Centre the rendered image inside the placeholder frame. For most types
    # render_w/h == placeholder.w/h so offsets are zero; for pie/donut the
    # render is squared, so we centre horizontally and vertically.
    frame_left, frame_top = placeholder.left, placeholder.top
    frame_w, frame_h = placeholder.width, placeholder.height
    img_left = frame_left + (frame_w - render_w) // 2
    img_top = frame_top + (frame_h - render_h) // 2
    import io
    pic = slide.shapes.add_picture(io.BytesIO(png_bytes), img_left, img_top,
                                    width=render_w, height=render_h)
    placeholder._element.getparent().remove(placeholder._element)
    # Chantier 18 — drop shadow on the chart picture
    _apply_drop_shadow(pic._element)


# -----------------------------------------------------------------------------
# Image placeholders
# -----------------------------------------------------------------------------
# Chantier 12 — slides whose `{{IMAGE}}` zone is meant to be diagonally
# cropped by the layout's design. The moteur copies the layout placeholder's
# custGeom onto the inserted picture so the photo respects the cut.
SLIDES_WITH_DIAGONAL_OVERLAY = {
    'cover', 'agenda_diagonal', 'section_diagonal', 'closing_diagonal',
}


def _process_image_placeholders(slide, spec: dict) -> None:
    """Replace any `{{IMAGE}}` shape with an image.

    Resolution order:
      1. `spec.image` (absolute path) — user-provided local file.
      2. Auto-fetch from Pexels/Picsum if `spec._auto_images` is True
         (default). Keyword comes from `spec.image_keyword`, then derived
         from `spec.title`, then a layout-based default.

    For layouts in `SLIDES_WITH_DIAGONAL_OVERLAY`, the inserted picture's
    geometry is replaced by the layout's custGeom so the photo is cut by
    the design's diagonal (Chantier 12).
    """
    image_path = spec.get('image')
    if image_path:
        path = Path(image_path)
        if not path.exists():
            return
        _insert_image_for_each_placeholder(slide, lambda: str(path))
        return

    if not spec.get('_auto_images', True):
        return
    img_sh = _find_image_placeholder(slide)
    if img_sh is None:
        return

    try:
        import image_engine
    except ImportError:
        return

    keyword = spec.get('image_keyword')
    if not keyword:
        keyword = image_engine.extract_keyword_from_title(spec.get('title', ''))
    if not keyword:
        cSld = slide._element.find(qn('p:cSld'))
        layout_name = cSld.get('name') if cSld is not None else None
        keyword = image_engine.LAYOUT_DEFAULT_KEYWORDS.get(layout_name, 'corporate')

    try:
        img_bytes = image_engine.fetch_image_for_slide(
            keyword, img_sh.width, img_sh.height, timeout=8.0,
        )
    except Exception:
        return
    if not img_bytes:
        return

    import io
    left, top, w, h = img_sh.left, img_sh.top, img_sh.width, img_sh.height
    pic = slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, w, h)
    img_sh._element.getparent().remove(img_sh._element)

    # Chantier 12 — diagonal overlay: copy the layout placeholder's custGeom
    # onto the freshly inserted picture so the photo is clipped by the design.
    cSld = slide._element.find(qn('p:cSld'))
    layout_name = cSld.get('name') if cSld is not None else None
    if layout_name in SLIDES_WITH_DIAGONAL_OVERLAY:
        _apply_layout_custgeom_to_picture(slide, pic)


def _find_image_placeholder(slide):
    for sh in _iter_top_level_shapes(slide):
        if sh.name == '{{IMAGE}}':
            return sh
    return None


def _insert_image_for_each_placeholder(slide, path_factory):
    """Replace every `{{IMAGE}}` shape with a picture from `path_factory()`."""
    to_remove = []
    for sh in list(_iter_top_level_shapes(slide)):
        if sh.name != '{{IMAGE}}':
            continue
        left, top, w, h = sh.left, sh.top, sh.width, sh.height
        slide.shapes.add_picture(path_factory(), left, top, w, h)
        to_remove.append(sh)
    for sh in to_remove:
        sh._element.getparent().remove(sh._element)


def _apply_layout_custgeom_to_picture(slide, pic_shape) -> None:
    """Find the layout's Image placeholder (or any shape carrying a custGeom
    that covers the picture area) and copy its `<a:custGeom>` onto the
    picture's `<p:spPr>`, replacing the default `<a:prstGeom prst="rect"/>`.

    No-op if no suitable source geom is found on the slide layout.
    """
    layout = slide.slide_layout
    src_custgeom = None
    for lsh in layout.shapes:
        spPr = lsh._element.find(qn('p:spPr'))
        if spPr is None:
            continue
        cust = spPr.find(qn('a:custGeom'))
        if cust is not None:
            src_custgeom = cust
            break
    if src_custgeom is None:
        return

    pic_spPr = pic_shape._element.find(qn('p:spPr'))
    if pic_spPr is None:
        return
    # Remove the default prstGeom
    for prst in pic_spPr.findall(qn('a:prstGeom')):
        pic_spPr.remove(prst)
    # And drop any pre-existing custGeom (defensive)
    for existing in pic_spPr.findall(qn('a:custGeom')):
        pic_spPr.remove(existing)
    # Insert the layout's custGeom (deep-copy so we don't mutate the layout)
    pic_spPr.append(deepcopy(src_custgeom))


# -----------------------------------------------------------------------------
# Utilities
# -----------------------------------------------------------------------------
def _iter_top_level_shapes(slide):
    """Yield top-level shapes from a slide. Avoids re-fetching `slide.shapes`
    after mutation (`list(...)` at the call site if needed)."""
    return slide.shapes


# =============================================================================
# Chantier 14 — canvas_blank freeform composition engine
# =============================================================================
# Layout coordinates (EMU). The canvas_blank slide is 10.00" × 5.625".
_CB_MARGIN_L = 365760           # 0.40" left margin for blocks area
_CB_FULL_H = 5143500            # 5.625" slide height
_CB_CONTENT_W = 8412480         # 9.20" usable width (= 10.0 - 0.4 left - 0.4 right)

# Chantier 16 — footer reserve so blocks never collide with the master's
# logo + orange bar + page number at the bottom of the slide.
# Master's first footer element (Rectangle 74) starts at 5.09" — we reserve
# 0.95" to leave ~0.40" of breathing room above the visible footer.
_CB_FOOTER_RESERVE = 869280     # 0.95"
_CB_CONTENT_BOTTOM = _CB_FULL_H - _CB_FOOTER_RESERVE  # ≈ 4.675" → 4274220 EMU

# Limits / warnings
_CB_MAX_BLOCKS = 6
_CB_MAX_BULLETS = 5
_CB_TITLE_WARN_LEN = 120


def _render_canvas_blank_freeform(slide, spec: dict) -> None:
    """Compose a freeform canvas_blank slide from `spec.blocks`.

    Chantier 15 — respects the template's actual composition. The
    `{{TITLE}}` placeholder is filled by the upstream simple-placeholder
    processor (with its template-defined position/font/color). This
    function ONLY positions the `blocks` array in the area below the title.

    If spec carries `eyebrow` / `takeaway` / `source` but the template has
    no matching `{{EYEBROW}}` / `{{TAKEAWAY}}` / `{{SOURCE}}` placeholder,
    those fields are ignored (with a warning on stderr) — the user can
    enrich the template manually if desired.
    """
    import sys

    # Warn for spec fields that have no matching placeholder in the template.
    # By this point, the simple processor has filled the existing ones; the
    # missing ones are simply absent from the slide.
    present_template_names = {
        sh.name for sh in _iter_top_level_shapes(slide)
        if isinstance(sh.name, str) and sh.name.startswith('{{')
    }
    for field, ph_name in [('eyebrow',  '{{EYEBROW}}'),
                           ('takeaway', '{{TAKEAWAY}}'),
                           ('source',   '{{SOURCE}}')]:
        if spec.get(field) and ph_name not in present_template_names:
            print(f"canvas_blank: {field} ignored — no {ph_name} placeholder "
                  f"in template", file=sys.stderr)

    title = spec.get('title', '')
    if title and len(title) > _CB_TITLE_WARN_LEN:
        print(f"canvas_blank: title length {len(title)} exceeds "
              f"{_CB_TITLE_WARN_LEN} (will likely wrap)", file=sys.stderr)

    blocks = spec.get('blocks') or []
    if not blocks:
        return
    if len(blocks) > _CB_MAX_BLOCKS:
        print(f"canvas_blank: {len(blocks)} blocks > {_CB_MAX_BLOCKS}, "
              f"truncating to first {_CB_MAX_BLOCKS}", file=sys.stderr)
        blocks = blocks[:_CB_MAX_BLOCKS]

    # Compute the blocks area dynamically from the actual {{TITLE}} position
    blocks_top, blocks_h = _cb_compute_blocks_area(slide)
    rects = _cb_compute_block_rects(blocks, _CB_MARGIN_L, blocks_top,
                                     _CB_CONTENT_W, blocks_h)

    # Chantier 19/20 — unify kpi_card value font sizes across all KPI blocks
    # on this slide : compute the per-card optimal via the dynamic layout
    # helper (which accounts for label height), take min, inject into each.
    kpi_indices = [i for i, b in enumerate(blocks) if b.get('type') == 'kpi_card']
    if len(kpi_indices) >= 2:
        per_max = []
        for i in kpi_indices:
            v = str(blocks[i].get('value', ''))
            lbl = str(blocks[i].get('label', ''))
            sz, _, _ = _cb_kpi_card_dynamic_layout(
                v, lbl, rects[i], override_size_pt=None, floor_pt=14)
            per_max.append(sz)
        unified_pt = min(per_max)
        for i in kpi_indices:
            blocks[i]['_unified_size_pt'] = unified_pt

    for block, rect in zip(blocks, rects):
        _cb_render_block(slide, block, rect)


def _cb_compute_blocks_area(slide):
    """Return (blocks_top, blocks_height) computed from the slide's
    `{{TITLE}}` placeholder position. The bottom of the area is capped at
    `_CB_CONTENT_BOTTOM` so blocks never collide with the master's footer
    (logo + orange bar + page number)."""
    title_sh = None
    for sh in _iter_top_level_shapes(slide):
        if sh.name == '{{TITLE}}':
            title_sh = sh
            break

    SAFETY_FLOOR = int(1.10 * 914400)        # 1.10" min from top
    GAP_BELOW_TITLE = int(0.30 * 914400)     # 0.30"

    if title_sh is not None:
        title_bottom = title_sh.top + title_sh.height
        blocks_top = max(title_bottom + GAP_BELOW_TITLE, SAFETY_FLOOR)
    else:
        blocks_top = SAFETY_FLOOR

    blocks_h = max(_CB_CONTENT_BOTTOM - blocks_top, int(2.0 * 914400))
    return blocks_top, blocks_h


# Header-band helpers (eyebrow/title/takeaway/source) removed in Chantier 15.
# The canvas_blank template carries only `{{TITLE}}`, which is filled by the
# standard simple-placeholder processor. If the user later adds {{EYEBROW}},
# {{TAKEAWAY}} or {{SOURCE}} placeholders to the template, the simple
# processor will pick them up automatically — no special code needed.


# -----------------------------------------------------------------------------
# Grid layout
# -----------------------------------------------------------------------------
def _cb_compute_block_rects(blocks, base_left, base_top, area_w, area_h):
    """Return a list of (L, T, W, H) tuples for each block.

    Layouts:
      - 1 block  : 1×1 full
      - 2 blocks : 1×2 (side by side)
      - 3 blocks : 1×3 columns
      - 4 blocks : 2×2 grid
      - 5 blocks : 3 top + 2 bottom (asymmetric)
      - 6 blocks : 3×2 grid

    Special: if exactly one image/chart present with 2+ blocks, that visual
    takes the right half, other blocks stack on the left half.
    """
    n = len(blocks)
    gap = int(0.15 * 914400)  # 0.15" inter-block

    # Detect a single visual asymmetric case
    visual_indices = [i for i, b in enumerate(blocks)
                      if b.get('type') in ('image', 'chart')]
    if n >= 2 and len(visual_indices) == 1:
        # Split: visual right (45%), others left
        right_w = int(area_w * 0.45)
        left_w = area_w - right_w - gap
        v_idx = visual_indices[0]
        # Left column: all NON-visual blocks
        non_visual = [i for i in range(n) if i != v_idx]
        m = len(non_visual)
        left_block_h = (area_h - gap * (m - 1)) // max(1, m)
        rects = [None] * n
        for j, i in enumerate(non_visual):
            rects[i] = (base_left, base_top + j * (left_block_h + gap),
                        left_w, left_block_h)
        rects[v_idx] = (base_left + left_w + gap, base_top, right_w, area_h)
        return rects

    if n == 1:
        return [(base_left, base_top, area_w, area_h)]
    if n == 2:
        w = (area_w - gap) // 2
        return [(base_left, base_top, w, area_h),
                (base_left + w + gap, base_top, w, area_h)]
    if n == 3:
        w = (area_w - 2 * gap) // 3
        return [(base_left + i * (w + gap), base_top, w, area_h) for i in range(3)]
    if n == 4:
        w = (area_w - gap) // 2
        h = (area_h - gap) // 2
        return [
            (base_left,             base_top,           w, h),
            (base_left + w + gap,   base_top,           w, h),
            (base_left,             base_top + h + gap, w, h),
            (base_left + w + gap,   base_top + h + gap, w, h),
        ]
    if n == 5:
        # 3 on top, 2 on bottom (centered)
        w3 = (area_w - 2 * gap) // 3
        w2 = (area_w - gap) // 2
        h = (area_h - gap) // 2
        rects = []
        for i in range(3):
            rects.append((base_left + i * (w3 + gap), base_top, w3, h))
        bottom_y = base_top + h + gap
        for i in range(2):
            rects.append((base_left + i * (w2 + gap), bottom_y, w2, h))
        return rects
    # n == 6: 3×2 grid
    w = (area_w - 2 * gap) // 3
    h = (area_h - gap) // 2
    rects = []
    for row in range(2):
        for col in range(3):
            rects.append((base_left + col * (w + gap),
                          base_top + row * (h + gap), w, h))
    return rects


# -----------------------------------------------------------------------------
# Block renderers
# -----------------------------------------------------------------------------
def _cb_render_block(slide, block: dict, rect) -> None:
    block_type = block.get('type', 'text')
    renderer = {
        'kpi_card': _cb_render_kpi_card,
        'bullets':  _cb_render_bullets,
        'text':     _cb_render_text,
        'image':    _cb_render_image,
        'chart':    _cb_render_chart,
        'quote':    _cb_render_quote,
    }.get(block_type)
    if renderer is None:
        return  # unknown block type, silently skip
    renderer(slide, block, rect)


def _cb_kpi_card_dynamic_layout(value: str, label: str, rect,
                                override_size_pt=None,
                                floor_pt: int = 14):
    """Chantier 20 — return (value_size_pt, value_rect, label_rect) for a
    kpi_card. The value+label block is centered vertically when there is
    headroom, tight-packed from the top otherwise. The value font size is
    chosen to leave room for the (estimated) label height, preventing
    value/label overlap and never overflowing the card bottom.
    """
    L, T, W, H = rect
    EMU_PER_PT = 12700
    pad = int(0.10 * 914400)        # 0.10" inner padding
    spacing = int(0.05 * 914400)    # 0.05" gap between value and label
    available_w = W - 2 * pad

    # Estimate label height (10pt uppercase, char width ≈ 0.60 × pt)
    label_text = (label or '').upper()
    n_chars = max(1, len(label_text))
    chars_per_line = max(1, int(available_w / (10 * EMU_PER_PT * 0.60)))
    n_lines = max(1, (n_chars + chars_per_line - 1) // chars_per_line)
    label_h = int(n_lines * 10 * EMU_PER_PT * 1.25)

    inner_h = H - 2 * pad
    # Cap label to at most 50% of inner height (so the value keeps room)
    max_label_h = max(int(10 * EMU_PER_PT * 1.25), inner_h // 2)
    if label_h > max_label_h:
        label_h = max_label_h

    # Value's vertical budget
    value_max_h = inner_h - label_h - spacing
    if value_max_h < 14 * EMU_PER_PT:
        # Cramped card → ensure at least 14pt for the value
        value_max_h = max(int(14 * EMU_PER_PT), inner_h // 2)
        label_h = max(int(10 * EMU_PER_PT * 1.25),
                      inner_h - value_max_h - spacing)

    if override_size_pt is not None:
        value_size_pt = override_size_pt
    else:
        value_size_pt = _compute_max_kpi_font_size(
            value, available_w, value_max_h, floor_pt=floor_pt)

    # Actual value height (line height ~1.05)
    value_h = int(value_size_pt * EMU_PER_PT * 1.05)
    if value_h > value_max_h:
        value_h = value_max_h

    total_h = value_h + spacing + label_h
    if total_h <= inner_h:
        top_margin = pad + (inner_h - total_h) // 2
    else:
        # Tight pack — shrink label rect to fit inside the card
        top_margin = pad
        label_h = max(int(0.13 * 914400),
                      H - top_margin - value_h - spacing - pad)

    value_rect = (L + pad, T + top_margin, available_w, value_h)
    label_rect = (L + pad, T + top_margin + value_h + spacing,
                  available_w, label_h)
    return value_size_pt, value_rect, label_rect


def _cb_render_kpi_card(slide, block: dict, rect) -> None:
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    L, T, W, H = rect
    value = str(block.get('value', ''))
    label = str(block.get('label', ''))
    color_key = block.get('color', 'orange')
    color_hex = _resolve_color({
        'orange': 'orange', 'navy': 'navy',
        'green': 'accent5', 'red': 'accent6',
    }.get(color_key, 'orange'))

    # Chantier 18 — white card background with thin gray border + drop shadow
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, L, T, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string('FFFFFF')
    bg.line.color.rgb = RGBColor.from_string('E8E9F2')
    bg.line.width = Emu(int(0.5 * 12700))   # 0.5pt
    _apply_drop_shadow(bg._element)

    # Chantier 20 — dynamic in-card positioning (anti-overlap)
    override = block.get('_unified_size_pt')
    size_pt, value_rect, label_rect = _cb_kpi_card_dynamic_layout(
        value, label, rect, override_size_pt=override, floor_pt=14)

    # Value textbox
    vL, vT, vW, vH = value_rect
    tb = slide.shapes.add_textbox(vL, vT, vW, vH)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_top = Emu(0)
    tf.margin_right = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    from pptx.enum.text import PP_ALIGN
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = value
    run.font.name = "Arial"
    run.font.size = Pt(size_pt)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string(color_hex)

    # Label textbox
    lL, lT, lW, lH = label_rect
    tb2 = slide.shapes.add_textbox(lL, lT, lW, lH)
    tf2 = tb2.text_frame
    tf2.margin_left = Emu(0); tf2.margin_top = Emu(0)
    tf2.margin_right = Emu(0); tf2.margin_bottom = Emu(0)
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = label.upper()
    run2.font.name = "Arial"
    run2.font.size = Pt(10)
    run2.font.bold = True
    run2.font.color.rgb = RGBColor.from_string(_resolve_color('navy'))


def _cb_render_bullets(slide, block: dict, rect) -> None:
    import sys
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    L, T, W, H = rect
    items = block.get('items') or []
    if len(items) > _CB_MAX_BULLETS:
        print(f"canvas_blank bullets block: {len(items)} items > "
              f"{_CB_MAX_BULLETS}, truncating", file=sys.stderr)
        items = items[:_CB_MAX_BULLETS]
    tb = slide.shapes.add_textbox(L, T, W, H)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_top = Emu(0)
    orange = _resolve_color('orange')
    navy = _resolve_color('navy')
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        bullet_run = p.add_run()
        bullet_run.text = "•  "
        bullet_run.font.name = "Arial"
        bullet_run.font.size = Pt(12)
        bullet_run.font.bold = True
        bullet_run.font.color.rgb = RGBColor.from_string(orange)
        text_run = p.add_run()
        text_run.text = str(item)
        text_run.font.name = "Arial"
        text_run.font.size = Pt(12)
        text_run.font.color.rgb = RGBColor.from_string(navy)


def _cb_render_text(slide, block: dict, rect) -> None:
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    L, T, W, H = rect
    content = str(block.get('content', ''))
    tb = slide.shapes.add_textbox(L, T, W, H)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_top = Emu(0)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = content
    run.font.name = "Arial"
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor.from_string(_resolve_color('navy'))


def _cb_render_image(slide, block: dict, rect) -> None:
    L, T, W, H = rect
    path = block.get('path')
    keyword = block.get('keyword')
    pic = None
    if path:
        from pathlib import Path
        p = Path(path)
        if p.exists():
            pic = slide.shapes.add_picture(str(p), L, T, W, H)
    elif keyword:
        try:
            import image_engine
            img_bytes = image_engine.fetch_image_for_slide(keyword, W, H, timeout=8.0)
        except Exception:
            img_bytes = None
        if img_bytes:
            import io
            pic = slide.shapes.add_picture(io.BytesIO(img_bytes), L, T, W, H)
    # Chantier 18 — drop shadow on canvas_blank image blocks (rect, no diagonal)
    if pic is not None:
        _apply_drop_shadow(pic._element)


def _cb_render_chart(slide, block: dict, rect) -> None:
    L, T, W, H = rect
    chart_spec = block.get('chart_spec') or block.get('chart')
    if not chart_spec:
        return
    try:
        import chart_engine
        png_bytes, w_render, h_render = chart_engine.render_chart_to_png(chart_spec, W, H)
    except Exception:
        return
    import io
    # For pie/donut chart_engine returns a squared image — centre it in the rect
    img_left = L + (W - w_render) // 2
    img_top = T + (H - h_render) // 2
    pic = slide.shapes.add_picture(io.BytesIO(png_bytes), img_left, img_top,
                                    width=w_render, height=h_render)
    # Chantier 18 — drop shadow on canvas_blank chart blocks
    _apply_drop_shadow(pic._element)


def _cb_render_quote(slide, block: dict, rect) -> None:
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    L, T, W, H = rect
    content = str(block.get('content', ''))
    author = str(block.get('author', ''))

    # Quote text: italic navy, 90% of height
    quote_h = int(H * 0.80) if author else H
    tb = slide.shapes.add_textbox(L, T, W, quote_h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_top = Emu(0)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f'« {content} »'
    run.font.name = "Arial"
    run.font.size = Pt(14)
    run.font.italic = True
    run.font.color.rgb = RGBColor.from_string(_resolve_color('navy'))

    if author:
        tb2 = slide.shapes.add_textbox(L, T + quote_h, W, H - quote_h)
        tf2 = tb2.text_frame
        tf2.margin_left = Emu(0); tf2.margin_top = Emu(0)
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = author.upper()
        run2.font.name = "Arial"
        run2.font.size = Pt(9)
        run2.font.bold = True
        run2.font.color.rgb = RGBColor.from_string(_resolve_color('orange'))


# =============================================================================
# Chantier 15 — data_table dynamic rendering
# =============================================================================
_DT_MAX_COLS = 6
_DT_MAX_ROWS = 8
_DT_CELL_WARN_LEN = 30

# Position and size of the table area on the data_table slide
_DT_LEFT_EMU = int(0.5 * 914400)
_DT_TOP_EMU = int(1.2 * 914400)
_DT_WIDTH_EMU = int(9.0 * 914400)
_DT_MAX_HEIGHT_EMU = int(3.7 * 914400)

# Brand colors used by the table (hex without '#')
_DT_NAVY = "14163C"
_DT_WHITE = "FFFFFF"
_DT_ORANGE = "F26622"
_DT_OFFWHITE = "FAFAF7"
_DT_LIGHT_ORANGE = "FDF1EA"
_DT_GRAY = "4A4D6B"


def _process_data_table(slide, spec: dict) -> None:
    """Draw a styled table on the data_table slide based on ``spec['table']``.

    Accepted spec keys (under `table`):
        headers (list[str], required, ≤ 6)
        rows (list[list], required, ≤ 8)
        highlight_column (int, optional, 0-based)
        highlight_row (int, optional, 0-based on data rows)

    Style: navy header bg + white bold text ; alternated white/off-white body
    rows ; first column bold left-aligned ; other columns centered. Optional
    orange-bold column highlight + light-orange row highlight. Auto-shrink font
    + padding if the table overflows the 3.7" content area.
    """
    import sys
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    table_spec = spec.get('table') or {}
    headers = table_spec.get('headers') or []
    rows = table_spec.get('rows') or []
    highlight_col = table_spec.get('highlight_column')
    highlight_row = table_spec.get('highlight_row')

    # --- Validation --------------------------------------------------------
    if not headers:
        return
    if len(headers) > _DT_MAX_COLS:
        raise ValueError(
            f"data_table: max {_DT_MAX_COLS} columns, got {len(headers)}"
        )
    if len(rows) > _DT_MAX_ROWS:
        print(f"data_table: {len(rows)} rows > {_DT_MAX_ROWS}, truncating",
              file=sys.stderr)
        rows = rows[:_DT_MAX_ROWS]
    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header

    # Pad short rows with empty strings + warn on overlong cells
    for r_idx, row in enumerate(rows):
        for c_idx, cell in enumerate(row):
            if isinstance(cell, str) and len(cell) > _DT_CELL_WARN_LEN:
                print(f"data_table: cell ({r_idx+1},{c_idx}) has "
                      f"{len(cell)} chars (>{_DT_CELL_WARN_LEN}), may overflow",
                      file=sys.stderr)

    # --- Auto-shrink computation ------------------------------------------
    # Estimate row height = font_size(pt) + 2*padding(pt). At 1pt = 12700 EMU,
    # an inch is 914400 EMU. Convert max height to pt:
    max_h_pt = _DT_MAX_HEIGHT_EMU / 12700.0
    font_size = 10       # body font size in pt
    header_size = 11     # header font size in pt
    padding = 6          # cell padding in pt

    # Wrap factor: more columns → narrower cells → more likely to wrap.
    # 1.0 ≤ 3 cols, 1.5 for 4-5 cols, 2.0 for 6 cols.
    if n_cols <= 3:
        wrap_factor = 1.0
    elif n_cols <= 5:
        wrap_factor = 1.5
    else:
        wrap_factor = 2.0

    def estimated_height_pt(fs, pd):
        # Each body row roughly = (font_size + 2*padding) × wrap_factor (pt)
        header_h = header_size + 2 * pd
        body_h = (fs + 2 * pd) * wrap_factor
        return header_h + len(rows) * body_h

    while estimated_height_pt(font_size, padding) > max_h_pt:
        if padding > 4:
            padding -= 1
        elif font_size > 8:
            font_size -= 1
        else:
            print(f"data_table: still overflowing at fs=8pt pad=4pt — "
                  f"consider fewer rows/columns or shorter content",
                  file=sys.stderr)
            break

    # --- Column widths ----------------------------------------------------
    first_col_w = int(_DT_WIDTH_EMU * 0.30)
    other_cols_w = ((_DT_WIDTH_EMU - first_col_w) // (n_cols - 1)
                    if n_cols > 1 else _DT_WIDTH_EMU)

    # --- Initial row heights (rough — PowerPoint auto-fits actual content) -
    init_total_h = int(estimated_height_pt(font_size, padding) * 12700)
    init_total_h = min(init_total_h, _DT_MAX_HEIGHT_EMU)

    # --- Add table --------------------------------------------------------
    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        _DT_LEFT_EMU, _DT_TOP_EMU,
        _DT_WIDTH_EMU, init_total_h,
    )
    table = table_shape.table

    # Column widths
    table.columns[0].width = first_col_w
    for i in range(1, n_cols):
        table.columns[i].width = other_cols_w

    # --- Style helpers ----------------------------------------------------
    navy_rgb = RGBColor.from_string(_DT_NAVY)
    white_rgb = RGBColor.from_string(_DT_WHITE)
    orange_rgb = RGBColor.from_string(_DT_ORANGE)
    offwhite_rgb = RGBColor.from_string(_DT_OFFWHITE)
    light_orange_rgb = RGBColor.from_string(_DT_LIGHT_ORANGE)

    def style_cell(cell, *, bg_hex, text, font_size_pt, bold, italic,
                   color_hex, align, padding_pt):
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(bg_hex)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(padding_pt)
        tf.margin_right = Pt(padding_pt)
        tf.margin_top = Pt(padding_pt)
        tf.margin_bottom = Pt(padding_pt)
        # Clear existing paragraphs
        for p in list(tf.paragraphs):
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.name = "Arial"
        run.font.size = Pt(font_size_pt)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = RGBColor.from_string(color_hex)

    # --- Headers (row 0) --------------------------------------------------
    for j, h in enumerate(headers):
        align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        style_cell(
            table.cell(0, j),
            bg_hex=_DT_NAVY, text=h, font_size_pt=header_size,
            bold=True, italic=False, color_hex=_DT_WHITE,
            align=align, padding_pt=padding,
        )

    # --- Data rows --------------------------------------------------------
    for i, row in enumerate(rows):
        is_highlight_row = (highlight_row is not None and i == highlight_row)
        row_bg = (_DT_LIGHT_ORANGE if is_highlight_row
                  else (_DT_OFFWHITE if i % 2 == 1 else _DT_WHITE))
        for j in range(n_cols):
            val = row[j] if j < len(row) else ""
            is_highlight_col = (highlight_col is not None and j == highlight_col)
            # First column always bold left-aligned
            if j == 0:
                bold = True
                align = PP_ALIGN.LEFT
                color = _DT_ORANGE if is_highlight_col else _DT_NAVY
            else:
                bold = is_highlight_col
                align = PP_ALIGN.CENTER
                color = _DT_ORANGE if is_highlight_col else _DT_NAVY
            style_cell(
                table.cell(i + 1, j),
                bg_hex=row_bg, text=val, font_size_pt=font_size,
                bold=bold, italic=False, color_hex=color,
                align=align, padding_pt=padding,
            )


# =============================================================================
# Chantier 16 — uniform font size across REPEAT_ITEM copies + closing fix
# =============================================================================

# Per-layout: which sub-placeholder names to harmonize (= take the smallest
# observed sz across all copies and apply it to all). Prevents inconsistent
# sizing produced by the per-item auto-shrink (Chantier 10/11).
_UNIFORM_REPEAT_SHAPES = {
    'roadmap_styled':   ['{{ITEM_DATE}}', '{{ITEM_MILESTONE}}'],
    'next_steps':       ['{{ITEM_ACTION}}', '{{ITEM_OWNER}}', '{{ITEM_DATE}}'],
    'kpi_with_chart':   ['{{KPI_VALUE}}', '{{KPI_LABEL}}'],
    'agenda_diagonal':  ['{{ITEM_TITLE}}'],
    'process_steps':    ['{{ITEM_TITLE}}', '{{ITEM_TEXT}}'],
    'framework_3cards': ['{{ITEM_TITLE}}', '{{ITEM_BULLETS}}'],
    'text_dense_3cols': ['{{ITEM_TITLE}}', '{{ITEM_TEXT}}'],
}


def _apply_uniform_font_size_to_repeats(slide, shape_names: list) -> None:
    """Read the sz attribute of every `<a:rPr>` across all REPEAT_ITEM copies
    in `slide`, group by sub-placeholder name. For each group, take the
    smallest observed sz and apply it to every other copy of the same name.

    Result: all copies of e.g. ``{{ITEM_DATE}}`` share the same font size,
    determined by the longest text (which forced the most aggressive shrink).
    """
    target_names = set(shape_names)
    # Map: shape_name → (min_sz_observed, list_of_rPr_elements)
    by_name = {n: {'min_sz': None, 'rPrs': []} for n in target_names}

    for top_sh in _iter_top_level_shapes(slide):
        if top_sh.shape_type != 6:  # only GROUPs (REPEAT_ITEM copies)
            continue
        for sp in top_sh._element.iter(qn('p:sp')):
            cNvPr = sp.find('.//' + qn('p:cNvPr'))
            if cNvPr is None:
                continue
            sp_name = cNvPr.get('name', '')
            if sp_name not in target_names:
                continue
            for rPr in sp.iter(qn('a:rPr')):
                sz_str = rPr.get('sz')
                if not sz_str:
                    continue
                sz = int(sz_str)
                entry = by_name[sp_name]
                entry['rPrs'].append(rPr)
                if entry['min_sz'] is None or sz < entry['min_sz']:
                    entry['min_sz'] = sz

    # Apply min sz to all collected rPrs
    for name, entry in by_name.items():
        if entry['min_sz'] is None:
            continue
        target_sz = str(entry['min_sz'])
        for rPr in entry['rPrs']:
            rPr.set('sz', target_sz)


# -----------------------------------------------------------------------------
# closing_diagonal — aggressive title shrink for long titles
# -----------------------------------------------------------------------------
def _shrink_closing_title(slide) -> None:
    """Shrink {{TITLE}} on closing_diagonal if its text doesn't fit the shape
    width. Paliers of 4pt from 60→36pt (floor). Also enforces wrap=square +
    spAutoFit-off so PowerPoint breaks on word boundaries.

    Heuristic: estimate the width of the longest word at current size and
    reduce until it fits ``shape_width × 0.95`` (5% margin).
    """
    title_sh = None
    for sh in _iter_top_level_shapes(slide):
        if sh.name == '{{TITLE}}':
            title_sh = sh
            break
    if title_sh is None or not title_sh.has_text_frame:
        return

    shape_w_emu = title_sh.width
    target_emu = int(shape_w_emu * 0.95)

    # Ensure wrap is on (word boundaries)
    txBody = title_sh._element.find(qn('p:txBody'))
    if txBody is not None:
        bodyPr = txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('wrap', 'square')

    # Read current text + sz from first run
    text = title_sh.text_frame.text or ''
    if not text.strip():
        return
    rPr_list = list(title_sh._element.iter(qn('a:rPr')))
    if not rPr_list:
        return
    first_rPr = rPr_list[0]
    sz_str = first_rPr.get('sz')
    if not sz_str:
        return
    sz = int(sz_str)
    is_bold = first_rPr.get('b') == '1'
    factor = 0.55 if is_bold else 0.50

    # Estimate width of the LONGEST word at a given sz
    longest_word = max(text.split(), key=len) if text.split() else text
    n_chars = len(longest_word)

    def word_width_emu(sz_val):
        char_emu = (sz_val / 100.0) * factor * 12700
        return int(n_chars * char_emu)

    # Shrink by 400 (4pt) steps from current sz down to 3600 (36pt) floor
    floor_sz = 3600
    while sz > floor_sz and word_width_emu(sz) > target_emu:
        sz -= 400

    # Apply new sz to all runs
    for rPr in rPr_list:
        rPr.set('sz', str(sz))


# =============================================================================
# Chantier 18 — Premium effects : drop shadow + XXL KPI + chart border
# =============================================================================
def _apply_drop_shadow(shape_element, blur_pt: int = 8,
                       distance_pt: int = 3, alpha_pct: int = 25,
                       angle_deg: int = 45) -> None:
    """Add an outerShdw drop shadow to a shape (via XML). The effect goes
    inside `<p:spPr><a:effectLst>`.

    Args:
        shape_element: the `<p:sp>` or `<p:pic>` lxml element.
        blur_pt: shadow blur in points.
        distance_pt: shadow offset distance in points.
        alpha_pct: shadow opacity 0-100 (% in 1000ths of percent for OOXML).
        angle_deg: shadow direction in degrees (45° = down-right).
    """
    spPr = shape_element.find(qn('p:spPr'))
    if spPr is None:
        return
    # Remove any existing effectLst (idempotent)
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)

    blur_emu = blur_pt * 12700
    dist_emu = distance_pt * 12700
    dir_60k = angle_deg * 60000  # OOXML angles are in 1/60000 degree
    alpha_1k = alpha_pct * 1000   # 25 % → 25000

    effectLst = etree.SubElement(spPr, qn('a:effectLst'))
    outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', str(blur_emu))
    outerShdw.set('dist', str(dist_emu))
    outerShdw.set('dir', str(dir_60k))
    outerShdw.set('algn', 'tl')
    outerShdw.set('rotWithShape', '0')
    srgb = etree.SubElement(outerShdw, qn('a:srgbClr'))
    srgb.set('val', '000000')
    alpha = etree.SubElement(srgb, qn('a:alpha'))
    alpha.set('val', str(alpha_1k))


# Layouts where {{IMAGE}} gets diagonal custGeom — no shadow on these (the
# diagonal cut would render badly with a rectangular shadow).
_NO_SHADOW_LAYOUTS = {'cover', 'agenda_diagonal', 'section_diagonal',
                      'closing_diagonal'}


# =============================================================================
# Chantier 19 — Height-aware KPI sizing (anti-overlap)
# =============================================================================
def _compute_max_kpi_font_size(value_text: str,
                               available_w_emu: int,
                               available_h_emu: int,
                               ceiling_pt: int = 60,
                               floor_pt: int = 24) -> int:
    """Compute the largest font size (in pt) that fits a KPI value in the
    given rectangle. Snaps to multiples of 6 (60/54/48/42/36/30/24).

    Heuristic constants:
      - line height factor ≈ 1.0 (single-line Arial Bold).
      - char width factor ≈ 0.55 (avg Arial Bold including space/digits).

    Returns an int in [floor_pt, ceiling_pt].
    """
    EMU_PER_PT = 12700.0
    # Vertical constraint: a single line takes ≈ font_size pt of height
    max_h_pt = available_h_emu / EMU_PER_PT
    # Horizontal constraint: text width estimate
    n_chars = max(1, len(value_text))
    max_w_pt = available_w_emu / (n_chars * 0.55 * EMU_PER_PT)
    # Combined max (also capped at ceiling)
    raw_pt = min(max_h_pt, max_w_pt, ceiling_pt)
    # Snap down to multiple of 6 (drop fractional)
    snapped = int(raw_pt // 6) * 6
    return max(floor_pt, min(ceiling_pt, snapped))


def _read_sp_dimensions(sp_element):
    """Return (width_emu, height_emu) for a <p:sp> by reading its xfrm/ext,
    or (None, None) if missing."""
    spPr = sp_element.find(qn('p:spPr'))
    if spPr is None:
        return (None, None)
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is None:
        return (None, None)
    ext = xfrm.find(qn('a:ext'))
    if ext is None:
        return (None, None)
    return (int(ext.get('cx', '0')), int(ext.get('cy', '0')))
